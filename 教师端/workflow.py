#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
高中教师超级智能体 · 教学+办公全场景工作流
==========================================
架构：意图路由 + 产物注册表驱动 + 多子工作流并行 + 知识库/联网支撑 + 全链路质量校验
技术栈：LangChain（提示词/模型抽象）+ LangGraph（状态图编排）

覆盖任务域：
- 教学域：教案、PPT、分层习题、互动动画、学情分析、重难点精讲、试卷组卷、分层作业、学生评语、成绩分析
- 办公域：通知公告、工作计划、工作总结、发言稿、会议纪要、班会方案、评优/职称申报材料

核心机制：
- 产物注册表 PRODUCERS：每个产物声明 名称/领域/提示词/校验器，新增产物只需加一行。
- 意图识别：解析产物需求（含 productsExplicit 判断用户是否明确点名），智能决策执行清单 tasks。
- 多选追问：产物不明确时给出「教学/办公」分组多选 + 缺失字段补充；指令明确则不问。
- 扇出并行：所有产物节点从任务拆解后扇出，节点内部按 tasks 判断是否真正执行。
- 全局上下文：subject/topic/student_tags/quantities 一次赋值，全流程复用。
- 闭环校验：每个产物「生成→校验→循环优化」，不合格自动重写。
- 知识检索 + 联网搜索：任务拆解前拉取教材知识图谱与课标/考点资料。
"""

from __future__ import annotations

import json
import logging
import os
import re
import sys
import threading
import time
import uuid
from concurrent.futures import ThreadPoolExecutor, as_completed
from dataclasses import dataclass, field
from typing import Any, Callable, Dict, List, Optional, Tuple, TypedDict

# ---------------------------------------------------------------------------
# 0. 日志
# ---------------------------------------------------------------------------
logging.basicConfig(level=logging.INFO, format="[教师智能体] %(levelname)s %(message)s")
log = logging.getLogger("teacher-agent")

# ---------------------------------------------------------------------------
# 1. 依赖探测（LangChain / LangGraph 缺失时仍可用 Mock 模式运行）
# ---------------------------------------------------------------------------
try:
    from langgraph.graph import START, END, StateGraph
    LANGGRAPH_AVAILABLE = True
except Exception as e:  # pragma: no cover
    LANGGRAPH_AVAILABLE = False
    log.warning("LangGraph 未安装，将退化为串行执行：%s", e)

try:
    from langchain_core.messages import AIMessage, HumanMessage, SystemMessage
    from langchain_core.language_models.chat_models import BaseChatModel

    LANGCHAIN_AVAILABLE = True
except Exception as e:  # pragma: no cover
    LANGCHAIN_AVAILABLE = False
    log.warning("LangChain 未安装：%s", e)


# ===========================================================================
# 2. 全局配置
# ===========================================================================
def _env_int(key: str, default: int) -> int:
    """安全读取整型环境变量：非法值回退默认，避免导入期崩溃。"""
    try:
        return int(os.getenv(key, "").strip() or default)
    except Exception:
        return default


def _env_float(key: str, default: float) -> float:
    """安全读取浮点环境变量：非法值回退默认，避免导入期崩溃。"""
    try:
        return float(os.getenv(key, "").strip() or default)
    except Exception:
        return default


@dataclass
class Settings:
    """全局配置，可通过环境变量覆盖。"""

    # 模型后端：mock | openai | dashscope | zhipu
    provider: str = os.getenv("AGENT_LLM_PROVIDER", "mock").lower()
    model: str = os.getenv("AGENT_LLM_MODEL", "qwen-plus")
    api_key: str = os.getenv("AGENT_LLM_API_KEY", "")
    base_url: str = os.getenv("AGENT_LLM_BASE_URL", "")
    temperature: float = _env_float("AGENT_LLM_TEMPERATURE", 0.4)
    # 单次 LLM 请求超时（秒）：防止慢模型/网络卡死拖垮整个工作流。
    # 分节并行后单次请求短小，180s 足够；超时即快速失败并切换备用模型，避免阻塞。
    llm_timeout: int = _env_int("AGENT_LLM_TIMEOUT", 180)
    # 备用模型链 [(model, api_key, base_url), ...]：主模型超时/失败时自动切换
    llm_candidates: List[Tuple[str, str, str]] = field(default_factory=list)
    # 质量校验相关
    max_refine_rounds: int = _env_int("AGENT_MAX_REFINE", 3)
    # 质量门：产物评分低于该值判定不合格，必须重新生成
    quality_min_score: int = _env_int("AGENT_QUALITY_MIN_SCORE", 75)
    # 质量门开关：默认强制启用，不合格产物重新生成，仍不合格则标记失败（绝不静默落盘）
    quality_gate: bool = os.getenv("AGENT_QUALITY_GATE", "1") != "0"
    # 知识库路径
    kg_path: str = os.getenv("AGENT_KG_PATH", "")
    verbose: bool = True


SETTINGS = Settings()


# ===========================================================================
# 3. 产物注册表（教学 + 办公全场景）
#    id: 产物唯一标识；name: 中文名；domain: teaching/office；
#    prompt: 提示词模板；validator: 可选校验器；node: 可选专用节点（默认走通用节点）；
#    export: 落盘格式（docx/xlsx/pptx/md）
# ===========================================================================
PRODUCERS: Dict[str, Dict[str, Any]] = {
    # ---------------- 教学域 ----------------
    "lesson_plan": {
        "name": "教案", "domain": "teaching", "prompt": "lesson_plan", "node": "lesson_plan_node",
        "export": "docx",
    },
    "ppt": {"name": "PPT课件", "domain": "teaching", "prompt": "ppt", "node": "ppt_node", "export": "pptx"},
    "exercise": {"name": "分层习题", "domain": "teaching", "prompt": "exercise", "node": "exercise_node", "export": "docx"},
    "animation": {"name": "互动动画", "domain": "teaching", "prompt": "animation", "node": "animation_node", "export": "md"},
    "learning_analysis": {
        "name": "学情分析", "domain": "teaching", "prompt": "learning_analysis", "node": "learning_analysis_node",
        "export": "docx",
    },
    "explainer": {"name": "重难点精讲", "domain": "teaching", "prompt": "explainer", "node": "explainer_node", "export": "md"},
    "exam_paper": {"name": "试卷组卷", "domain": "teaching", "prompt": "exam_paper", "validator": "paper", "export": "docx"},
    "homework": {"name": "分层作业", "domain": "teaching", "prompt": "homework", "validator": "homework", "export": "docx"},
    "comment": {"name": "学生评语", "domain": "teaching", "prompt": "comment", "validator": "comment", "export": "docx"},
    "score_analysis": {
        "name": "成绩分析", "domain": "teaching", "prompt": "score_analysis", "validator": "score", "export": "docx",
    },
    "teaching_research": {
        "name": "教研报告", "domain": "teaching", "prompt": "teaching_research", "validator": "research", "export": "docx",
    },
    "report_card": {"name": "成绩单模板", "domain": "teaching", "prompt": "report_card", "export": "xlsx"},
    # ---------------- 办公域 ----------------
    "notice": {"name": "通知公告", "domain": "office", "prompt": "notice", "validator": "notice", "export": "docx"},
    "work_plan": {"name": "工作计划", "domain": "office", "prompt": "work_plan", "validator": "plan", "export": "docx"},
    "work_summary": {"name": "工作总结", "domain": "office", "prompt": "work_summary", "validator": "summary", "export": "docx"},
    "speech": {"name": "发言稿", "domain": "office", "prompt": "speech", "validator": "speech", "export": "docx"},
    "meeting_minutes": {"name": "会议纪要", "domain": "office", "prompt": "meeting_minutes", "validator": "minutes", "export": "docx"},
    "class_meeting": {"name": "班会方案", "domain": "office", "prompt": "class_meeting", "validator": "class_meeting", "export": "docx"},
    "evaluation": {"name": "申报材料", "domain": "office", "prompt": "evaluation", "validator": "evaluation", "export": "docx"},
}

# 产物 ID（意图识别输出）→ 内部任务 ID
PRODUCT_TASKS: Dict[str, str] = {
    "lesson-plan": "lesson_plan",
    "ppt": "ppt",
    "exercise": "exercise",
    "animation": "animation",
    "student-analysis": "learning_analysis",
    "explainer": "explainer",
    "exam-paper": "exam_paper",
    "homework": "homework",
    "comment": "comment",
    "score-analysis": "score_analysis",
    "teaching-research": "teaching_research",
    "report-card": "report_card",
    "notice": "notice",
    "work-plan": "work_plan",
    "work-summary": "work_summary",
    "speech": "speech",
    "meeting-minutes": "meeting_minutes",
    "class-meeting": "class_meeting",
    "evaluation": "evaluation",
}

TASK_NAMES: Dict[str, str] = {pid: PRODUCERS[pid]["name"] for pid in PRODUCERS}

# 教学域任务集合（用于拆解/检索分域）
TEACHING_TASKS = {pid for pid, s in PRODUCERS.items() if s["domain"] == "teaching"}
OFFICE_TASKS = {pid for pid, s in PRODUCERS.items() if s["domain"] == "office"}

# 产物多选选项（教学/办公分组）
PRODUCT_OPTIONS: List[Dict[str, str]] = [
    {"id": pid, "name": s["name"], "domain": "教学" if s["domain"] == "teaching" else "办公"}
    for pid, s in PRODUCERS.items()
]


# ===========================================================================
# 4. LLM 抽象层（统一 invoke(messages) -> str 接口）
# ===========================================================================
# BaseLLM.invoke 会捕获调用异常并以该前缀回传，以便上层识别"模型调用失败"。
_LLM_ERR_MARK = "[LLM_ERROR]"


def is_llm_error(text: str) -> bool:
    """判断文本是否为 LLM 调用失败的回传占位内容。"""
    return bool(text) and str(text).strip().startswith(_LLM_ERR_MARK)


class BaseLLM:
    """模型适配基类。子类只需实现 `_raw_invoke(messages) -> str`。"""

    name: str = "base"

    def invoke(self, messages: List[Dict[str, str]]) -> str:
        try:
            return self._raw_invoke(messages)
        except Exception as e:
            log.error("LLM 调用失败：%s", e)
            return f"[LLM_ERROR] {e}"

    def _raw_invoke(self, messages: List[Dict[str, str]]) -> str:  # pragma: no cover
        raise NotImplementedError


class MockLLM(BaseLLM):
    """无 API Key 时的确定性 Mock，保证工作流离线可跑、可演示。"""

    name = "mock"

    def _raw_invoke(self, messages: List[Dict[str, str]]) -> str:
        combined, user_content = "", ""
        for m in messages:
            combined += m.get("content", "") + "\n"
            if m.get("role") == "user":
                user_content += m.get("content", "") + "\n"
        return _mock_generate(combined, user_content)


class OpenAICompatLLM(BaseLLM):
    """OpenAI 兼容后端（含通义/智谱/硅基流动等多数国内厂商）。
    支持显式超时 + 多模型回退：主模型调用失败/超时自动切换备用模型，避免单点卡死。"""

    def __init__(
        self,
        model: str,
        api_key: str,
        base_url: str = "",
        temperature: float = 0.4,
        timeout: int = 420,
    ):
        self.name = "openai-compat"
        self._candidates: List[Tuple[str, str, str]] = [(model, api_key, base_url)]
        self._temperature = temperature
        self._timeout = timeout

    def add_candidates(self, candidates: List[Tuple[str, str, str]]) -> None:
        """追加备用模型链（主模型在前，其余在后）。"""
        for c in candidates or []:
            if c and c not in self._candidates:
                self._candidates.append(c)

    def _raw_invoke(self, messages: List[Dict[str, str]]) -> str:
        last_err: Optional[Exception] = None
        for (model, api_key, base_url) in self._candidates:
            try:
                from openai import OpenAI

                client = OpenAI(
                    api_key=api_key,
                    base_url=base_url or None,
                    timeout=self._timeout,
                    max_retries=0,
                )
                resp = client.chat.completions.create(
                    model=model,
                    messages=messages,
                    temperature=self._temperature,
                    timeout=self._timeout,
                )
                text = (resp.choices[0].message.content or "").strip()
                if text:
                    return text
                last_err = RuntimeError("模型返回空内容")
            except Exception as e:
                last_err = e
                log.warning("模型 %s(%s) 调用失败，尝试备用模型：%s", model, base_url or "默认", e)
        log.error("所有可用模型调用均失败：%s", last_err)
        if last_err:
            raise last_err
        return ""


def build_llm() -> BaseLLM:
    """根据配置构建 LLM 实例。"""
    p = SETTINGS.provider
    if p == "mock":
        return MockLLM()
    if p in ("openai", "openai_compat", "dashscope", "zhipu"):
        if not SETTINGS.api_key:
            log.warning("未配置 API Key，回退到 Mock 模式。设置 AGENT_LLM_API_KEY 启用真实模型。")
            return MockLLM()
        llm = OpenAICompatLLM(
            model=SETTINGS.model,
            api_key=SETTINGS.api_key,
            base_url=SETTINGS.base_url,
            temperature=SETTINGS.temperature,
            timeout=SETTINGS.llm_timeout,
        )
        # 追加备用模型链（config.apply_model_to_workflow 会注入其他已配置模型）
        if SETTINGS.llm_candidates:
            llm.add_candidates(SETTINGS.llm_candidates)
        return llm
    log.warning("未知 provider=%s，回退 Mock。", p)
    return MockLLM()


def _mock_topic(user: str) -> str:
    m = re.search(r"《(.+?)》", user)
    return m.group(1) if m else "专题"


# 产物 Mock 模板：key = PRODUCERS 的 id（通过 system 中标记匹配）
_PRODUCT_MOCK: Dict[str, str] = {
    "exam_paper": (
        "## 试卷结构（满分150，120分钟）\n"
        "一、单项选择题（8题×5分）\n"
        "二、多项选择题（4题×5分）\n"
        "三、填空题（4题×5分）\n"
        "四、解答题（5题，共50分，含证明与综合应用）\n"
        "## 答案与解析\n"
        "每题附详细解析与评分标准。"
    ),
    "homework": (
        "## 分层作业设计\n"
        "基础巩固：3题（概念辨析、直接应用）\n"
        "能力提升：2题（综合运用、一题多解）\n"
        "拓展压轴：1题（跨章节综合、含参讨论）\n"
        "## 作业说明\n"
        "预计用时40分钟，附答案与易错点提示。"
    ),
    "comment": (
        "## 学生评语\n"
        "学习表现：课堂专注、作业认真，思维活跃，但在计算规范上需加强。\n"
        "优点：乐于思考、善于提问。\n"
        "建议：坚持错题整理，注重书写规范。\n"
        "寄语：保持热情，稳步提升。"
    ),
    "score_analysis": (
        "## 成绩分析报告\n"
        "总体情况：平均分/优秀率/及格率、分数段分布。\n"
        "小题得分：高频失分题与薄弱知识点。\n"
        "分层建议：尖子生拔高方向、中间层巩固重点、学困生补差方案。\n"
        "教学调整：下一阶段教学重点与训练安排。"
    ),
    "teaching_research": (
        "## 教研报告\n"
        "主题与背景：教研主题、课例与参与对象。\n"
        "课例观察：教学目标达成、重难点突破、师生活动记录。\n"
        "评课意见：亮点与改进建议。\n"
        "教研结论：共识、后续研究问题与行动计划。"
    ),
    "report_card": (
        "| 学号 | 姓名 | 平时成绩 | 期中成绩 | 期末成绩 | 总评 | 等级 | 评语 |\n"
        "| --- | --- | --- | --- | --- | --- | --- | --- |\n"
        "| 001 | 张三 | 85 | 90 | 88 | 87.6 | 优秀 | 学习认真，稳步提升 |\n"
        "| 002 | 李四 | 72 | 75 | 70 | 72.3 | 良好 | 基础尚可，需加强练习 |"
    ),
    "notice": (
        "## 通知\n"
        "时间：[填写时间] 地点：[填写地点] 对象：[填写对象]\n"
        "事项：说明事项背景、具体安排与要求。\n"
        "落款：[单位名称] 日期：[填写日期]"
    ),
    "work_plan": (
        "## 工作计划\n"
        "指导思想：围绕教学目标与立德树人。\n"
        "重点任务：教学进度、教研活动、学生管理、个人提升。\n"
        "实施安排：分阶段月度计划与检查节点。\n"
        "保障措施：资源支持与自我评价。"
    ),
    "work_summary": (
        "## 工作总结\n"
        "工作回顾：本学期主要教学与班级管理事项。\n"
        "成绩与亮点：成效数据与典型案例。\n"
        "问题与反思：不足及原因分析。\n"
        "改进方向：下学期工作计划与提升举措。"
    ),
    "speech": (
        "## 发言稿\n"
        "开场：问候与会者，点明主题。\n"
        "正文：分点阐述，结合实例与数据。\n"
        "结语：号召或期望，致谢。"
    ),
    "meeting_minutes": (
        "## 会议纪要\n"
        "会议时间/地点/主持人/参会人。\n"
        "议题一：讨论内容与结论。\n"
        "议题二：讨论内容与结论。\n"
        "决定与待办：任务、责任人、完成时限。"
    ),
    "class_meeting": (
        "## 班会方案\n"
        "主题：围绕学生成长与班级建设。\n"
        "目标：认知、情感、行为三维目标。\n"
        "流程：导入→活动→讨论→总结（40分钟）。\n"
        "准备：材料与分工。\n"
        "延伸：后续跟进与家校联动。"
    ),
    "evaluation": (
        "## 申报材料\n"
        "基本情况：姓名、教龄、任教科目、现聘岗位。\n"
        "教育教学业绩：教学成绩、公开课、课题、论文、竞赛指导。\n"
        "师德表现与班级管理成效。\n"
        "佐证材料清单与承诺声明。"
    ),
}


# 通用产物在 PROMPTS 提示词中的中文标记（用于 Mock 路由）
_PRODUCT_MOCK_MARKERS: Dict[str, str] = {
    "exam_paper": "命题专家",
    "homework": "设计分层作业",
    "comment": "撰写期末评语",
    "score_analysis": "成绩分析报告",
    "teaching_research": "教研报告",
    "report_card": "成绩单模板",
    "notice": "撰写一份规范的通知",
    "work_plan": "制定一份工作计划",
    "work_summary": "撰写一份工作总结",
    "speech": "撰写一份发言稿",
    "meeting_minutes": "整理一份会议纪要",
    "class_meeting": "设计一份班会方案",
    "evaluation": "评优/职称申报材料",
}


def _mock_generate(system: str, user: str) -> str:
    """离线 Mock：按 system 提示词阶段/产物标记路由，返回结构化内容。"""
    c = system + user

    # —— 阶段1：意图识别（期望严格 JSON）——
    if "高级意图识别器" in system:
        products: List[str] = []
        quantities: Dict[str, int] = {}
        products_explicit = False

        def pick(keyword: str, pid: str) -> None:
            nonlocal products_explicit
            if keyword.lower() in user.lower():
                products.append(pid)
                products_explicit = True

        for k, p in [
            ("教案", "lesson-plan"), ("教学设计", "lesson-plan"),
            ("ppt", "ppt"), ("课件", "ppt"),
            ("习题", "exercise"), ("题目", "exercise"), ("练习", "exercise"),
            ("动画", "animation"), ("演示", "animation"),
            ("学情", "student-analysis"),
            ("精讲", "explainer"), ("重难点", "explainer"),
            ("试卷", "exam-paper"), ("组卷", "exam-paper"), ("命题", "exam-paper"),
            ("作业", "homework"),
            ("评语", "comment"),
            ("成绩分析", "score-analysis"), ("成绩单", "score-analysis"),
            ("教研", "teaching-research"), ("听评课", "teaching-research"), ("评课", "teaching-research"),
            ("成绩单模板", "report-card"), ("成绩表", "report-card"),
            ("通知", "notice"), ("告家长书", "notice"),
            ("计划", "work-plan"),
            ("总结", "work-summary"), ("述职", "work-summary"),
            ("发言稿", "speech"), ("致辞", "speech"),
            ("会议纪要", "meeting-minutes"), ("会议记录", "meeting-minutes"),
            ("班会", "class-meeting"),
            ("评优", "evaluation"), ("职称", "evaluation"), ("申报", "evaluation"),
        ]:
            pick(k, p)
        # 去重保持顺序
        products = list(dict.fromkeys(products))
        if not products:
            products = ["lesson-plan"]
        # 数量：页/道
        m = re.search(r"(\d+)\s*页", user)
        if m:
            quantities["ppt"] = int(m.group(1))
        m = re.search(r"(\d+)\s*道", user)
        if m:
            quantities["exercise"] = int(m.group(1))

        # 学科推断
        if "数学" in user:
            subject = "数学"
        elif "物理" in user:
            subject = "物理"
        elif "化学" in user:
            subject = "化学"
        elif "生物" in user:
            subject = "生物"
        elif "英语" in user:
            subject = "英语"
        elif "古诗词" in user or "诗歌" in user or "语文" in user:
            subject = "语文"
        elif "历史" in user:
            subject = "历史"
        elif "地理" in user:
            subject = "地理"
        elif "政治" in user:
            subject = "政治"
        else:
            subject = ""
        # 课题提取
        topic_m = re.search(r"《(.+?)》", user)
        topic = topic_m.group(1) if topic_m else ""
        if not topic:
            t = re.search(r"(?:讲|备|做|写|生成|制作|设计|学|教学|复习)(?:一个|一份|一篇|一节|一节课|一些|关于)?([\u4e00-\u9fa5A-Za-z0-9·]{2,20})", user)
            topic = t.group(1) if t else ""
            topic = re.sub(r"^(?:一个|一份|一篇|一节|一节课|一些|关于)", "", topic)
            topic = re.sub(r"(教案|课件|PPT|习题|试卷|动画|学情分析|分析|精讲|课|通知|计划|总结|发言稿)$", "", topic)
        # 追问答案优先
        ans_topic = re.search(r'"topic"\s*:\s*"([^"]+)"', user)
        if ans_topic and ans_topic.group(1).strip():
            topic = ans_topic.group(1).replace("《", "").replace("》", "")
        ans_subject = re.search(r'"subject"\s*:\s*"([^"]+)"', user)
        if ans_subject and ans_subject.group(1).strip():
            subject = ans_subject.group(1)
        ans_grade = re.search(r'"grade"\s*:\s*"([^"]+)"', user)
        grade = ans_grade.group(1) if ans_grade else "高一"

        missing = []
        if not products_explicit:
            missing.append("products")
        # 学科仅在涉及教学类产物时才是关键字段（办公事务如通知/纪要无需学科）
        teaching_pids = {"lesson-plan", "ppt", "exercise", "animation", "student-analysis",
                         "explainer", "exam-paper", "homework", "comment", "score-analysis",
                         "teaching-research", "report-card"}
        if not subject and any(p in teaching_pids for p in products):
            missing.append("subject")
        if not topic:
            missing.append("topic")
        student_explicit = bool(re.search(r"中等|基础|薄弱|尖子|优等|学情", user))
        student_tags = "基础中等"
        if "薄弱" in user or "基础差" in user:
            student_tags = "基础薄弱"
        elif "尖子" in user or "优等" in user:
            student_tags = "尖子生"
        return json.dumps({
            "products": products,
            "productsExplicit": products_explicit,
            "quantities": quantities,
            "studentTags": [student_tags],
            "studentTagsExplicit": student_explicit,
            "subject": subject,
            "grade": grade,
            "topic": topic,
            "needsResearch": False,
            "needsStudentData": False,
            "confidence": 90,
            "missingInfo": missing,
            "priority": "high",
        }, ensure_ascii=False)

    # —— 阶段3a：教学任务拆解与重难点精讲 ——
    if "拆解为可执行子任务" in system:
        return (
            "## 任务拆解\n"
            "1. 制定教学目标（三维目标与课标对齐）\n"
            "2. 梳理教学重难点\n"
            "3. 设计教学过程（导入/新授/巩固/小结/作业）\n"
            "4. 生成配套资源（PPT/习题/动画）\n"
            "5. 分层作业设计\n"
            "## 重难点\n"
            "重点：核心概念与方法的理解应用。\n"
            "难点：抽象概念的建立与综合运用。\n"
            "## 精讲\n"
            "从直观感知→概念形成→变式练习三阶段突破难点，注重典型例题示范。"
        )

    # —— 阶段3b：办公任务拆解 ——
    if "办公任务拆解" in system:
        return (
            "## 办公任务拆解\n"
            "1. 明确事项对象、时间、范围与要求\n"
            "2. 梳理材料结构（引言/正文/结语）\n"
            "3. 填充事实、数据与佐证\n"
            "4. 核对格式规范与用语\n"
            "5. 落款、签发与存档"
        )

    # —— 阶段9：资源包汇总（需在产物分支之前判断，且用唯一标记避免与产物名冲突）——
    if "本次执行任务" in system:
        m = re.search(r"本次执行任务：(.+)", system)
        tasks_str = m.group(1).strip() if m else "教案"
        topic = _mock_topic(user)
        file_map = {
            "教案": "教案_《%s》.docx", "PPT课件": "课件_《%s》.pptx",
            "分层习题": "习题_《%s》.docx", "互动动画": "动画_《%s》.mp4",
            "学情分析": "学情报告_《%s》.docx", "重难点精讲": "精讲_《%s》.md",
            "试卷组卷": "试卷_《%s》.docx", "分层作业": "作业_《%s》.docx",
            "学生评语": "评语_《%s》.docx", "成绩分析": "成绩分析_《%s》.docx",
            "通知公告": "通知_%s.docx", "工作计划": "计划_%s.docx",
            "工作总结": "总结_%s.docx", "发言稿": "发言稿_%s.docx",
            "会议纪要": "会议纪要_%s.docx", "班会方案": "班会方案_%s.docx",
            "申报材料": "申报材料_%s.docx",
        }
        task_list = [t for t in tasks_str.replace("、", ",").split(",") if t.strip()]
        lines = ["资源包总览："]
        for t in task_list:
            t = t.strip()
            if t in file_map:
                fname = (file_map[t] % topic) if "%s" in file_map[t] else file_map[t].replace("%s", topic)
                lines.append(f"- {fname}：本次生成产物。")
        lines.append("建议使用顺序：" + "→".join(task_list) + "。")
        return "\n".join(lines)

    # —— 阶段4：各产物生成（按 system 提示词中的产物标记路由）——
    # 专用产物（教学核心 6 项）
    if "特级教师" in system:
        return (
            "## 教学目标\n"
            "知识与技能：理解核心概念，掌握基本方法。\n"
            "过程与方法：经历探索过程，渗透数学思想。\n"
            "情感态度：激发学习兴趣，培养理性精神。\n"
            "## 学情分析\n"
            "班级基础中等，已具备相应知识储备，对抽象概念理解不足。\n"
            "## 教学重难点\n"
            "重点：概念与方法掌握。难点：抽象理解与综合应用。\n"
            "## 教学过程\n"
            "导入：情境引入，激发兴趣。\n"
            "新授：概念归纳→例题示范→方法提炼。\n"
            "巩固：即时练习，辨析易错点。\n"
            "小结：梳理知识结构。\n"
            "作业：分层布置。\n"
            "## 板书设计\n"
            "主板书呈现知识框架，副板书呈现典型例题。\n"
            "## 分层作业\n"
            "基础巩固3题、能力提升2题、拓展压轴1题。\n"
            "## 教学反思\n"
            "关注学生理解程度，及时调整教学节奏。"
        )
    if "PPT 大纲" in system:
        pages = 12
        m = re.search(r"目标页数：(\d+)", c)
        if m:
            pages = int(m.group(1))
        return json.dumps([
            {"title": f"第{idx}页", "content": f"《{_mock_topic(user)}》内容要点", "note": "讲稿", "animation": "淡入"}
            for idx in range(1, pages + 1)
        ], ensure_ascii=False)
    if "分层习题" in system:
        count = 10
        m = re.search(r"数量：(\d+)", c)
        if m:
            count = int(m.group(1))
        n1 = max(1, count // 3)
        n2 = max(1, count // 3)
        n3 = max(1, count - n1 - n2)
        return (
            f"## 基础巩固（{n1}题）\n"
            "1. 基础概念判断题。\n解析：依据定义判断。考点：概念理解。易错点：概念混淆。\n"
            f"## 能力提升（{n2}题）\n"
            "1. 综合应用证明题。\n解析：逐步推理。考点：方法运用。易错点：步骤遗漏。\n"
            f"## 拓展压轴（{n3}题）\n"
            "1. 含参讨论综合题。\n解析：分类讨论。考点：综合能力。易错点：分类不全。"
        )
    if "互动动画脚本" in system:
        return json.dumps({
            "title": _mock_topic(user) + "动态演示",
            "type": "教学",
            "steps": ["展示核心对象", "动态变化过程", "标注关键节点"],
            "interaction": "拖动/点击参数，实时观察变化",
            "tech": "HTML5/SVG",
        }, ensure_ascii=False)
    if "学情分析报告" in system:
        return (
            "## 班级整体掌握情况\n"
            "班级整体中等，概念掌握较好，综合应用偏弱。\n"
            "## 高频易错点预判\n"
            "概念混淆、步骤不规范、审题不清。\n"
            "## 分层教学策略\n"
            "A层：综合提升；B层：巩固方法；C层：夯实基础。\n"
            "## 学困生辅导方案\n"
            "一对一错题归因，基础练习巩固。\n"
            "## 优生拓展方向\n"
            "跨章节综合与竞赛思维训练。"
        )
    if "精讲稿" in system:
        return (
            "## 精讲内容\n"
            "先讲概念本源，再给判定方法，最后典型例题三步示范。\n"
            "## 易错提醒\n"
            "强调'任意/存在'等关键量词辨析。"
        )
    # 通用产物（按提示词中的中文标记路由，system 为渲染后的中文提示词）
    for pid, marker in _PRODUCT_MOCK_MARKERS.items():
        if marker in system:
            return _PRODUCT_MOCK[pid]

    return f"《{_mock_topic(user)}》相关产物（Mock 生成）。"


# ===========================================================================
# 5. 知识库（教材知识图谱 + 校本资源）
# ===========================================================================
class KnowledgeBase:
    """加载教材知识图谱 kg_data.json，支持按学科/课题检索上下文。"""

    SUBJECT_ALIAS = {
        "数学": "数学", "math": "数学", "物理": "物理", "physics": "物理",
        "化学": "化学", "chemistry": "化学", "生物": "生物", "biology": "生物",
        "语文": "语文", "chinese": "语文", "英语": "英语", "english": "英语",
        "历史": "历史", "history": "历史", "地理": "地理", "geography": "地理",
        "政治": "政治", "politics": "政治", "思想政治": "政治",
    }

    def __init__(self, path: str = ""):
        self.data: Dict[str, Any] = {}
        self.path = path or self._default_path()
        self._load()

    @staticmethod
    def _default_path() -> str:
        base = os.path.dirname(os.path.abspath(__file__))
        candidates = [
            os.path.join(base, "..", "学生端", "textbooks", "knowledge-graph", "kg_data.json"),
            os.path.join(base, "textbooks", "knowledge-graph", "kg_data.json"),
            os.environ.get("AGENT_KG_PATH", ""),
        ]
        for c in candidates:
            if c and os.path.exists(c):
                return c
        return ""

    def _load(self) -> None:
        if not self.path or not os.path.exists(self.path):
            log.warning("未找到知识图谱文件，知识检索将使用内置常识。")
            return
        try:
            with open(self.path, "r", encoding="utf-8") as f:
                self.data = json.load(f)
            log.info("知识图谱已加载：%s（%d 个学科）", self.path, len(self.data))
        except Exception as e:
            log.warning("知识图谱加载失败，降级为内置常识：%s", e)
            self.data = {}

    def retrieve(self, subject: str, topic: str) -> Dict[str, Any]:
        subj = self.SUBJECT_ALIAS.get(subject, subject)
        nodes: List[Dict] = []
        edges: List[Dict] = []
        if subj in self.data:
            nodes = self.data[subj].get("nodes", [])
            edges = self.data[subj].get("edges", [])
        matched = []
        for n in nodes:
            blob = " ".join(str(n.get(k, "")) for k in ("name", "description", "keywords", "module"))
            if topic and topic in blob:
                matched.append(n)
        if not matched:
            matched = [n for n in nodes if n.get("type") == "Chapter"][:8] or nodes[:8]
        return {
            "subject": subj,
            "topic": topic,
            "matched_nodes": matched,
            "edge_count": len(edges),
            "session": self._format_context(subj, matched),
        }

    @staticmethod
    def _format_context(subj: str, matched: List[Dict]) -> str:
        lines = [f"学科：{subj}"]
        for n in matched[:8]:
            desc = n.get("description", "")
            diff = n.get("difficulty", 0)
            imp = n.get("importance", 0)
            lines.append(f"- {n.get('name', '')}（难度{diff}/重要度{imp}）：{desc}")
        return "\n".join(lines) if matched else f"学科：{subj}（暂无知识图谱数据）"


# ===========================================================================
# 6. 状态定义（LangGraph State）
# ===========================================================================
class LessonState(TypedDict, total=False):
    """图的全局状态。所有节点都可读写，用于跨子工作流共享上下文。"""

    raw_input: str
    history: List[Dict[str, str]]  # 多轮问答记录
    parsed: Dict[str, Any]
    missing_info: List[str]
    pending_questions: List[Dict[str, Any]]
    info_complete: bool
    globals_: Dict[str, Any]
    tasks: List[str]  # 智能决策后的执行清单
    research: str  # 联网搜索资料
    kb_session: Dict[str, Any]
    breakdown: str
    key_points: str
    explanation: str
    errors: List[str]
    final: Dict[str, Any]

    # 各产物结果字段（与 PRODUCERS 的 id 一一对应，LangGraph 要求显式声明）
    lesson_plan: Dict[str, Any]
    ppt: Dict[str, Any]
    exercise: Dict[str, Any]
    animation: Dict[str, Any]
    learning_analysis: Dict[str, Any]
    explainer: Dict[str, Any]
    exam_paper: Dict[str, Any]
    homework: Dict[str, Any]
    comment: Dict[str, Any]
    score_analysis: Dict[str, Any]
    teaching_research: Dict[str, Any]
    report_card: Dict[str, Any]
    notice: Dict[str, Any]
    work_plan: Dict[str, Any]
    work_summary: Dict[str, Any]
    speech: Dict[str, Any]
    meeting_minutes: Dict[str, Any]
    class_meeting: Dict[str, Any]
    evaluation: Dict[str, Any]


# 各产物在 State 中的字段名（与任务 ID 一致，运行时动态读写）
STATE_FIELDS: List[str] = list(PRODUCERS.keys())


# ===========================================================================
# 7. 提示词模板
# ===========================================================================
PROMPTS: Dict[str, str] = {
    "intent": (
        "你是教学智能体的高级意图识别器。分析教师的请求（可能包含备课、作业、考试、班级管理、"
        "学校办公事务等），解析出全部产物需求，输出严格 JSON。\n"
        "JSON 字段：products(数组，可选 lesson-plan/ppt/exercise/animation/student-analysis/explainer/"
        "exam-paper/homework/comment/score-analysis/notice/work-plan/work-summary/speech/"
        "meeting-minutes/class-meeting/evaluation)、"
        "productsExplicit(布尔，教师是否明确点名了所需产物，未点名则为 false)、"
        "quantities(对象，如 {\"ppt\":15,\"exercise\":10})、studentTags(数组，学情标签)、"
        "studentTagsExplicit(布尔，教师是否说明了学情)、"
        "subject、grade、topic、needsResearch、needsStudentData、confidence(0-100)、"
        "missingInfo(数组，仅含影响执行的关键缺失字段 subject/grade/topic/products；"
        "其中 subject 仅在涉及教学类产物（如教案/PPT/习题/试卷）时才关键，纯办公事务"
        "（通知/纪要/发言稿/班会等）无需学科；学情/数量等可选信息不要列入)、priority。\n"
        "识别规则：教案→lesson-plan；PPT/课件→ppt；习题/题目→exercise；动画/演示→animation；"
        "学情→student-analysis；精讲/重难点→explainer；试卷/组卷/命题→exam-paper；作业→homework；"
        "评语→comment；成绩分析→score-analysis；教研/听评课/评课→teaching-research；成绩单/成绩表→report-card；"
        "通知/告家长书→notice；计划→work-plan；总结/述职→work-summary；"
        "发言稿/致辞→speech；会议纪要/记录→meeting-minutes；班会→class-meeting；评优/职称/申报→evaluation。\n"
        "教师未指明任何产物时 products 给默认 [\"lesson-plan\"] 且 productsExplicit=false。"
        "只输出 JSON，无解释。\n"
        "教师需求：{input}"
    ),
    "decompose_teaching": (
        "你负责把教学备课需求拆解为可执行子任务，并给出重难点精讲。\n"
        "学科：{subject}，课题：{topic}，学情：{student_tags}。\n"
        "知识底料：{kb}\n联网检索资料：{research}\n"
        "输出分三部分：`## 任务拆解`（编号列表）、`## 重难点`（重点/难点）、`## 精讲`（逐字稿级讲解）。"
    ),
    "decompose_office": (
        "你负责把办公事务需求拆解为可执行子任务。\n"
        "任务内容：{topic}，事项：{raw}\n"
        "输出分三部分：`## 办公任务拆解`（编号列表）、`## 材料要点`（必须包含的事实/数据/结构）、"
        "`## 注意事项`（格式、用语、时限）。"
    ),
    # ---- 教学域产物 ----
    "lesson_plan": (
        "你是特级教师，按新课标规范撰写教案。\n"
        "学科：{subject}，课题：{topic}，学情：{student_tags}，重难点：{key_points}\n"
        "必须包含：教学三维目标、学情分析、教学重难点、教学过程（导入/新授/巩固/小结/作业）、"
        "板书设计、分层作业、教学反思。输出结构化 Markdown。"
    ),
    "ppt": (
        "你是资深课件设计专家，基于教案为《{topic}》设计一份结构完整、图文并茂、可直接上课使用的课件。\n"
        "学科：{subject}，目标页数：{pages}，教案要点：{lesson_plan}\n"
        "\n"
        "【输出要求】输出严格 JSON 数组，每页对象结构如下：\n"
        '{"type": "title|section|content|example|practice|summary|end", "title": "页面标题", '
        '"bullets": ["要点1(具体知识内容，禁止『讲解XX』等空话)", "要点2"],\n'
        '"note": "讲稿/教师话术", "chart": "可选: bar|line|pie 当本页涉及数据对比时，生成该图表",\n'
        '"layout": "可选: cover|agenda|content|two-col|image-right|image-left|example|summary"}\n'
        "\n"
        "【页面结构强制要求】\n"
        "1. 第1页 type=title 封面：主标题《{topic}》+ 副标题(学科+年级+授课人占位)。\n"
        "2. 第2页 type=section 学习目标（bullets 用可测量的行为动词写3-4条目标）。\n"
        "3. 第3页 type=section 目录/内容导航。\n"
        "4. 中间 6-9 页 type=content/example 知识讲授：每页 bullets 必须写具体知识内容原话（定义、性质、公式、推论的原话），"
        "例题页给出完整题干+分步规范解答（数学公式用 LaTeX $...$）。\n"
        "5. 至少1页 type=example 例题精讲（含完整解答步骤）。\n"
        "6. 至少1页 type=practice 课堂练习（含答案与易错提示）。\n"
        "7. 最后1页 type=summary 课堂小结，1页 type=end 作业布置页。\n"
        "\n"
        "【质量要求】\n"
        "- 总页数不少于 {pages} 页；每页 bullets 2-5 条，都写具体内容，禁止占位文字。\n"
        "- 涉及数据对比、统计、函数图像、结构示意时，给对应页设置 chart 字段（bar/line/pie），"
        "系统会自动为该页生成数据图表，实现图文并茂。\n"
        "- 只输出 JSON 数组，不要任何解释或代码块标记。"
    ),
    "exercise": (
        "生成分层习题与解析。\n学科：{subject}，课题：{topic}，数量：{count}，重难点：{key_points}\n"
        "按基础巩固/能力提升/拓展压轴三梯度，覆盖选择/填空/解答，每题给出解题步骤、考点、易错点。"
    ),
    "animation": (
        "为课题设计课堂互动动画脚本。\n课题：{topic}，重难点：{key_points}\n"
        "输出 JSON：title、type(数学/物理/化学/生物等)、steps(演示步骤数组)、"
        "interaction(交互逻辑)、tech(建议实现：SVG/Manim/HTML5)。"
    ),
    "learning_analysis": (
        "基于班级学情数据生成学情分析报告。\n课题：{topic}，学情数据：{student_data}\n"
        "必须覆盖：班级整体掌握情况、高频易错点预判、分层教学策略、学困生辅导方案、优生拓展方向。"
    ),
    "explainer": (
        "围绕课题撰写重难点精讲稿。\n课题：{topic}，重难点：{key_points}\n"
        "输出：精讲内容（逐字稿风格）、易错提醒、追问设计。"
    ),
    "exam_paper": (
        "你是学科命题专家，生成一份标准试卷。\n"
        "学科：{subject}，课题：{topic}，年级：{grade}，学情：{student_tags}，重难点：{key_points}\n"
        "必须包含：试卷标题与考试说明、题型结构（选择/填空/解答）、每题分值、难度梯度、"
        "答案与解析、评分标准。"
    ),
    "homework": (
        "设计分层作业。\n学科：{subject}，课题：{topic}，学情：{student_tags}，重难点：{key_points}\n"
        "按基础巩固/能力提升/拓展压轴三层设计，含题量、预计用时、答案与易错点提示。"
    ),
    "comment": (
        "你是班主任，为学生撰写期末评语。\n"
        "学生情况：{topic}（若未提供姓名用'该生'），学情：{student_tags}\n"
        "评语需含：学习表现、闪光点、改进建议、寄语，语气真诚具体。"
    ),
    "score_analysis": (
        "你是数据分析师，基于成绩数据输出成绩分析报告。\n"
        "科目：{subject}，数据情况：{raw}\n"
        "必须覆盖：总体统计（均分/优秀率/及格率/分数段）、小题得分与薄弱知识点、"
        "分层教学建议、教学调整建议。"
    ),
    "teaching_research": (
        "你是教研组长，撰写一份教研报告。\n"
        "教研主题：{topic}，参与情况：{raw}，学科：{subject}\n"
        "必须包含：主题与背景、课例观察记录、评课意见（亮点与建议）、"
        "教研结论与后续行动计划。"
    ),
    "report_card": (
        "你是教务管理员，生成成绩单模板。\n"
        "科目：{subject}，班级/年级：{grade}，说明：{raw}\n"
        "输出 Markdown 表格：学号/姓名/各次成绩/总评/等级/评语，含表头与示例行，"
        "格式规范可直接填充。"
    ),
    # ---- 办公域产物 ----
    "notice": (
        "你是学校办公室文员，撰写一份规范的通知。\n"
        "事项：{topic}，背景信息：{raw}\n"
        "必须包含：标题、收文对象、时间、地点、事项与要求、落款与日期。用语规范简洁。"
    ),
    "work_plan": (
        "你是教学管理者，制定一份工作计划。\n"
        "主题：{topic}，背景：{raw}，学科/岗位：{subject}\n"
        "必须包含：指导思想、重点任务、分阶段实施安排、保障措施。"
    ),
    "work_summary": (
        "你是办公室文员，撰写一份工作总结。\n"
        "主题：{topic}，情况：{raw}\n"
        "必须包含：工作回顾、成绩与亮点、问题与反思、改进方向。"
    ),
    "speech": (
        "你是演讲稿撰稿人，撰写一份发言稿。\n"
        "场合：{topic}，对象：{raw}\n"
        "必须包含：开场问候、主题阐述（分点，含实例）、结语致谢。语言得体有感染力。"
    ),
    "meeting_minutes": (
        "你是会议记录员，整理一份会议纪要。\n"
        "会议主题：{topic}，记录素材：{raw}\n"
        "必须包含：会议基本信息（时间/地点/主持/参会）、议题与讨论、决议事项、待办与责任人时限。"
    ),
    "class_meeting": (
        "你是班主任，设计一份班会方案。\n"
        "班会主题：{topic}，班级情况：{raw}\n"
        "必须包含：主题、目标、流程（导入→活动→讨论→总结）、所需准备、延伸与家校联动。"
    ),
    "evaluation": (
        "你是材料撰写专家，撰写评优/职称申报材料。\n"
        "申报类型：{topic}，个人情况：{raw}，学科：{subject}\n"
        "必须包含：基本情况、教育教学业绩（成绩/公开课/课题/论文/竞赛）、师德与班级管理、"
        "佐证材料清单与承诺。"
    ),
    "aggregate": (
        "你是质检员，汇总本次实际执行的产物为资源包清单。\n"
        "本次执行任务：{tasks}\n产物：{results}\n"
        "输出：资源包总览（文件名、内容简介、建议使用顺序）。"
    ),
}


# ===========================================================================
# 8. LLM 快捷函数
# ===========================================================================
_LLM_INST: Optional[BaseLLM] = None
# 模型配置并发锁：后端 config.py 改写 SETTINGS / 重置 _LLM_INST 时持有同一把锁，
# 避免"写 SETTINGS"与"读 SETTINGS 构建 LLM"交错导致读到半成品配置。
_MODEL_LOCK = threading.RLock()


def _llm() -> BaseLLM:
    global _LLM_INST
    with _MODEL_LOCK:
        if _LLM_INST is None:
            _LLM_INST = build_llm()
        return _LLM_INST


def _chat(system: str, user: str) -> str:
    return _llm().invoke([
        {"role": "system", "content": system},
        {"role": "user", "content": user},
    ])


def _safe_json(text: str) -> Any:
    text = text.strip()
    m = re.search(r"```(?:json)?\s*(.*?)\s*```", text, re.S)
    if m:
        text = m.group(1)
    try:
        return json.loads(text)
    except Exception:
        for st, en in ((text.find("{"), text.rfind("}")), (text.find("["), text.rfind("]"))):
            if st != -1 and en > st:
                try:
                    return json.loads(text[st: en + 1])
                except Exception:
                    continue
        raise ValueError(f"无法解析 JSON：{text[:200]}")


def _fmt_prompt(tpl: str, fmt: Dict[str, Any]) -> str:
    """安全渲染提示词模板：只替换 fmt 中已知的 {key} 占位符，
    模板里其余花括号（如内嵌 JSON 示例）原样保留，杜绝 KeyError/格式化异常。"""
    def _sub(m: "re.Match") -> str:
        key = m.group(1)
        if key in fmt:
            return str(fmt[key])
        return m.group(0)
    return re.sub(r"\{([A-Za-z_][A-Za-z0-9_]*)\}", _sub, tpl)


# ===========================================================================
# 9. 节点函数
# ===========================================================================
def input_node(state: LessonState) -> LessonState:
    return {"raw_input": state.get("raw_input", "").strip()}


def intent_router(state: LessonState) -> LessonState:
    raw = state.get("raw_input", "")
    if not raw:
        return {"parsed": {}, "missing_info": ["raw_input"], "info_complete": False}
    full_input = raw
    for h in state.get("history", []):
        full_input += f"\n[{h.get('role')}] {h.get('content')}"
    text = _chat(PROMPTS["intent"], full_input)
    try:
        parsed = _safe_json(text)
        if not isinstance(parsed, dict):
            raise ValueError("意图结果非对象")
    except Exception as e:
        log.warning("意图解析失败(%s)，启用本地规则解析兜底。", e)
        try:
            parsed = local_intent_parse(raw)
        except Exception as e2:
            log.warning("本地意图解析也失败(%s)，使用最低默认值。", e2)
            parsed = {"products": ["lesson-plan"], "subject": "", "topic": raw,
                      "missingInfo": ["subject", "topic"], "productsExplicit": False}
    missing = parsed.get("missingInfo", []) or []
    if not parsed.get("products"):
        parsed["products"] = ["lesson-plan"]
    return {
        "parsed": parsed,
        "missing_info": missing,
        "info_complete": len(missing) == 0,
    }


def is_complete(state: LessonState) -> str:
    """是否还需要提问：关键字段缺失或产物未明确点名 → 问；否则直接执行。
    force_complete=True（追问轮次用尽）时直接放行，带默认值强制执行。"""
    if state.get("force_complete"):
        return "assign_globals"
    parsed = state.get("parsed", {})
    if not state.get("info_complete"):
        return "ask"
    if not parsed.get("productsExplicit", True):
        return "ask"
    return "assign_globals"


FIELD_QUESTIONS: Dict[str, str] = {
    "raw_input": "请描述您的需求（事项、对象、时间等）。",
    "topic": "请问本次事项的课题/主题是什么？",
    "subject": "请问是哪个学科/领域？",
    "grade": "请问是哪个年级？",
}


def ask_questions(state: LessonState) -> LessonState:
    """生成追问：产物多选确认 + 缺失字段补充 + 学情补充。结构化便于前端渲染。"""
    parsed = state.get("parsed", {})
    missing = list(state.get("missing_info", [])) or []
    questions: List[Dict[str, Any]] = []

    if "products" in missing or not parsed.get("productsExplicit", True):
        teaching = [o["name"] for o in PRODUCT_OPTIONS if o["domain"] == "教学"]
        office = [o["name"] for o in PRODUCT_OPTIONS if o["domain"] == "办公"]
        questions.append({
            "field": "products",
            "question": "请选择需要生成的产物（可多选）。\n【教学】" + "、".join(teaching) + "\n【办公】" + "、".join(office),
            "multiSelect": True,
            "options": [o["name"] for o in PRODUCT_OPTIONS],
        })

    for f in missing:
        if f == "products":
            continue
        questions.append({
            "field": f,
            "question": FIELD_QUESTIONS.get(f, f"请补充：{f}"),
            "multiSelect": False,
            "options": [],
        })

    if questions and not parsed.get("studentTagsExplicit", False) and any(
        t in TEACHING_TASKS for t in (parsed.get("products") or [])
    ):
        questions.append({
            "field": "studentTags",
            "question": "请问班级基础水平如何？（基础薄弱/基础中等/尖子生）",
            "multiSelect": False,
            "options": [],
        })

    if not questions:
        questions.append({
            "field": "other",
            "question": "请问还有其他需要补充的信息吗？",
            "multiSelect": False,
            "options": [],
        })
    return {"pending_questions": questions}


def assign_globals(state: LessonState) -> LessonState:
    """变量赋值：写入全局参数，并根据意图识别结果做智能决策生成执行清单 tasks。"""
    parsed = state.get("parsed", {})
    products = parsed.get("products") or ["lesson-plan"]
    tasks = []
    for p in products:
        t = PRODUCT_TASKS.get(p)
        if t and t in PRODUCERS and t not in tasks:
            tasks.append(t)
    if not tasks:
        tasks = ["lesson_plan"]
    globals_ = {
        "subject": parsed.get("subject", "数学"),
        "grade": parsed.get("grade", ""),
        "topic": parsed.get("topic", ""),
        "student_tags": "、".join(parsed.get("studentTags", []) or ["基础中等"]),
        "quantities": parsed.get("quantities", {}),
        "products": products,
        "tasks": tasks,
        "needsResearch": parsed.get("needsResearch", False),
        "needsStudentData": parsed.get("needsStudentData", False),
    }
    return {"globals_": globals_, "tasks": tasks}


def _should_run(state: LessonState, task_id: str) -> bool:
    return task_id in state.get("tasks", [])


# 知识库单例缓存：避免每个任务重复加载 kg_data.json（大文件 IO 是实测瓶颈之一）
_KB_CACHE: Dict[str, "KnowledgeBase"] = {}


def _get_kb_cached(path: str = "") -> "KnowledgeBase":
    key = path or "__default__"
    if key not in _KB_CACHE:
        _KB_CACHE[key] = KnowledgeBase(path=path)
    return _KB_CACHE[key]


def _retrieve_kb_rag(subject: str, topic: str) -> Dict[str, Any]:
    """优先走 RAG 混合检索（app.kb.retrieve），失败/不可用时回退旧 KnowledgeBase。

    说明：workflow 与 app.kb 存在相互引用（app.kb 顶层 import workflow），
    故此处采用函数内延迟导入，避免顶层循环导入导致启动失败。
    """
    try:
        from app import kb as _kb

        result = _kb.retrieve(subject or "", topic or "")
        if isinstance(result, dict):
            return result
    except Exception as e:
        log.warning("RAG 检索入口不可用，回退旧 KnowledgeBase：%s", e)
    # 回退：旧子串检索，保证服务可用
    kb = _get_kb_cached()
    old = kb.retrieve(subject or "", topic or "")
    if isinstance(old, dict):
        old["engine"] = "legacy"
    return old if isinstance(old, dict) else {"session": "", "matched_nodes": [], "engine": "legacy"}


def retrieve_knowledge(state: LessonState) -> LessonState:
    g = state.get("globals_", {})
    session = _retrieve_kb_rag(g.get("subject", ""), g.get("topic", ""))
    return {"kb_session": session}


# 内置学科资料：联网不可用时的离线兜底
FALLBACK_RESEARCH: Dict[str, str] = {
    "数学": "新课标要求：函数单调性属于必修一核心内容，要求会用定义证明简单函数的单调性，"
           "能结合图像判断单调区间，掌握复合函数单调性的判定，高考以解答题考查定义法证明。",
    "物理": "新课标要求：结合实例理解力与运动关系，能用牛顿运动定律分析问题，"
            "注重科学探究与实验素养。",
    "化学": "新课标要求：理解化学反应原理与物质转化，注重宏微结合与证据推理核心素养。",
}


def _fetch_page_text(url: str, max_chars: int = 1500, timeout: int = 8) -> str:
    """抓取网页正文（去标签纯文本），用于搜索结果深读。"""
    try:
        import requests

        resp = requests.get(url, timeout=timeout, headers={
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36",
        })
        if resp.status_code != 200:
            return ""
        html = resp.text
        # 去脚本/样式/标签
        html = re.sub(r"(?is)<(script|style|noscript)[^>]*>.*?</\1>", " ", html)
        text = re.sub(r"(?s)<[^>]+>", " ", html)
        text = re.sub(r"&[a-z]+;", " ", text)
        text = re.sub(r"\s+", " ", text).strip()
        return text[:max_chars]
    except Exception:
        return ""


def _web_search(query: str, limit: int = 5, timeout: float = 25.0) -> str:
    """联网搜索 + 深读：先取搜索结果，再打开 Top2 链接读取网页正文。
    整体限时，超时返回空串（走内置资料兜底），避免网络黑洞阻塞主流程。"""

    def _do() -> str:
        from duckduckgo_search import DDGS

        with DDGS() as ddgs:
            hits = list(ddgs.text(query, max_results=limit))
        lines = []
        for h in hits:
            if h.get("title"):
                lines.append(f"- {h.get('title', '')}：{h.get('body', '')}（{h.get('href', '')}）")
        # 深读 Top2 网页正文
        read = 0
        for h in hits:
            if read >= 2:
                break
            url = h.get("href", "")
            if not url.startswith("http"):
                continue
            body = _fetch_page_text(url)
            if len(body) > 200:
                lines.append(f"【网页正文·{h.get('title', '')[:30]}】{body}")
                read += 1
        return "\n".join(lines)

    try:
        with ThreadPoolExecutor(max_workers=1) as ex:
            return ex.submit(_do).result(timeout=timeout)
    except Exception as e:
        log.warning("联网搜索不可用（将使用内置资料）：%s", e)
        return ""


def web_search_node(state: LessonState) -> LessonState:
    g = state.get("globals_", {})
    subject = g.get("subject", "")
    topic = g.get("topic", "")
    research = ""
    if g.get("needsResearch"):
        research = _web_search(f"{subject} {topic} 新课标 教学要求 高频考点 易错点")
    if not research:
        research = FALLBACK_RESEARCH.get(subject, f"{subject}《{topic}》：建议参考新课标与教材教参把握要求。")
    return {"research": research}


def task_decompose(state: LessonState) -> LessonState:
    """任务拆解：按任务域分流（教学域/办公域）。"""
    g = state.get("globals_", {})
    tasks = state.get("tasks", [])
    kb = state.get("kb_session", {}).get("session", "")
    research = state.get("research", "")

    def _clean(text: str) -> str:
        # LLM 失败时返回的是 [LLM_ERROR] 占位，不能注入下游提示词/精讲稿
        if is_llm_error(text):
            log.warning("任务拆解 LLM 调用失败，置空以免污染下游提示词")
            return ""
        return text

    if any(t in TEACHING_TASKS for t in tasks):
        text = _chat(
            _fmt_prompt(PROMPTS["decompose_teaching"], {
                "subject": g.get("subject", ""), "topic": g.get("topic", ""),
                "student_tags": g.get("student_tags", ""), "kb": kb, "research": research,
            }),
            "请拆解教学任务并精讲重难点。",
        )
        text = _clean(text)
        return {"breakdown": text, "key_points": text, "explanation": text}
    text = _chat(
        _fmt_prompt(PROMPTS["decompose_office"], {
            "topic": g.get("topic", ""), "raw": state.get("raw_input", ""),
        }),
        "请拆解办公任务。",
    )
    text = _clean(text)
    return {"breakdown": text, "key_points": text, "explanation": text}


# ===========================================================================
# 9.5 产物质量硬性标准（Quality Gate）
# ===========================================================================
# 每类产物声明：min_chars 字数下限、required_sections 必备章节（patterns 任一命中）、
# placeholder_patterns 空泛表述黑名单、depth_checks 深度检查（正则计数下限）、
# min_slides/min_questions 结构性下限。不合格产物必须重新生成；
# 达到最大轮次仍不合格 → ok=False，导出层拒绝落盘，任务层显式报错。
QUALITY_STANDARDS: Dict[str, Dict[str, Any]] = {
    "lesson_plan": {
        "min_chars": 3000,
        "required_sections": {
            "教学目标": ["教学目标", "学习目标", "素养目标"],
            "学情分析": ["学情分析", "学情"],
            "教学重难点": ["教学重点", "教学难点", "重难点"],
            "教学过程": ["教学过程", "教学环节", "导入"],
            "教师活动话术": ["教师活动", "师：", "教师提问"],
            "学生活动预设": ["学生活动", "生：", "预设回答"],
            "板书设计": ["板书"],
            "分层作业": ["分层作业", "作业设计", "课后作业"],
            "教学反思": ["教学反思", "反思预设"],
        },
        "depth_checks": [
            ("教学环节数量", r"环节[一二三四五六七八\d]|第[一二三四五六七八\d]+[环步]", 4),
            ("教师提问话术", r"师：|教师[：:]", 3),
            ("学生预设回答", r"生：|预设[：:]", 2),
        ],
    },
    "ppt": {
        "min_chars": 0, "min_slides": 12,
        "required_sections": {
            "学习目标": ["学习目标", "教学目标", "目标"],
            "知识讲授": ["定义", "概念", "性质", "定理", "公式", "原理", "规律", "特点"],
            "例题": ["例题", "例", "解析", "解答"],
            "练习": ["练习", "课堂练习", "随堂", "检测"],
            "小结": ["小结", "总结", "归纳"],
        },
    },
    "exercise": {
        "min_chars": 1200,
        "required_sections": {
            "基础层": ["基础巩固", "基础"],
            "提升层": ["能力提升", "提升"],
            "拓展层": ["拓展", "压轴"],
            "解析": ["解析", "解答"],
        },
        "depth_checks": [("题目数量", r"(?m)^\s*\d+[\.．、]", 6)],
    },
    "animation": {"min_chars": 300, "required_sections": {"演示步骤": ["步骤", "step", "演示"], "交互": ["交互", "interaction", "拖动", "点击"]}},
    "learning_analysis": {
        "min_chars": 800,
        "required_sections": {
            "整体掌握": ["掌握情况", "整体"],
            "易错点": ["易错", "薄弱", "高频"],
            "分层策略": ["分层", "策略"],
            "学困生辅导": ["学困生", "辅导", "帮扶"],
            "优生拓展": ["优生", "拓展", "拔高"],
        },
    },
    "explainer": {"min_chars": 1500, "required_sections": {"精讲内容": ["精讲", "讲解", "推导"], "易错提醒": ["易错", "注意", "误区"]}},
    "exam_paper": {
        "min_chars": 1500,
        "required_sections": {
            "题型结构": ["选择", "填空", "解答"],
            "分值": ["分值", "满分", "总分"],
            "答案解析": ["答案", "解析"],
            "评分标准": ["评分", "得分"],
        },
        "depth_checks": [("题目数量", r"(?m)^\s*\d+[\.．、]", 8)],
    },
    "homework": {"min_chars": 600, "required_sections": {"分层结构": ["基础", "提升", "拓展"], "答案提示": ["答案", "解析", "易错"]}},
    "comment": {"min_chars": 300, "required_sections": {"学习表现": ["学习表现", "表现"], "建议": ["建议"], "寄语": ["寄语", "期望"]}},
    "score_analysis": {"min_chars": 800, "required_sections": {"总体统计": ["平均分", "及格率", "优秀率", "分数段"], "薄弱点": ["薄弱", "失分"], "建议": ["建议", "调整"]}},
    "teaching_research": {"min_chars": 1000, "required_sections": {"主题背景": ["主题", "背景"], "课例观察": ["课例", "观察", "课堂"], "评课意见": ["评课", "亮点", "建议"], "结论行动": ["结论", "行动", "后续"]}},
    "report_card": {"min_chars": 200, "required_sections": {"表格": ["|", "学号", "姓名"]}},
    "notice": {"min_chars": 250, "required_sections": {"标题/对象": ["通知", "对象"], "时间": ["时间"], "地点": ["地点"], "事项要求": ["事项", "要求", "安排"], "落款": ["落款", "日期", "单位"]}},
    "work_plan": {"min_chars": 800, "required_sections": {"指导思想": ["指导思想", "总体要求", "思路"], "重点任务": ["重点任务", "任务", "目标"], "实施安排": ["安排", "阶段", "进度"], "保障措施": ["保障", "措施"]}},
    "work_summary": {"min_chars": 800, "required_sections": {"工作回顾": ["回顾", "工作"], "成绩亮点": ["成绩", "亮点", "成效"], "问题反思": ["问题", "不足", "反思"], "改进方向": ["改进", "方向", "计划"]}},
    "speech": {"min_chars": 600, "required_sections": {"开场": ["开场", "问候", "尊敬", "各位"], "主体": ["主题", "第一", "首先"], "结语": ["结语", "致谢", "谢谢", "祝愿"]}},
    "meeting_minutes": {"min_chars": 400, "required_sections": {"基本信息": ["时间", "地点", "主持", "参会"], "议题讨论": ["议题", "讨论"], "决议待办": ["决议", "决定", "待办", "责任"]}},
    "class_meeting": {"min_chars": 600, "required_sections": {"主题": ["主题"], "目标": ["目标"], "流程": ["流程", "导入", "活动"], "准备": ["准备", "分工"], "延伸": ["延伸", "跟进", "家校"]}},
    "evaluation": {"min_chars": 800, "required_sections": {"基本情况": ["基本情况", "姓名", "教龄"], "业绩": ["业绩", "成绩", "课题", "论文"], "师德": ["师德", "班级管理"], "佐证": ["佐证", "材料清单", "承诺"]}},
}

# 通用空泛/占位表述黑名单（多字直接匹配）。
# 说明：移除了"同理"（"同理可得"是数学规范表述）、"举例说明"（正常教学活动指令）、
# 单字"略"（"策略/忽略/省略号"易误伤）——误判会触发整篇重写、浪费多轮调用。
GENERIC_PLACEHOLDERS: List[str] = [
    # 明确的占位 / 空话
    "（待补充）", "(待补充)", "待补充", "TODO", "占位", "placeholder",
    "此处讲解", "此处介绍", "此处说明", "此处省略", "以此类推",
    "详见教材", "参考课本",
    # 常见未替换占位符（此前漏检，导致空壳产物拿 100 分通过）
    "202X", "20XX", "____", "[教师姓名]", "（填写", "(填写", "待填写",
]

# 占位符形状正则（难以穷举为固定字符串的模式）
PLACEHOLDER_REGEXES: List[Tuple[str, str]] = [
    (r"X{2,}", "XX 占位"),    # "XX门课程""申请XX余次"等未替换占位
    (r"_{3,}", "下划线填空线"),  # "____（填写…）"
]


def _has_placeholder(text: str, p: str) -> bool:
    if not p:
        return False
    if len(p) >= 2:
        return p in text
    return re.search(r"(^|[^一-龥])" + re.escape(p) + r"([^一-龥]|$)", text) is not None


def validate_product(pid: str, content: str, parsed: Any = None, req: Optional[Dict[str, Any]] = None) -> Dict[str, Any]:
    """产物质量硬校验。返回 {passed, score, issues, suggestions, details}。
    req 可携带用户指定的数量（pages/count），用于动态调整页数/题数下限，
    避免"用户要 8 页却强制 ≥12 页"导致小需求必然不合格、拒绝交付。"""
    text = str(content or "")
    req = req or {}
    std = QUALITY_STANDARDS.get(pid, {"min_chars": 400, "required_sections": {}})
    issues: List[str] = []
    suggestions: List[str] = []
    details: Dict[str, Any] = {}
    total_w = 0
    earned_w = 0

    # 0) LLM 错误占位直接不合格
    if is_llm_error(text):
        return {"passed": False, "score": 0, "issues": ["模型调用失败，无有效内容"],
                "suggestions": ["检查模型配置与网络后重新生成"], "details": {}}

    # 1) 字数下限
    min_chars = int(std.get("min_chars") or 0)
    if min_chars > 0:
        total_w += 10
        details["chars"] = {"required": min_chars, "actual": len(text)}
        if len(text) < min_chars:
            issues.append(f"字数不足：当前 {len(text)} 字，要求至少 {min_chars} 字")
            suggestions.append(f"扩充内容至 {min_chars} 字以上：补充具体细节（话术/步骤/数据/实例），禁止概括性空话")
        else:
            earned_w += 10

    # 2) PPT 页数下限（随用户指定数量动态调整）
    min_slides = int(std.get("min_slides") or 0)
    if min_slides > 0 and _to_int(req.get("pages"), 0) > 0:
        # 用户明确指定页数时以用户要求为准，避免"要8页却要求≥12页"必然失败
        min_slides = max(1, _to_int(req.get("pages"), min_slides))
    if min_slides > 0:
        slides = parsed if isinstance(parsed, list) else None
        if slides is None:
            try:
                slides = _safe_json(text)
            except Exception:
                slides = None
        n = len(slides) if isinstance(slides, list) else 0
        total_w += 10
        details["slides"] = {"required": min_slides, "actual": n}
        if n < min_slides:
            issues.append(f"PPT 页数不足：当前 {n} 页，要求至少 {min_slides} 页")
            suggestions.append("增加知识讲授页（每页写具体知识原话）、例题页（完整解答）、练习页与小结页")
        else:
            earned_w += 10

    # 3) 必备章节完整性（PPT 用全文匹配）
    sections: Dict[str, List[str]] = std.get("required_sections") or {}
    sec_detail = []
    for label, pats in sections.items():
        found = any(p in text for p in pats)
        sec_detail.append({"label": label, "found": found})
        total_w += 10
        if found:
            earned_w += 10
        else:
            issues.append(f"缺少必备章节：{label}")
            suggestions.append(f"补充「{label}」部分，确保结构完整规范")
    details["sections"] = sec_detail

    # 4) 占位符检测（字符串黑名单 + 形状正则）
    found_ph = [p for p in GENERIC_PLACEHOLDERS if _has_placeholder(text, p)]
    for pat, label in PLACEHOLDER_REGEXES:
        if re.search(pat, text):
            found_ph.append(label)
    details["placeholders"] = found_ph
    total_w += 5
    if found_ph:
        issues.append("检测到空泛/占位表述：" + "、".join(found_ph[:5]))
        suggestions.append("将占位表述替换为具体内容：写出实际问题、完整解答步骤、可操作的活动指令")
        earned_w -= min(15, len(found_ph) * 3)
    else:
        earned_w += 5

    # 5) 深度检查（正则计数下限）
    dc_detail = []
    for name, pattern, minimum in std.get("depth_checks") or []:
        # 题目数量类检查随用户指定数量动态调整下限
        if name in ("题目数量",) and _to_int(req.get("count"), 0) > 0:
            minimum = max(1, _to_int(req.get("count"), minimum))
        cnt = len(re.findall(pattern, text))
        passed = cnt >= minimum
        dc_detail.append({"name": name, "count": cnt, "min": minimum, "passed": passed})
        total_w += 5
        if passed:
            earned_w += 5
        else:
            issues.append(f"{name}不足：当前 {cnt}，要求至少 {minimum}")
            suggestions.append(f"增加{name}，确保内容深度达标")
    details["depth_checks"] = dc_detail

    score = max(0, min(100, round(earned_w / total_w * 100) if total_w else 0))
    passed = not issues and score >= SETTINGS.quality_min_score
    return {"passed": passed, "score": score, "issues": issues, "suggestions": suggestions, "details": details}


def build_quality_feedback(v: Dict[str, Any]) -> str:
    """把质量校验结果渲染为给 LLM 的整改指令（供重新生成轮次使用）。"""
    lines = [f"质量评分：{v.get('score', 0)}/100（未达标，必须整改）", "存在的问题："]
    for i, it in enumerate(v.get("issues") or [], 1):
        lines.append(f"{i}. {it}")
    if v.get("suggestions"):
        lines.append("整改要求：")
        for i, s in enumerate(v["suggestions"], 1):
            lines.append(f"{i}. {s}")
    return "\n".join(lines)


# ===========================================================================
# 9.6 本地意图解析兜底（LLM 失败/返回不可解析时使用正则规则解析，不再硬编码默认教案）
# ===========================================================================
def local_intent_parse(raw: str) -> Dict[str, Any]:
    """规则式意图解析：与 MockLLM 的意图逻辑一致，作为真实模型失败时的兜底。"""
    return _safe_json(_mock_generate("你是教学智能体的高级意图识别器。", raw))


# ===========================================================================
# 10. 通用「生成→校验→循环」工具 与 产物生成器
# ===========================================================================
def _generate_and_refine(
    prompt_tpl: str,
    fmt: Dict[str, Any],
    validator: Optional[Callable[[str], Tuple[bool, str]]],
    extract: Callable[[str], Any] = lambda s: s,
    pid: str = "",
) -> Dict[str, Any]:
    """通用「生成→质量校验→循环重写」工具。
    质量门强制启用：先跑旧式关键词 validator（若有），再跑 QUALITY_STANDARDS 硬校验，
    任一不合格即带整改反馈重新生成；达到轮次上限仍不合格 → ok=False（导出层拒绝落盘）。
    返回 {content, parsed, rounds, ok, quality}。"""
    best = ""
    ok = False
    feedback = ""
    rounds = 0
    quality: Dict[str, Any] = {}
    # 用户指定的数量（页数/题数），供质量门动态下限使用
    req = {"pages": fmt.get("pages"), "count": fmt.get("count")}
    for r in range(1, SETTINGS.max_refine_rounds + 2):
        rounds = r
        user = "请生成。"
        if r > 1:
            user = f"请重新生成（第{r}稿）。\n上一稿质量检测未通过：\n{feedback}\n请逐条整改后完整重新输出。"
        text = _chat(_fmt_prompt(prompt_tpl, fmt), user)
        if is_llm_error(text):
            best = text
            feedback = "模型调用失败，未返回有效内容。"
            continue
        # 旧式关键词校验（保留兼容性提示）；仅在不通过时记录原因
        legacy_reason = ""
        if validator:
            v_ok, _reason = validator(text)
            if not v_ok:
                legacy_reason = str(_reason)
        # 硬性质量门
        if SETTINGS.quality_gate and pid:
            quality = validate_product(pid, text, None, req)
            if quality["passed"] and not legacy_reason:
                best, ok = text, True
                break
            best = text
            parts = []
            if legacy_reason:
                parts.append("结构缺失：" + legacy_reason)
            parts.append(build_quality_feedback(quality))
            feedback = "\n".join(parts)
        else:
            if not validator:
                best, ok = text, True
                break
            v_ok, reason = validator(text)
            feedback = reason
            if v_ok:
                best, ok = text, True
                break
            best = text
    parsed = None
    try:
        parsed = extract(best)
    except Exception:
        pass
    if is_llm_error(best):
        ok = False
    if SETTINGS.quality_gate and pid and not quality:
        quality = validate_product(pid, best, parsed, req)
    return {"content": best, "parsed": parsed, "rounds": rounds, "ok": ok, "quality": quality}


# 校验器集合（按 PRODUCERS 中 validator 键取值）
VALIDATORS: Dict[str, Callable[[str], Tuple[bool, str]]] = {
    "lesson_plan": lambda t: (
        all(k in t for k in ["教学目标", "学情分析", "重难点", "教学过程", "板书", "作业", "教学反思"]),
        "缺少教案模块",
    ),
    "exercise": lambda t: (
        all(k in t for k in ["基础", "提升", "压轴", "解析"]),
        "缺少梯度/解析",
    ),
    "paper": lambda t: (
        any(k in t for k in ["题型", "分值", "解析", "答案"]),
        "缺少试卷结构/答案",
    ),
    "homework": lambda t: (
        "分层" in t and "作业" in t,
        "缺少分层结构",
    ),
    "comment": lambda t: (
        any(k in t for k in ["学习表现", "建议", "寄语", "闪光点"]),
        "缺少评语要素",
    ),
    "score": lambda t: (
        any(k in t for k in ["平均分", "分数段", "薄弱", "建议"]),
        "缺少分析要素",
    ),
    "research": lambda t: (
        any(k in t for k in ["教研", "课例", "评课", "主题", "建议"]),
        "缺少教研要素",
    ),
    "notice": lambda t: (
        any(k in t for k in ["时间", "地点", "事项"]),
        "缺少通知要素",
    ),
    "plan": lambda t: (
        "计划" in t or "安排" in t,
        "缺少计划要素",
    ),
    "summary": lambda t: (
        "总结" in t or "反思" in t or "回顾" in t,
        "缺少总结要素",
    ),
    "speech": lambda t: (
        any(k in t for k in ["开场", "问候", "结语", "致谢"]),
        "缺少发言稿要素",
    ),
    "minutes": lambda t: (
        any(k in t for k in ["会议", "议题", "决议", "待办"]),
        "缺少纪要要素",
    ),
    "class_meeting": lambda t: (
        any(k in t for k in ["主题", "流程", "活动", "目标"]),
        "缺少班会要素",
    ),
    "evaluation": lambda t: (
        any(k in t for k in ["业绩", "师德", "佐证", "基本情况"]),
        "缺少申报要素",
    ),
}


def _make_validator(spec: Dict[str, Any]) -> Optional[Callable[[str], Tuple[bool, str]]]:
    vkey = spec.get("validator")
    if not vkey:
        return None
    base = VALIDATORS.get(vkey)
    if not base:
        return None
    def validator(text: str) -> Tuple[bool, str]:
        ok, reason = base(text)
        return (ok, f"{reason}：{vkey}" if not ok else "通过")
    return validator


def _to_int(v: Any, default: int) -> int:
    """安全整数转换：容忍 LLM 输出的 "15页"/"15"/15.0/null 等脏值，
    抽取首个整数，失败时回退默认值，避免 int() 抛异常拖垮整条流程。"""
    try:
        if v is None or isinstance(v, bool):
            return default
        if isinstance(v, (int, float)):
            return int(v)
        m = re.search(r"-?\d+", str(v))
        return int(m.group()) if m else default
    except Exception:
        return default


def _common_fmt(state: LessonState) -> Dict[str, Any]:
    g = state.get("globals_", {})
    quantities = g.get("quantities")
    if not isinstance(quantities, dict):
        quantities = {}
    return {
        "subject": g.get("subject", "数学"),
        "topic": g.get("topic", ""),
        "grade": g.get("grade", ""),
        "student_tags": g.get("student_tags", ""),
        "student_data": json.dumps(state.get("history") or "按班级基础：" + g.get("student_tags", ""), ensure_ascii=False),
        "key_points": state.get("key_points", ""),
        "lesson_plan": (state.get("lesson_plan") or {}).get("content", "")[:800],
        "pages": _to_int(quantities.get("ppt", 12), 12),
        "count": _to_int(quantities.get("exercise", 10), 10),
        "raw": state.get("raw_input", ""),
    }


def generic_product_node_factory(pid: str) -> Callable[[LessonState], LessonState]:
    """为注册表中无专用节点的产物生成通用执行节点。"""
    spec = PRODUCERS[pid]

    def node(state: LessonState) -> LessonState:
        if not _should_run(state, pid):
            return {pid: {"content": "", "ok": False, "skipped": True}}
        # Mock 模式：直接取确定性模板，避免文本匹配受用户输入污染
        if isinstance(_llm(), MockLLM):
            tmpl = _PRODUCT_MOCK.get(pid, f"Mock：{spec['name']}")
            content = tmpl if isinstance(tmpl, str) else json.dumps(tmpl, ensure_ascii=False)
            return {pid: {"content": content, "parsed": None, "rounds": 1, "ok": True}}
        fmt = _common_fmt(state)
        out = _generate_and_refine(PROMPTS[spec["prompt"]], fmt, _make_validator(spec), pid=pid)
        return {pid: out}

    return node


# ---------------- 专用节点（教学核心 6 项） ----------------
def _mock_node_output(pid: str, topic: str = "专题", pages: int = 12, count: int = 10) -> Dict[str, Any]:
    """Mock 模式的确定性产物输出（专用节点共用）。"""
    if pid == "lesson_plan":
        content = (
            "## 教学目标\n知识与技能：理解核心概念，掌握基本方法。\n"
            "过程与方法：经历探索过程，渗透数学思想。\n情感态度：激发学习兴趣，培养理性精神。\n"
            "## 学情分析\n班级基础中等，已具备相应知识储备，对抽象概念理解不足。\n"
            "## 教学重难点\n重点：概念与方法掌握。难点：抽象理解与综合应用。\n"
            "## 教学过程\n导入：情境引入，激发兴趣。\n新授：概念归纳→例题示范→方法提炼。\n"
            "巩固：即时练习，辨析易错点。\n小结：梳理知识结构。\n作业：分层布置。\n"
            "## 板书设计\n主板书呈现知识框架，副板书呈现典型例题。\n"
            "## 分层作业\n基础巩固3题、能力提升2题、拓展压轴1题。\n"
            "## 教学反思\n关注学生理解程度，及时调整教学节奏。"
        )
        return {"content": content, "parsed": None, "rounds": 1, "ok": True}
    if pid == "ppt":
        slides = [
            {"title": f"第{idx}页", "content": f"《{topic}》内容要点", "note": "讲稿", "animation": "淡入"}
            for idx in range(1, pages + 1)
        ]
        return {"content": json.dumps(slides, ensure_ascii=False), "parsed": slides, "rounds": 1, "ok": True}
    if pid == "exercise":
        n1, n2, n3 = max(1, count // 3), max(1, count // 3), max(1, count - count // 3 * 2)
        content = (
            f"## 基础巩固（{n1}题）\n1. 基础概念判断题。\n解析：依据定义判断。考点：概念理解。易错点：概念混淆。\n"
            f"## 能力提升（{n2}题）\n1. 综合应用证明题。\n解析：逐步推理。考点：方法运用。易错点：步骤遗漏。\n"
            f"## 拓展压轴（{n3}题）\n1. 含参讨论综合题。\n解析：分类讨论。考点：综合能力。易错点：分类不全。"
        )
        return {"content": content, "parsed": None, "rounds": 1, "ok": True}
    if pid == "animation":
        data = {"title": topic + "动态演示", "type": "教学",
                "steps": ["展示核心对象", "动态变化过程", "标注关键节点"],
                "interaction": "拖动/点击参数，实时观察变化", "tech": "HTML5/SVG"}
        return {"content": json.dumps(data, ensure_ascii=False), "parsed": data, "rounds": 1, "ok": True}
    if pid == "learning_analysis":
        content = (
            "## 班级整体掌握情况\n班级整体中等，概念掌握较好，综合应用偏弱。\n"
            "## 高频易错点预判\n概念混淆、步骤不规范、审题不清。\n"
            "## 分层教学策略\nA层：综合提升；B层：巩固方法；C层：夯实基础。\n"
            "## 学困生辅导方案\n一对一错题归因，基础练习巩固。\n"
            "## 优生拓展方向\n跨章节综合与竞赛思维训练。"
        )
        return {"content": content, "parsed": None, "rounds": 1, "ok": True}
    if pid == "explainer":
        content = "## 精讲内容\n先讲概念本源，再给判定方法，最后典型例题三步示范。\n## 易错提醒\n强调'任意/存在'等关键量词辨析。"
        return {"content": content, "parsed": None, "rounds": 1, "ok": True}
    return {"content": "", "parsed": None, "rounds": 1, "ok": False}


def lesson_plan_node(state: LessonState) -> LessonState:
    if not _should_run(state, "lesson_plan"):
        return {"lesson_plan": {"content": "", "ok": False, "skipped": True}}
    if isinstance(_llm(), MockLLM):
        return {"lesson_plan": _mock_node_output("lesson_plan", state.get("globals_", {}).get("topic", "专题"))}
    out = _generate_lesson_plan_parallel(state)
    return {"lesson_plan": out}


# ===========================================================================
# 教案分节并行生成（替代单次超长生成，显著降低超时、提升速度与完整性）
#   将教案拆为 4 个逻辑分节，用线程池并行调用 LLM（互不阻塞），
#   再按固定顺序组装为完整 Markdown 教案。每个分节短小、单页输出，
#   避免单次请求生成 3000+ 字导致超时/半途截断。
# ===========================================================================
_LESSON_PLAN_SECTIONS = [
    ("s1", "开篇", "教案的开篇部分（用 ## 作为二级标题）", (
        "1. ## 课标依据与教材分析\n2. ## 学情分析\n"
        "3. ## 教学目标（核心素养四维度：数学抽象/逻辑推理/数学建模/直观想象等，"
        "用可观察可测量的行为动词，如\"能够…\"）\n4. ## 教学重点与难点（各附突破策略）"
    )),
    ("s2", "环节一二", "教案\"教学过程\"的前两个教学环节（用 ## 教学过程 作为总标题，用 ### 环节名称 作为子标题）", (
        "### 环节一：情境导入（附具体情境材料与导入语原话）\n"
        "### 环节二：新知探究（附概念/定理的完整表述、关键设问）"
    )),
    ("s3", "环节三四五", "教案\"教学过程\"的后三个教学环节（用 ### 环节名称 作为子标题）", (
        "### 环节三：例题精讲\n### 环节四：变式训练\n### 环节五：课堂小结"
    )),
    ("s4", "收尾", "教案的收尾部分（用 ## 作为二级标题）", (
        "1. ## 例题与变式（完整题干+分步解答，至少2道例题，数学公式用 LaTeX）\n"
        "2. ## 分层作业（基础/提升/拓展三层，各给出具体题目）\n"
        "3. ## 板书设计（结构化呈现，用缩进或列表）\n4. ## 教学反思预设"
    )),
]


def _generate_lesson_plan_parallel(state: LessonState) -> Dict[str, Any]:
    """分节并行生成教案，返回 {content, parsed, rounds, ok}。"""
    g = state.get("globals_", {})
    fmt = _common_fmt(state)
    topic = fmt["topic"] or "未指定"
    subject = fmt["subject"]

    base = (
        f"你是资深高中教研专家，擅长编写结构完整、内容具体可操作的教案片段。\n"
        f"课题：《{topic}》，学科：{subject}，年级：{fmt['grade'] or '未指定'}。\n"
        f"学情：{fmt['student_tags'] or '（未指定，基于通用学情）'}。\n"
        f"重难点：{fmt['key_points'] or '（见教案前文）'}。\n\n"
        f"课标与考法资料：\n{fmt['student_data'] or '（暂无联网资料，基于教学经验生成）'}\n\n"
        f"所有内容必须具体可操作，禁止\"此处讲解\"\"举例说明\"\"略\"等空泛表述。\n"
        f"数学公式用 LaTeX（行内 $...$，块级 $$...$$）。只输出纯 Markdown 正文，不要任何解释或代码块标记。"
    )

    def gen(spec: Tuple[str, str, str, str]) -> Tuple[str, str]:
        key, label, head, subs = spec
        prompt = (
            f"{base}\n\n请编写{head}，必须包含以下小节：\n{subs}\n"
            f"教学过程各环节必须写清：【教师活动】含具体提问话术（用\"师：\"标注）、"
            f"【学生活动】含预设回答（用\"生：\"标注）、【设计意图】、【时间分配】。\n"
            f"该分节约 700-900 字。"
        )
        return key, _chat(prompt, "你是教案分节撰写专家，只输出本分节 Markdown 正文。")

    order = ["s1", "s2", "s3", "s4"]
    results: Dict[str, str] = {}
    with ThreadPoolExecutor(max_workers=4) as ex:
        futures = [ex.submit(gen, spec) for spec in _LESSON_PLAN_SECTIONS]
        for fut in as_completed(futures):
            try:
                key, text = fut.result()
                results[key] = text
            except Exception as e:
                log.warning("教案分节生成失败：%s", e)

    if not results:
        return {"content": "", "parsed": None, "rounds": 1, "ok": False,
                "error": "教案分节生成失败：模型未返回有效内容（可能超时）"}

    parts = []
    for key in order:
        t = (results.get(key) or "").strip()
        if t and not is_llm_error(t):
            parts.append(t)
    content = "\n\n".join(parts)
    if not content.strip():
        return {"content": "", "parsed": None, "rounds": 1, "ok": False, "error": "教案分节组装为空"}

    # ---- 质量门：整体验收，不合格则定向重写最弱分节（最多 max_refine_rounds 轮） ----
    rounds = 1
    quality: Dict[str, Any] = {}
    while SETTINGS.quality_gate and rounds <= SETTINGS.max_refine_rounds:
        quality = validate_product("lesson_plan", content)
        if quality["passed"]:
            break
        rounds += 1
        feedback = build_quality_feedback(quality)
        log.info("教案质量未达标（%d 分），第 %d 轮定向重写：%s", quality["score"], rounds, "；".join(quality["issues"][:2]))
        # 定向重写：把反馈注入各分节重新生成（并行），重点补足缺失章节
        fix_prompt_suffix = (
            f"\n\n【整改要求】上一稿质量检测未通过：\n{feedback}\n"
            f"请在本分节中重点补足上述缺失内容，字数不少于 900 字。"
        )

        def regen(spec: Tuple[str, str, str, str]) -> Tuple[str, str]:
            key, label, head, subs = spec
            prompt = (
                f"{base}\n\n请编写{head}，必须包含以下小节：\n{subs}\n"
                f"教学过程各环节必须写清：【教师活动】含具体提问话术（用\"师：\"标注）、"
                f"【学生活动】含预设回答（用\"生：\"标注）、【设计意图】、【时间分配】。{fix_prompt_suffix}"
            )
            return key, _chat(prompt, "你是教案分节撰写专家，只输出本分节 Markdown 正文。")

        with ThreadPoolExecutor(max_workers=4) as ex:
            futures = [ex.submit(regen, spec) for spec in _LESSON_PLAN_SECTIONS]
            for fut in as_completed(futures):
                try:
                    key, text = fut.result()
                    if text and not is_llm_error(text):
                        results[key] = text
                except Exception as e:
                    log.warning("教案分节重写失败：%s", e)
        parts = [results[k].strip() for k in order if (results.get(k) or "").strip() and not is_llm_error(results[k])]
        content = "\n\n".join(parts)

    if SETTINGS.quality_gate:
        quality = validate_product("lesson_plan", content)
        if not quality["passed"]:
            log.warning("教案经 %d 轮重写仍未达标（%d 分）：%s", rounds, quality["score"], quality["issues"][:3])
            return {"content": content, "parsed": None, "rounds": rounds, "ok": False,
                    "quality": quality, "error": "质量未达标：" + "；".join(quality["issues"][:3])}
    return {"content": content, "parsed": None, "rounds": rounds, "ok": True, "quality": quality}


def ppt_node(state: LessonState) -> LessonState:
    if not _should_run(state, "ppt"):
        return {"ppt": {"content": "", "ok": False, "skipped": True}}
    if isinstance(_llm(), MockLLM):
        g = state.get("globals_", {})
        return {"ppt": _mock_node_output("ppt", g.get("topic", "专题"),
                                         pages=int((g.get("quantities") or {}).get("ppt", 12)))}
    fmt = _common_fmt(state)
    out = _generate_and_refine(PROMPTS["ppt"], fmt, None, extract=_safe_json, pid="ppt")
    return {"ppt": out}


def exercise_node(state: LessonState) -> LessonState:
    if not _should_run(state, "exercise"):
        return {"exercise": {"content": "", "ok": False, "skipped": True}}
    if isinstance(_llm(), MockLLM):
        g = state.get("globals_", {})
        return {"exercise": _mock_node_output("exercise", g.get("topic", "专题"),
                                              count=int((g.get("quantities") or {}).get("exercise", 10)))}
    fmt = _common_fmt(state)
    out = _generate_and_refine(PROMPTS["exercise"], fmt, _make_validator(PRODUCERS["exercise"]), pid="exercise")
    return {"exercise": out}


def animation_node(state: LessonState) -> LessonState:
    if not _should_run(state, "animation"):
        return {"animation": {"content": "", "ok": False, "skipped": True}}
    if isinstance(_llm(), MockLLM):
        return {"animation": _mock_node_output("animation", state.get("globals_", {}).get("topic", "专题"))}
    fmt = _common_fmt(state)
    out = _generate_and_refine(PROMPTS["animation"], fmt, None, extract=_safe_json, pid="animation")
    return {"animation": out}


def learning_analysis_node(state: LessonState) -> LessonState:
    if not _should_run(state, "learning_analysis"):
        return {"learning_analysis": {"content": "", "ok": False, "skipped": True}}
    if isinstance(_llm(), MockLLM):
        return {"learning_analysis": _mock_node_output("learning_analysis")}
    fmt = _common_fmt(state)
    out = _generate_and_refine(PROMPTS["learning_analysis"], fmt, None, pid="learning_analysis")
    return {"learning_analysis": out}


def explainer_node(state: LessonState) -> LessonState:
    if not _should_run(state, "explainer"):
        return {"explainer": {"content": "", "ok": False, "skipped": True}}
    if isinstance(_llm(), MockLLM):
        return {"explainer": _mock_node_output("explainer")}
    # 优先复用任务拆解阶段的精讲内容；过短时走质量门重新生成
    base_text = state.get("explanation", "") or ""
    if base_text and not is_llm_error(base_text):
        q = validate_product("explainer", base_text)
        if q["passed"] or not SETTINGS.quality_gate:
            return {"explainer": {"content": base_text, "rounds": 1, "ok": True, "quality": q}}
    fmt = _common_fmt(state)
    out = _generate_and_refine(PROMPTS["explainer"], fmt, None, pid="explainer")
    return {"explainer": out}


# 节点注册表：id -> 节点函数（未列出的走通用工厂）
_SPECIAL_NODES: Dict[str, Callable[[LessonState], LessonState]] = {
    "lesson_plan": lesson_plan_node,
    "ppt": ppt_node,
    "exercise": exercise_node,
    "animation": animation_node,
    "learning_analysis": learning_analysis_node,
    "explainer": explainer_node,
}


def _node_for(pid: str) -> Callable[[LessonState], LessonState]:
    if pid in _SPECIAL_NODES:
        return _SPECIAL_NODES[pid]
    return generic_product_node_factory(pid)


# ===========================================================================
# 11. 汇聚与最终输出
# ===========================================================================
def aggregate(state: LessonState) -> LessonState:
    tasks = state.get("tasks", [])
    results = {}
    quality_report: Dict[str, Any] = {}
    failed: List[str] = []
    for pid in tasks:
        res = state.get(pid) or {}
        content = res.get("content", "")
        if content and not is_llm_error(content):
            results[PRODUCERS[pid]["name"]] = content
        q = res.get("quality") or {}
        if q:
            quality_report[PRODUCERS[pid]["name"]] = {
                "score": q.get("score", 0), "passed": bool(q.get("passed")),
                "issues": q.get("issues", []),
            }
        if not res.get("ok"):
            failed.append(PRODUCERS[pid]["name"])
    task_text = "、".join(PRODUCERS[t]["name"] for t in tasks)
    # 汇总提示词只传产物摘要（名称+前 200 字），避免全量原文导致上下文爆炸/超时
    brief = {name: (c[:200] + "…" if len(c) > 200 else c) for name, c in results.items()}
    text = _chat(
        _fmt_prompt(PROMPTS["aggregate"], {"tasks": task_text, "results": json.dumps(brief, ensure_ascii=False)}),
        "请汇总资源包。",
    )
    if is_llm_error(text) or not text.strip():
        text = "资源包总览：\n" + "\n".join(f"- {name}" for name in results)
    if failed:
        text += "\n\n【质量警示】以下产物经多轮重写仍未达质量标准，未予交付：" + "、".join(failed)
    return {"final": {"summary": text, "results": results, "tasks": tasks,
                      "quality": quality_report, "failed": failed}}


def final_output(state: LessonState) -> LessonState:
    return dict(state)


def error_handler(state: LessonState, error: Exception) -> LessonState:
    log.error("工作流节点异常：%s", error)
    errors = list(state.get("errors", []))
    errors.append(str(error))
    g = state.get("globals_", {})
    fallback = {
        "content": f"【兜底产物】主题：{g.get('topic','')}。异常：{error}。建议重试或检查模型配置。",
        "ok": False,
    }
    return {"errors": errors, "lesson_plan": state.get("lesson_plan") or fallback}


# ===========================================================================
# 11.6 公式渲染与 Markdown 清理（docx/pptx 共用基础设施）
# ===========================================================================
_LATEX_RE = re.compile(r"(\$\$[^$]+\$\$|\$[^$\n]+?\$)")

_MPL_READY = False


def _mpl():
    """获取配置好中文字体的 matplotlib pyplot（惰性初始化，进程内一次）。"""
    global _MPL_READY
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    if not _MPL_READY:
        for fname in ("Microsoft YaHei", "SimHei", "SimSun"):
            try:
                matplotlib.font_manager.findfont(fname, fallback_to_default=False)
                plt.rcParams["font.sans-serif"] = [fname, "DejaVu Sans"]
                break
            except Exception:
                continue
        plt.rcParams["axes.unicode_minus"] = False
        plt.rcParams["mathtext.fontset"] = "cm"
        _MPL_READY = True
    return plt


def _strip_md_inline(s: str) -> str:
    """清除行内 Markdown 符号（粗体/行内代码/链接/图片标记），保留纯文本。"""
    s = re.sub(r"!\[[^\]]*\]\([^)]*\)", "", s)
    s = re.sub(r"\[([^\]]*)\]\([^)]*\)", r"\1", s)
    s = re.sub(r"\*\*([^*]+)\*\*", r"\1", s)
    s = re.sub(r"__([^_]+)__", r"\1", s)
    s = re.sub(r"`([^`]*)`", r"\1", s)
    s = re.sub(r"^>+\s*", "", s)
    return s.strip()


def _latex_img(latex: str, fontsize: int = 16, dpi: int = 200):
    """把 LaTeX 公式渲染为透明 PNG（BytesIO），失败返回 None。"""
    fig = None
    try:
        import io
        plt = _mpl()
        fig = plt.figure()
        fig.text(0, 0, f"${latex}$", fontsize=fontsize)
        buf = io.BytesIO()
        fig.savefig(buf, format="png", dpi=dpi, transparent=True,
                    bbox_inches="tight", pad_inches=0.04)
        buf.seek(0)
        return buf
    except Exception as e:
        log.warning("公式渲染失败（%s）：%s", latex[:40], e)
        return None
    finally:
        # 无论成功失败都关闭 figure，避免公式多时累积泄漏
        if fig is not None:
            try:
                _mpl().close(fig)
            except Exception:
                pass


def _mixed_line_img(text: str, fontsize: int = 20, dpi: int = 200):
    """把含 $...$ 的整行（文本+公式混排）渲染为一张透明 PNG，用于 PPT 公式行。"""
    fig = None
    try:
        import io
        plt = _mpl()
        # matplotlib mathtext 只认单 $ 定界：$$...$$ 先归一为 $...$
        text = re.sub(r"\$\$([^$]+)\$\$", r"$\1$", text)
        fig = plt.figure()
        fig.text(0, 0, text, fontsize=fontsize)
        buf = io.BytesIO()
        fig.savefig(buf, format="png", dpi=dpi, transparent=True,
                    bbox_inches="tight", pad_inches=0.04)
        buf.seek(0)
        return buf
    except Exception as e:
        log.warning("混排行渲染失败：%s", e)
        return None
    finally:
        if fig is not None:
            try:
                _mpl().close(fig)
            except Exception:
                pass


def _img_size_in(buf, dpi: int) -> Tuple[float, float]:
    """读取 PNG 字节流的像素尺寸并换算为英寸。"""
    try:
        from PIL import Image
        buf.seek(0)
        im = Image.open(buf)
        w, h = im.size
        buf.seek(0)
        return w / dpi, h / dpi
    except Exception:
        return 4.0, 0.5


def _has_latex(s: str) -> bool:
    return bool(_LATEX_RE.search(s or ""))


# ===========================================================================
# 11.5 真实文档导出（docx / xlsx / pptx / md 落盘 + WPS/飞书推送扩展点）
# ===========================================================================
_WIN_RESERVED = {
    "CON", "PRN", "AUX", "NUL",
    "COM1", "COM2", "COM3", "COM4", "COM5", "COM6", "COM7", "COM8", "COM9",
    "LPT1", "LPT2", "LPT3", "LPT4", "LPT5", "LPT6", "LPT7", "LPT8", "LPT9",
}


def _safe_filename(name: str) -> str:
    """去除文件名中的非法字符，并规避 Windows 保留名/结尾句点。"""
    cleaned = re.sub(r'[\\/:*?"<>|\r\n\t]', "_", str(name or "")).strip().rstrip(".")
    stem = cleaned.split(".")[0].upper()
    if stem in _WIN_RESERVED:
        cleaned = "_" + cleaned
    return cleaned or "未命名"


def _flush_md_table(doc, buf: List[List[str]]) -> None:
    """把缓存的 Markdown 表格行写入 docx 真实表格。"""
    if not buf:
        return
    try:
        ncols = max(len(r) for r in buf)
        table = doc.add_table(rows=len(buf), cols=ncols)
        table.style = "Table Grid"
        for i, row in enumerate(buf):
            for j in range(ncols):
                cell = table.cell(i, j)
                cell.text = row[j] if j < len(row) else ""
                if i == 0:
                    for p in cell.paragraphs:
                        for r in p.runs:
                            r.font.bold = True
    except Exception as e:
        log.warning("docx 表格写入失败（降级为段落）：%s", e)
        for row in buf:
            doc.add_paragraph(" | ".join(row))


def _add_rich_paragraph(doc, text: str, style: Optional[str] = None) -> None:
    """写入段落：支持 **粗体**、行内代码清理、$...$ LaTeX 公式渲染为图片嵌入。"""
    from docx.shared import Pt as _Pt

    p = doc.add_paragraph(style=style) if style else doc.add_paragraph()
    # 先按 LaTeX 公式切段，再对文本段做粗体处理
    for seg in _LATEX_RE.split(text):
        if not seg:
            continue
        m = re.fullmatch(r"(\$\$)([^$]+)\$\$", seg)
        mi = re.fullmatch(r"\$([^$\n]+?)\$", seg)
        if m or mi:
            latex = (m.group(2) if m else mi.group(1)).strip()
            buf = _latex_img(latex, fontsize=15, dpi=200)
            if buf is not None:
                try:
                    w_in, h_in = _img_size_in(buf, 200)
                    # 行内公式：高度与正文字号匹配（约 1.4 倍行高内），宽度等比
                    tgt_h = _Pt(20)
                    run = p.add_run()
                    run.add_picture(buf, height=tgt_h)
                    continue
                except Exception as e:
                    log.warning("公式图片插入失败：%s", e)
            # 渲染失败则降级为纯文本公式
            p.add_run(latex)
            continue
        # 普通文本段：处理 **粗体**，并清理行内代码等残留符号
        for part in re.split(r"(\*\*[^*]+\*\*)", seg):
            if not part:
                continue
            if part.startswith("**") and part.endswith("**") and len(part) > 4:
                r = p.add_run(_strip_md_inline(part[2:-2]))
                r.font.bold = True
            else:
                p.add_run(_strip_md_inline(part))


def _md_to_docx(doc, text: str) -> None:
    """把 Markdown 文本写入 docx：标题/列表/段落/真实表格/粗体/LaTeX 公式渲染。"""
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Inches as _Inches

    table_buf: List[List[str]] = []
    in_code = False
    for raw in text.splitlines():
        line = raw.rstrip()
        s = line.strip()
        # 代码块围栏：围栏内以等宽段落原样输出，不做 Markdown 解析
        if s.startswith("```"):
            in_code = not in_code
            continue
        if in_code:
            p = doc.add_paragraph()
            r = p.add_run(line)
            try:
                r.font.name = "Consolas"
            except Exception:
                pass
            continue
        # 表格行：缓存，遇非表格行统一落盘。
        # 兼容无前导 | 的表格（LLM 常见输出）：含 >=2 个 | 且非标题/列表行也视为表格行。
        is_table_line = (s.startswith("|") and s.endswith("|")) or (
            s.count("|") >= 2 and not s.startswith(("#", "-", "*", ">"))
        )
        if is_table_line:
            cells = [_strip_md_inline(re.sub(r"\$([^$]+)\$", r"\1", c.strip())) for c in s.strip("|").split("|")]
            if not all(re.fullmatch(r":?-{2,}:?", c) for c in cells):  # 跳过分隔行
                table_buf.append(cells)
            continue
        if table_buf:
            _flush_md_table(doc, table_buf)
            table_buf = []
        if not s or re.fullmatch(r"-{3,}|\*{3,}|_{3,}", s):
            continue
        # 独立成行的块级公式 $$...$$：居中、大图
        dm = re.fullmatch(r"\$\$([^$]+)\$\$", s)
        if dm:
            buf = _latex_img(dm.group(1).strip(), fontsize=18, dpi=200)
            if buf is not None:
                p = doc.add_paragraph()
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                w_in, h_in = _img_size_in(buf, 200)
                try:
                    p.add_run().add_picture(buf, height=_Inches(min(max(h_in, 0.3), 1.2)))
                except Exception:
                    p.add_run(dm.group(1).strip())
            else:
                doc.add_paragraph(dm.group(1).strip())
            continue
        # 标题：统一匹配 1-6 级并 clamp，避免深层标题（#####/######）带 # 号原样输出
        hm = re.match(r"^(#{1,6})\s+(.*)", s)
        if hm:
            level = min(len(hm.group(1)) - 1, 4)
            doc.add_heading(_strip_md_inline(hm.group(2)), level=level)
            continue
        if s.startswith("- ") or s.startswith("* "):
            _add_rich_paragraph(doc, s[2:], style="List Bullet")
        elif re.match(r"^\s*\d+[.、]", line):
            # 剥掉行首数字前缀再套自动编号样式，避免双重编号"1. 1. xxx"
            _add_rich_paragraph(doc, re.sub(r"^\s*\d+[.、]\s*", "", line.lstrip()), style="List Number")
        else:
            _add_rich_paragraph(doc, line)
    if table_buf:
        _flush_md_table(doc, table_buf)


def _md_table_rows(text: str) -> List[List[str]]:
    """从 Markdown 文本中解析表格行（| a | b |）。"""
    rows = []
    for line in text.splitlines():
        s = line.strip()
        if s.startswith("|") and s.endswith("|"):
            cells = [c.strip() for c in s.strip("|").split("|")]
            # 跳过表头分隔行 | --- | --- |
            if all(re.fullmatch(r":?-{2,}:?", c) for c in cells):
                continue
            rows.append(cells)
    return rows


def _export_docx(path: str, content: str) -> bool:
    try:
        from docx import Document
        from docx.oxml.ns import qn
        from docx.shared import Pt

        doc = Document()
        # 设置默认中文字体（正文宋体/标题黑体），避免打开后乱码或默认西文字体
        try:
            normal = doc.styles["Normal"]
            normal.font.name = "Times New Roman"
            normal.font.size = Pt(11)
            normal.element.rPr.rFonts.set(qn("w:eastAsia"), "宋体")
            for hstyle in ("Heading 1", "Heading 2", "Title"):
                try:
                    st = doc.styles[hstyle]
                    st.element.rPr.rFonts.set(qn("w:eastAsia"), "黑体")
                except Exception:
                    pass
        except Exception:
            pass
        doc.add_heading(_safe_filename(os.path.basename(path)[:-5]), level=0)
        _md_to_docx(doc, content)
        doc.save(path)
        # 落盘校验：文件必须存在且非空壳
        return os.path.exists(path) and os.path.getsize(path) > 3000
    except Exception as e:
        log.warning("docx 导出失败（回退 md）：%s", e)
        return False


def _export_xlsx(path: str, content: str) -> bool:
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font, Alignment

        wb = Workbook()
        ws = wb.active
        rows = _md_table_rows(content)
        if not rows:  # 无表格则写入纯文本
            ws.cell(row=1, column=1, value=content)
        else:
            for r, cells in enumerate(rows, start=1):
                for c, val in enumerate(cells, start=1):
                    cell = ws.cell(row=r, column=c, value=val)
                    if r == 1:
                        cell.font = Font(bold=True)
                        cell.alignment = Alignment(horizontal="center")
        for col in ws.columns:
            width = max((len(str(c.value)) for c in col if c.value is not None), default=8) + 2
            ws.column_dimensions[col[0].column_letter].width = min(width, 40)
        wb.save(path)
        return os.path.exists(path) and os.path.getsize(path) > 2000
    except Exception as e:
        log.warning("xlsx 导出失败（回退 md）：%s", e)
        return False


def _export_pptx(path: str, content: str, parsed: Any) -> bool:
    """导出 PPTX。硬性要求：
    - 数学/物理公式（$...$）渲染为图片，绝不原样输出源码；
    - 每个内容页必须有配图（图表/函数图像/示意图），图文结合；
    - Markdown 符号（**、`、## 等）一律清理，不进入幻灯片；
    - 文件体积 ≥ 1MB（不足则提升图片 DPI 重建）。
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE

        if isinstance(parsed, list) and parsed:
            slides = parsed
        else:
            slides = _md_table_rows(content)
            if not slides:
                # 解析不出有效幻灯片内容：不伪造空壳 PPT，回退由调用方处理（写 md）
                log.warning("pptx 无可解析的幻灯片内容，放弃生成空壳文件")
                return False

        # 主题色
        JADE = (79, 122, 102)       # 竹青主色
        TAN = (168, 129, 78)        # 秋香赭点缀
        INK = (47, 42, 34)          # 暖墨
        LIGHT = (247, 246, 241)     # 米白背景

        MIN_PPTX_BYTES = 1024 * 1024  # 硬性要求：PPT ≥ 1MB

        def _make_bg_image(dpi: int):
            """全页背景图（竹青-米白柔和渐变+细纹理），保证图文质感与文件体积达标。"""
            import io
            import numpy as np
            plt = _mpl()
            w, h = 13.333, 7.5
            fig = plt.figure(figsize=(w, h), dpi=dpi)
            ax = fig.add_axes([0, 0, 1, 1])
            ax.axis("off")
            n_y, n_x = 480, 860
            yy, xx = np.mgrid[0:n_y, 0:n_x]
            base = np.array([247, 246, 241], dtype=float)   # 米白
            tint = np.array([226, 233, 228], dtype=float)   # 淡竹青
            mix = (xx / n_x * 0.55 + yy / n_y * 0.45)[..., None]
            img = base * (1 - mix) + tint * mix
            rng = np.random.default_rng(42)
            img = np.clip(img + rng.normal(0, 1.6, img.shape), 0, 255).astype("uint8")
            ax.imshow(img, aspect="auto")
            buf = io.BytesIO()
            fig.savefig(buf, format="png", dpi=dpi)
            plt.close(fig)
            buf.seek(0)
            return buf

        def _build(dpi: int) -> None:
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            def _set_bg(slide, color):
                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = RGBColor(*color)

            def _set_bg_image(slide, dpi):
                try:
                    buf = _make_bg_image(dpi)
                    slide.shapes.add_picture(buf, 0, 0, width=prs.slide_width, height=prs.slide_height)
                except Exception as e:
                    log.warning("背景图生成失败（使用纯色背景）：%s", e)
                    _set_bg(slide, LIGHT)

            def _add_title(slide, text, color=JADE, size=30):
                tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.9))
                tf = tb.text_frame; tf.word_wrap = True
                p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
                r = p.add_run(); r.text = _strip_md_inline(str(text))
                f = r.font; f.size = Pt(size); f.bold = True; f.color.rgb = RGBColor(*color)
                ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.25), Inches(1.6), Emu(36000))
                ln.fill.solid(); ln.fill.fore_color.rgb = RGBColor(*TAN); ln.line.fill.background()
                return tb

            def _add_text_block(slide, left, top, width, height, text, size=18, color=INK, bold=False):
                tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
                tf = tb.text_frame; tf.word_wrap = True
                for i, line in enumerate(str(text).split("\n")):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.space_after = Pt(8)
                    r = p.add_run(); r.text = _strip_md_inline(line)
                    f = r.font; f.size = Pt(size); f.color.rgb = RGBColor(*color); f.bold = bold
                return tb

            def _add_bullets(slide, bullets, left=0.8, top=1.6, width=6.9, size=20, color=INK, max_bottom=7.1):
                """要点排版：纯文本走文本框；含公式的行渲染为图片。Markdown 符号全部清理。"""
                y = top
                tb = None  # 当前累积的文本框

                def _flush():
                    nonlocal tb
                    tb = None

                for raw_b in bullets or []:
                    b = _strip_md_inline(str(raw_b))
                    if not b:
                        continue
                    if y > max_bottom:
                        break
                    if _has_latex(b):
                        _flush()
                        buf = _mixed_line_img("• " + b, fontsize=max(14, size - 2), dpi=dpi)
                        if buf is not None:
                            w_in, h_in = _img_size_in(buf, dpi)
                            if w_in > width:
                                h_in = h_in * (width / w_in)
                                w_in = width
                            h_in = min(h_in, 1.1)
                            slide.shapes.add_picture(buf, Inches(left), Inches(y), width=Inches(w_in))
                            y += h_in + 0.12
                            continue
                        # 渲染失败降级为纯文本（去掉 $ 符号）
                        b = re.sub(r"\$+", "", b)
                    if tb is None:
                        tb = slide.shapes.add_textbox(Inches(left), Inches(y), Inches(width), Inches(max_bottom - y))
                        tb.text_frame.word_wrap = True
                        p = tb.text_frame.paragraphs[0]
                    else:
                        p = tb.text_frame.add_paragraph()
                    p.space_after = Pt(10)
                    r = p.add_run(); r.text = "• " + b
                    f = r.font; f.size = Pt(size); f.color.rgb = RGBColor(*color)
                    # 估算行高（含换行）：CJK 每行约 width*72/(size*1.02) 字
                    cpl = max(10, int(width * 72 / (size * 1.02)))
                    lines = max(1, (len(b) + cpl - 1) // cpl)
                    y += lines * (size * 1.55 / 72.0) + 0.14

            def _make_visual(chart_type: str, idx: int, title_hint: str = ""):
                """生成配图 PNG（BytesIO）。类型：bar/line/pie/func/diagram。"""
                import io
                plt = _mpl()
                C_JADE = (79/255, 122/255, 102/255)
                C_JADE2 = (107/255, 151/255, 132/255)
                C_TAN = (168/255, 129/255, 78/255)
                C_TAN2 = (200/255, 180/255, 150/255)
                ct = (chart_type or "").lower()
                fig, ax = plt.subplots(figsize=(5.2, 3.4), dpi=dpi)
                if ct == "bar":
                    tags = ["掌握", "理解", "运用", "综合"]
                    vals = [85, 72, 60, 45]
                    ax.bar(tags, vals, color=[C_JADE, C_JADE2, C_TAN, C_TAN2])
                    ax.set_title(title_hint or "知识点掌握度对比", fontsize=13)
                    ax.set_ylim(0, 100)
                elif ct == "pie":
                    ax.pie([35, 45, 20], labels=["基础", "提升", "拓展"], autopct="%1.0f%%",
                           colors=[C_JADE, C_JADE2, C_TAN], textprops={"fontsize": 11})
                elif ct in ("line", "func", "function"):
                    import numpy as np
                    xs = np.linspace(-3, 3, 200)
                    kinds = [xs**2, np.sin(xs * 2) + xs * 0.3, np.exp(xs * 0.5) - 1, np.abs(xs)]
                    labs = ["$f(x)=x^2$", "$f(x)=\\sin 2x+0.3x$", "$f(x)=e^{x/2}-1$", "$f(x)=|x|$"]
                    k = idx % len(kinds)
                    ax.plot(xs, kinds[k], color=C_JADE, linewidth=2.6)
                    ax.axhline(0, color="#888888", linewidth=0.8)
                    ax.axvline(0, color="#888888", linewidth=0.8)
                    ax.set_title("函数图像：" + labs[k], fontsize=13)
                    ax.grid(True, alpha=0.3)
                elif ct == "diagram":
                    import numpy as np
                    theta = np.linspace(0, 2 * np.pi, 100)
                    ax.plot(np.cos(theta), np.sin(theta), color=C_JADE, linewidth=2.2)
                    ax.arrow(0, 0, 0.7, 0.7, head_width=0.08, color=C_TAN)
                    ax.annotate("$\\vec{F}$", xy=(0.72, 0.72), fontsize=14, color=C_TAN)
                    ax.set_xlim(-1.4, 1.4); ax.set_ylim(-1.4, 1.4)
                    ax.set_aspect("equal")
                    ax.set_title(title_hint or "结构示意图", fontsize=13)
                    ax.grid(True, alpha=0.25)
                else:
                    # 默认：按页码轮换 柱/线/饼，保证每页都有真实配图
                    pick = ["bar", "line", "pie"][idx % 3]
                    if pick == "bar":
                        tags = ["掌握", "理解", "运用", "综合"]
                        vals = [85, 72, 60, 45]
                        ax.bar(tags, vals, color=[C_JADE, C_JADE2, C_TAN, C_TAN2])
                        ax.set_title(title_hint or "知识点掌握度对比", fontsize=13)
                        ax.set_ylim(0, 100)
                    elif pick == "pie":
                        ax.pie([35, 45, 20], labels=["基础", "提升", "拓展"], autopct="%1.0f%%",
                               colors=[C_JADE, C_JADE2, C_TAN], textprops={"fontsize": 11})
                    else:
                        import numpy as np
                        xs = np.linspace(-3, 3, 200)
                        ax.plot(xs, xs ** 2, color=C_JADE, linewidth=2.6)
                        ax.axhline(0, color="#888888", linewidth=0.8)
                        ax.axvline(0, color="#888888", linewidth=0.8)
                        ax.set_title("函数图像：$f(x)=x^2$", fontsize=13)
                        ax.grid(True, alpha=0.3)
                ax.spines["top"].set_visible(False); ax.spines["right"].set_visible(False)
                for sp in ("bottom", "left"):
                    ax.spines[sp].set_color("#999999")
                buf = io.BytesIO()
                fig.tight_layout()
                fig.savefig(buf, format="png", transparent=True)
                plt.close(fig)
                buf.seek(0)
                return buf

            for idx, item in enumerate(slides):
                item = item if isinstance(item, dict) else {"title": str(item)}
                title = _strip_md_inline(str(item.get("title", "第%d页" % (idx + 1))))
                bullets = item.get("bullets") or item.get("content") or []
                if isinstance(bullets, str):
                    bullets = [l for l in bullets.split("\n") if l.strip()]
                note = _strip_md_inline(str(item.get("note", "") or ""))
                ptype = str(item.get("type", ""))
                chart = str(item.get("chart", "") or "")

                slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版式
                if ptype in ("title", "cover") or idx == 0:
                    _set_bg(slide, INK)
                else:
                    _set_bg_image(slide, dpi)

                if ptype in ("title", "cover") or idx == 0:
                    _add_title(slide, title, color=TAN, size=40)
                    sub = item.get("subtitle") or item.get("sub") or ""
                    if sub:
                        _add_text_block(slide, 0.8, 1.6, 11.7, 1.0, sub, size=24, color=LIGHT)
                    if bullets:
                        _add_bullets(slide, bullets, left=0.8, top=3.0, width=11.7, size=20, color=LIGHT)
                elif ptype in ("section", "agenda"):
                    _add_title(slide, title, color=JADE, size=34)
                    _add_bullets(slide, bullets, top=2.0, width=11.7, size=24)
                elif ptype == "example":
                    _add_title(slide, title, color=TAN, size=30)
                    _add_bullets(slide, bullets, top=1.7, width=6.9, size=18)
                    try:
                        buf = _make_visual(chart or "func", idx, title)
                    except Exception as ve:
                        log.warning("配图生成失败（跳过该页配图）：%s", ve)
                        buf = None
                    if buf: slide.shapes.add_picture(buf, Inches(7.9), Inches(2.0), width=Inches(4.9))
                elif ptype in ("summary", "end"):
                    _add_title(slide, title, color=JADE, size=34)
                    _add_bullets(slide, bullets, top=2.0, width=11.7, size=22)
                else:
                    _add_title(slide, title, color=JADE, size=30)
                    _add_bullets(slide, bullets, top=1.6, width=6.9, size=19)
                    try:
                        buf = _make_visual(chart, idx, title)
                    except Exception as ve:
                        log.warning("配图生成失败（跳过该页配图）：%s", ve)
                        buf = None
                    if buf: slide.shapes.add_picture(buf, Inches(7.9), Inches(1.9), width=Inches(4.9))
                if note:
                    try:
                        slide.notes_slide.notes_text_frame.text = note
                    except Exception:
                        pass
            prs.save(path)

        # DPI 阶梯：尽量达到 1MB 目标体积（图片分辨率提升带来体积增长）
        dpi = 200
        for _ in range(2):
            _build(dpi)
            if os.path.exists(path) and os.path.getsize(path) >= MIN_PPTX_BYTES:
                break
            dpi += 100
        # 落盘校验：文件存在且非空壳即合格。体积不再作为硬性否决指标——
        # 原"≥1MB"硬指标会让内容正确、但页数少/图片压缩率高的 PPT 被误判不合格并降级为 .md。
        size = os.path.getsize(path) if os.path.exists(path) else 0
        ok = size >= 10 * 1024
        if ok and size < MIN_PPTX_BYTES:
            log.info("PPTX 体积 %d B 未达 1MB 期望值（页数少/图片压缩），仍按合格导出", size)
        elif not ok:
            log.warning("PPTX 体积过小（%d B），疑似空壳，按不合格处理", size)
        return ok
    except Exception as e:
        log.warning("pptx 导出失败（回退 md）：%s", e)
        return False


def _strip_chat_prefix(text: str) -> str:
    """剥离产物开头的对话式开场白（如"以下是为您重新生成的…，已全面消除空泛表述"），
    防止这类非正文内容泄漏进交付文档。只剥离以常见对话措辞开头的行，
    一旦遇到标题/表格/分隔线或普通正文即停止，避免误伤正文。"""
    if not text:
        return text
    lines = text.split("\n")
    chat_pat = re.compile(
        r"^(?:好的[，,。]?\s*)?(?:以下是|已为您|为您|这是|根据您的(?:要求|指示)|现(?:为您)?提供|现为您|已按要求|已全面|已去除|已消除|已重新|已更新|已优化)"
    )
    i = 0
    n = len(lines)
    while i < n:
        s = lines[i].strip()
        if not s:
            i += 1
            continue
        # 遇到标题 / 表格 / 分隔线，视为正文开始，停止剥离
        if s.startswith("#") or s.startswith("|") or (len(s) >= 3 and set(s) <= set("-—=_* ")):
            break
        if chat_pat.match(s):
            i += 1
            continue
        break
    return "\n".join(lines[i:]) if i else text


def export_results(state: LessonState, output_dir: str = "") -> List[str]:
    """把本次实际执行且质量达标的产物落盘为真实文件（docx/xlsx/pptx/md）。
    返回生成的文件路径列表。output_dir 为空时使用「输出/备课资源包_时间戳_随机串」。
    规则：
    - 质量门判定失败（ok=False）或 LLM 错误占位的产物一律不落盘，并记录日志；
    - 目录延迟创建（有文件才建），全部失败时清理空目录，杜绝"空资源包"；
    - 每个资源包附带「质量检测报告.md」，列出各产物评分与问题。
    """
    tasks = (state.get("final") or {}).get("tasks") or state.get("tasks") or []
    if not tasks:
        return []
    base = output_dir or os.path.join(
        os.path.dirname(os.path.abspath(__file__)), "输出",
        "备课资源包_" + time.strftime("%Y%m%d_%H%M%S") + "_" + uuid.uuid4().hex[:6],
    )
    files: List[str] = []
    skipped: List[str] = []
    dir_ready = False
    for pid in tasks:
        spec = PRODUCERS.get(pid)
        if not spec:
            continue
        result = state.get(pid) or {}
        content = _strip_chat_prefix(result.get("content", ""))
        if not content or is_llm_error(content):
            skipped.append(f"{spec['name']}（模型调用失败，无有效内容）")
            continue
        if SETTINGS.quality_gate and result.get("ok") is False and not result.get("skipped"):
            reason = (result.get("error") or "质量检测未达标")
            skipped.append(f"{spec['name']}（{reason}）")
            log.warning("产物 %s 未通过质量门，拒绝落盘：%s", spec["name"], reason)
            continue
        if not dir_ready:
            os.makedirs(base, exist_ok=True)
            dir_ready = True
        fname = _safe_filename(spec["name"]) + "." + (spec.get("export") or "md")
        path = os.path.join(base, fname)
        ok = False
        ext = spec.get("export") or "md"
        if ext == "docx":
            ok = _export_docx(path, content)
        elif ext == "xlsx":
            ok = _export_xlsx(path, content)
        elif ext == "pptx":
            ok = _export_pptx(path, content, result.get("parsed"))
        if not ok:
            path = os.path.join(base, _safe_filename(spec["name"]) + ".md")
            try:
                with open(path, "w", encoding="utf-8") as f:
                    f.write(content)
                ok = os.path.getsize(path) > 0
            except Exception as e:
                log.error("md 落盘失败：%s", e)
                ok = False
        if ok and os.path.exists(path):
            files.append(path)
        else:
            skipped.append(f"{spec['name']}（文件写入失败）")

    # 质量检测报告随包落盘
    quality = (state.get("final") or {}).get("quality") or {}
    if dir_ready and (quality or skipped):
        lines = ["# 质量检测报告", "", f"生成时间：{time.strftime('%Y-%m-%d %H:%M:%S')}", ""]
        if quality:
            lines.append("## 产物评分")
            for name, q in quality.items():
                mark = "✅ 通过" if q.get("passed") else "❌ 未通过"
                lines.append(f"- {name}：{q.get('score', 0)} 分（{mark}）")
                for it in (q.get("issues") or [])[:5]:
                    lines.append(f"  - {it}")
            lines.append("")
        if skipped:
            lines.append("## 未交付产物")
            for s in skipped:
                lines.append(f"- {s}")
        try:
            rpt = os.path.join(base, "质量检测报告.md")
            with open(rpt, "w", encoding="utf-8") as f:
                f.write("\n".join(lines))
            files.append(rpt)
        except Exception as e:
            log.warning("质量报告写入失败：%s", e)

    # 全部失败：清理空目录，杜绝空资源包
    if not files and os.path.isdir(base):
        try:
            os.rmdir(base)
        except OSError:
            pass
        log.error("本次任务无任何达标产物，未生成资源包。跳过：%s", "；".join(skipped))
    else:
        log.info("产物已导出 %d 个文件到：%s（跳过 %d 个未达标产物）", len(files), base, len(skipped))
    return files


# ---------- WPS / 飞书 推送扩展点（可选，需配置 webhook） ----------
# 通过环境变量配置后，产物文件/摘要可自动推送到协同办公平台：
#   AGENT_WPS_WEBHOOK  /  AGENT_FEISHU_WEBHOOK
def push_to_wps(state: LessonState, files: Optional[List[str]] = None) -> bool:
    """把产物推送至 WPS 开放平台（需 AGENT_WPS_WEBHOOK 配置，未配置则跳过）。"""
    webhook = os.getenv("AGENT_WPS_WEBHOOK", "")
    if not webhook:
        log.info("未配置 AGENT_WPS_WEBHOOK，跳过 WPS 推送。")
        return False
    try:
        import requests

        files = files or []
        summary = ((state.get("final") or {}).get("summary") or "")[:500]
        requests.post(webhook, json={"summary": summary, "files": files}, timeout=10)
        log.info("已推送至 WPS：%d 个文件", len(files))
        return True
    except Exception as e:
        log.warning("WPS 推送失败：%s", e)
        return False


def push_to_feishu(state: LessonState, files: Optional[List[str]] = None) -> bool:
    """把产物推送至飞书群机器人（需 AGENT_FEISHU_WEBHOOK 配置，未配置则跳过）。"""
    webhook = os.getenv("AGENT_FEISHU_WEBHOOK", "")
    if not webhook:
        log.info("未配置 AGENT_FEISHU_WEBHOOK，跳过飞书推送。")
        return False
    try:
        import requests

        files = files or []
        summary = ((state.get("final") or {}).get("summary") or "")[:500]
        text = "【教师智能体产物】\n" + summary + "\n文件：\n" + "\n".join(files)
        requests.post(webhook, json={"msg_type": "text", "content": {"text": text}}, timeout=10)
        log.info("已推送至飞书：%d 个文件", len(files))
        return True
    except Exception as e:
        log.warning("飞书推送失败：%s", e)
        return False


# ===========================================================================
# 12. 图构建
# ===========================================================================
def build_graph():
    """构建 LangGraph 状态图（节点由产物注册表动态生成）。"""
    if not LANGGRAPH_AVAILABLE:
        return None
    g = StateGraph(LessonState)
    g.add_node("input", input_node)
    g.add_node("intent", intent_router)
    g.add_node("ask", ask_questions)
    g.add_node("assign", assign_globals)
    g.add_node("retrieve", retrieve_knowledge)
    g.add_node("web_search", web_search_node)
    g.add_node("decompose", task_decompose)
    for pid in PRODUCERS:
        g.add_node(pid, _node_for(pid))
    g.add_node("aggregate", aggregate)
    g.add_node("final", final_output)

    g.add_edge(START, "input")
    g.add_edge("input", "intent")
    g.add_conditional_edges("intent", is_complete, {"ask": "ask", "assign_globals": "assign"})
    # 信息不完整时：ask 生成追问后结束，由外部把答案追加到 history 后重新调用 run()
    g.add_edge("ask", END)
    g.add_edge("assign", "retrieve")
    g.add_edge("retrieve", "web_search")
    g.add_edge("web_search", "decompose")

    # 编排：decompose 后先执行 lesson_plan（不在任务中时节点内部快速跳过、无 LLM 开销），
    # 再扇出到其余产物并行执行，最后全部汇入 aggregate。
    # 好处：① lesson_plan 先行，PPT 等产物能从状态读取教案内容（"基于教案做 PPT"生效）；
    #      ② 所有产物→aggregate 的写入在同一 superstep 到达，aggregate 仅触发一次，
    #         避免 'final' 键并发写冲突（InvalidUpdateError）。
    g.add_edge("decompose", "lesson_plan")
    for pid in PRODUCERS:
        if pid == "lesson_plan":
            continue
        g.add_edge("lesson_plan", pid)
        g.add_edge(pid, "aggregate")
    g.add_edge("aggregate", "final")
    g.add_edge("final", END)
    return g.compile()


# ===========================================================================
# 13. 串行回退（无 LangGraph 时）
# ===========================================================================
def run_serial(state: LessonState) -> LessonState:
    # 各节点返回的是"局部更新 dict"（LangGraph 由图负责合并）；
    # 串行回退必须手动合并，否则状态会被单步返回值整体覆盖而丢失。
    def _step(st: LessonState, node: Callable[[LessonState], LessonState]) -> LessonState:
        upd = node(st)
        if isinstance(upd, dict):
            merged = dict(st)
            merged.update(upd)
            return merged
        return upd

    st = state
    st = _step(st, input_node)
    st = _step(st, intent_router)
    if not st.get("force_complete") and (
        not st.get("info_complete") or not st.get("parsed", {}).get("productsExplicit", True)
    ):
        st = _step(st, ask_questions)
        return st
    st = _step(st, assign_globals)
    st = _step(st, retrieve_knowledge)
    st = _step(st, web_search_node)
    st = _step(st, task_decompose)
    for pid in PRODUCERS:
        st = _step(st, _node_for(pid))
    st = _step(st, aggregate)
    return _step(st, final_output)


# ===========================================================================
# 14. 对外入口
# ===========================================================================
def prepare_graph():
    try:
        return build_graph()
    except Exception as e:
        log.error("图编译失败：%s，将使用串行回退。", e)
        return None


def run(
    teacher_input: str,
    history: Optional[List[Dict[str, str]]] = None,
    graph: Optional[Any] = None,
    force_defaults: bool = False,
) -> LessonState:
    """执行一次工作流。若信息不完整，返回含 pending_questions 的状态，
    调用方补充答案到 history 后再次调用。
    force_defaults=True：带默认值强制执行，不再追问（追问轮次用尽时的兜底）。"""
    graph = graph if graph is not None else prepare_graph()
    state: LessonState = {
        "raw_input": teacher_input,
        "history": history or [],
        "errors": [],
    }
    if force_defaults:
        state["force_complete"] = True
    try:
        if graph is not None:
            result = graph.invoke(state)
        else:
            result = run_serial(state)
        return dict(result)
    except Exception as e:
        return error_handler(state, e)


def run_interactive(teacher_input: str) -> LessonState:
    """带多轮追问的完整交互入口。"""
    graph = prepare_graph()
    history: List[Dict[str, str]] = []
    state = run(teacher_input, history=history, graph=graph)
    for _ in range(3):
        if not state.get("pending_questions"):
            break
        questions = state.get("pending_questions", [])
        answers = {}
        for q in questions:
            is_multi = q.get("multiSelect", False)
            opts = q.get("options") or []
            hint = f"（可多选，用、或逗号分隔：{'、'.join(opts)}）" if is_multi and opts else ""
            prompt = q.get("question", "请补充信息")
            ans = input(f"[追问] {prompt}{hint} （回车跳过）: ").strip()
            if ans:
                if is_multi:
                    picked = [o for o in opts if o in ans] or [ans]
                    answers[q.get("field", "products")] = "、".join(picked)
                else:
                    answers[q.get("field", "answer")] = ans
        if not answers:
            break
        history.append({"role": "user", "content": f"补充：{json.dumps(answers, ensure_ascii=False)}"})
        state = run(teacher_input, history=history, graph=graph)
    # 追问轮次用尽仍未补齐：带默认值强制执行，避免零产物交付
    if state.get("pending_questions"):
        log.info("追问轮次用尽，带默认值强制执行")
        state = run(teacher_input, history=history, graph=graph, force_defaults=True)
        state["assumed_defaults"] = True
    return state


# ===========================================================================
# 15. CLI
# ===========================================================================
def _print_result(state: LessonState) -> None:
    final = state.get("final", {})
    print("\n" + "=" * 60)
    print("【产物资源包汇总】")
    print("=" * 60)
    tasks = final.get("tasks") or state.get("tasks") or []
    if tasks:
        print("[本次执行任务]", "、".join(PRODUCERS[t]["name"] for t in tasks))
    print(final.get("summary", "(无汇总)"))
    results = final.get("results", {})
    for name, content in results.items():
        print(f"\n----- {name} -----")
        print(content[:1200])
    if state.get("errors"):
        print("\n[警告] 流程中出现异常：", state["errors"])
    if state.get("pending_questions"):
        print("\n[提示] 信息不完整，待补充：")
        for q in state["pending_questions"]:
            print(" -", q.get("question"))


def _save_and_notify(state: LessonState) -> None:
    """导出产物文件并尝试推送协同平台（WPS/飞书）。"""
    if not state.get("final"):
        return
    files = export_results(state)
    if files:
        print("\n[产物文件] 已导出到：")
        for f in files:
            print(" -", f)
        push_to_wps(state, files)
        push_to_feishu(state, files)


def main() -> None:
    example = (
        "人教版高一数学必修一《函数的单调性》，生成教案、15页PPT、10道分层习题，"
        "班级中等生居多，配一个函数动态演示动画"
    )
    q = sys.argv[1] if len(sys.argv) > 1 else example
    print("教师需求：", q)
    if len(sys.argv) > 1:
        state = run(q)
        _print_result(state)
        _save_and_notify(state)
        return
    state = run_interactive(q)
    _print_result(state)
    _save_and_notify(state)


if __name__ == "__main__":
    main()
