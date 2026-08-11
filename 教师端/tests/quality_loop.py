# -*- coding: utf-8 -*-
"""
智能体产物质量自动化检测与循环优化
==================================
让智能体（workflow.py 全链路）实际生成六类产物：
  教案 / PPT / 教学报告（教研报告）/ 学情分析 / 习题 / 办公任务文档
对每类产物执行质量门校验（QUALITY_STANDARDS）+ 落盘校验（文件存在、大小、可打开），
输出结构化检测报告。任何一类不达标即判定本轮失败，供修复后重跑。

用法：
  D:\\Python\\Python312\\python.exe tests\\quality_loop.py            # 全部场景
  D:\\Python\\Python312\\python.exe tests\\quality_loop.py --only lesson_plan,ppt
环境：
  AGENT_LLM_PROVIDER/AGENT_LLM_MODEL/AGENT_LLM_API_KEY/AGENT_LLM_BASE_URL 配置真实模型；
  未配置时走 Mock（用于验证流程机械正确性，质量门预期会拦截敷衍产物）。
"""
from __future__ import annotations

import argparse
import json
import os
import sys
import time

BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, BASE)

import workflow  # noqa: E402

# ---------------------------------------------------------------------------
# 六类检测场景：名称 → (请求文本, 期望任务id列表)
# ---------------------------------------------------------------------------
SCENARIOS = [
    ("lesson_plan", "教案", "人教版高一数学必修一《函数的单调性》，生成一份完整教案，班级中等生居多", ["lesson_plan"]),
    ("ppt", "PPT课件", "为高一数学《函数的单调性》生成15页PPT课件", ["ppt"]),
    ("teaching_research", "教学报告", "撰写一份教研报告：高一数学组《函数的单调性》公开课听评课教研活动", ["teaching_research"]),
    ("learning_analysis", "学情分析", "为高一数学《函数的单调性》生成学情分析报告，班级中等生居多", ["learning_analysis"]),
    ("exercise", "分层习题", "为高一数学《函数的单调性》生成10道分层习题，含解析", ["exercise"]),
    ("office_doc", "办公任务文档", "撰写一份本学期的教学工作总结", ["work_summary"]),
]

# 落盘文件大小下限（字节）：docx/pptx/xlsx 结构开销 + 实质内容
# PPT 硬性要求 ≥ 1MB（图文结合，每页配图，用户硬性指标）
MIN_FILE_SIZE = {"docx": 15000, "pptx": 1024 * 1024, "xlsx": 4000, "md": 400}


def run_scenario(sid: str, name: str, query: str, expect_tasks: list) -> dict:
    """执行一个场景：跑工作流 → 质量校验 → 落盘 → 文件校验。"""
    t0 = time.time()
    report = {"id": sid, "name": name, "query": query, "passed": False, "issues": [],
              "products": [], "files": [], "elapsed": 0.0}

    state = workflow.run(query)
    report["elapsed"] = round(time.time() - t0, 1)

    if state.get("pending_questions"):
        report["issues"].append("工作流进入追问分支（意图信息不足），未直接执行")
    tasks = state.get("tasks") or []
    report["tasks"] = tasks
    for et in expect_tasks:
        if et not in tasks:
            report["issues"].append(f"任务路由缺失：期望 {et}，实际 {tasks}")

    # 逐产物质量门
    for pid in expect_tasks:
        res = state.get(pid) or {}
        content = res.get("content") or ""
        prod = {"id": pid, "chars": len(content), "ok": res.get("ok"), "rounds": res.get("rounds")}
        if not content or workflow.is_llm_error(content):
            prod["verdict"] = "FAIL: 无有效内容（模型调用失败）"
            report["issues"].append(f"{pid}: 模型调用失败，无有效内容")
        else:
            v = workflow.validate_product(pid, content, res.get("parsed"))
            prod["score"] = v["score"]
            prod["passed"] = v["passed"]
            prod["issues"] = v["issues"]
            if not v["passed"]:
                report["issues"].append(f"{pid} 质量不达标({v['score']}分): " + "；".join(v["issues"][:3]))
        report["products"].append(prod)

    # 落盘 + 文件级校验
    out_dir = os.path.join(BASE, "备课产物", "质量检测_" + sid + "_" + time.strftime("%H%M%S"))
    files = workflow.export_results(state, output_dir=out_dir)
    for f in files:
        if not os.path.exists(f):
            report["issues"].append(f"文件缺失：{f}")
            continue
        size = os.path.getsize(f)
        ext = os.path.splitext(f)[1].lstrip(".").lower()
        entry = {"path": f, "ext": ext, "size": size}
        minimum = MIN_FILE_SIZE.get(ext, 200)
        if os.path.basename(f) != "质量检测报告.md" and size < minimum:
            entry["verdict"] = f"FAIL: 文件过小（{size} < {minimum}），疑似空壳"
            report["issues"].append(f"文件 {os.path.basename(f)} 过小（{size}B）")
        # 可打开性校验
        try:
            if ext == "docx":
                from docx import Document
                d = Document(f)
                entry["paragraphs"] = len(d.paragraphs)
            elif ext == "pptx":
                from pptx import Presentation
                p = Presentation(f)
                entry["slides"] = len(p.slides)
                if len(p.slides) < 12:
                    entry["verdict"] = f"FAIL: PPT 仅 {len(p.slides)} 页"
                    report["issues"].append(f"PPT 页数不足：{len(p.slides)} < 12")
            elif ext == "xlsx":
                from openpyxl import load_workbook
                wb = load_workbook(f)
                entry["sheets"] = len(wb.sheetnames)
        except Exception as e:
            entry["verdict"] = f"FAIL: 文件无法打开：{e}"
            report["issues"].append(f"文件 {os.path.basename(f)} 损坏：{e}")
        report["files"].append(entry)

    if not files:
        report["issues"].append("无任何文件落盘")

    report["passed"] = not report["issues"]
    return report


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--only", default="", help="只跑指定场景，逗号分隔")
    ap.add_argument("--json", default="", help="报告输出路径")
    args = ap.parse_args()

    only = [s for s in args.only.split(",") if s] if args.only else None
    print("=" * 70)
    print("智能体产物质量自动化检测")
    print(f"模型提供方：{workflow.SETTINGS.provider}（{'真实模型' if workflow.SETTINGS.provider != 'mock' else 'Mock 离线模式'}）")
    print(f"质量门：{'启用' if workflow.SETTINGS.quality_gate else '关闭'}，合格线 {workflow.SETTINGS.quality_min_score} 分")
    print("=" * 70)

    reports = []
    for sid, name, query, expect in SCENARIOS:
        if only and sid not in only:
            continue
        print(f"\n[{sid}] {name}：{query[:40]}…")
        rep = run_scenario(sid, name, query, expect)
        reports.append(rep)
        status = "✅ 通过" if rep["passed"] else "❌ 未通过"
        print(f"  {status}（耗时 {rep['elapsed']}s）")
        for p in rep["products"]:
            print(f"    产物 {p['id']}: {p.get('chars', 0)} 字, 评分 {p.get('score', '-')}, "
                  f"{'合格' if p.get('passed') else '不合格'}, 重写 {p.get('rounds', 1)} 轮")
        for f in rep["files"]:
            extra = f" slides={f['slides']}" if "slides" in f else ""
            print(f"    文件 {os.path.basename(f['path'])}: {f['size']}B{extra} {f.get('verdict', '')}")
        for it in rep["issues"]:
            print(f"    ⚠ {it}")

    passed = sum(1 for r in reports if r["passed"])
    print("\n" + "=" * 70)
    print(f"总结：{passed}/{len(reports)} 个场景通过")
    print("=" * 70)

    out = args.json or os.path.join(BASE, "tests", "quality_report.json")
    with open(out, "w", encoding="utf-8") as fh:
        json.dump({"time": time.strftime("%Y-%m-%d %H:%M:%S"),
                   "provider": workflow.SETTINGS.provider,
                   "passed": passed, "total": len(reports), "reports": reports},
                  fh, ensure_ascii=False, indent=2)
    print(f"报告已写入：{out}")
    return 0 if passed == len(reports) else 1


if __name__ == "__main__":
    sys.exit(main())
