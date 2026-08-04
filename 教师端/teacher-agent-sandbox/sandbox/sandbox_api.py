#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
教师智能体沙箱 · 服务器端 API 服务
==================================
基于 FastAPI 提供 Python 代码执行、文档生成（docx/pptx/xlsx/pdf）、
格式转换（pandoc/libreoffice）、OCR（tesseract）、图表生成（matplotlib）等能力。

浏览器端通过 sandbox-runtime.js 使用 Pyodide 实现等价功能，
本服务用于服务端真实环境，功能更完整（支持系统级工具调用）。

启动方式：
    uvicorn sandbox_api:app --host 0.0.0.0 --port 8000 --reload
或：
    python sandbox_api.py

对应目录结构（workspace 下）：
    workspace/
    ├── 教案/      .docx .pdf
    ├── 课件/      .pptx
    ├── 量规/      .xlsx
    └── 临时/      图表等临时文件
"""

# ============================================================
# 1. Imports
# ============================================================
import os
import sys
import io
import re
import json
import base64
import shutil
import signal
import subprocess
import tempfile
import traceback
import threading
import multiprocessing
from pathlib import Path
from datetime import datetime
from typing import Any, Optional, List, Dict, Union

# FastAPI
from fastapi import FastAPI, HTTPException, Request, Response, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, JSONResponse, StreamingResponse, PlainTextResponse
from pydantic import BaseModel, Field

# ---- 文档生成库 ----
# python-docx
from docx import Document as DocxDocument
from docx.shared import Pt, RGBColor, Inches, Cm, Emu
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml.ns import qn, nsdecls
from docx.oxml import parse_xml

# python-pptx
from pptx import Presentation as PptxPresentation
from pptx.util import Inches as PptxInches, Pt as PptxPt, Emu as PptxEmu
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# openpyxl
from openpyxl import Workbook
from openpyxl.styles import Font as XlFont, Alignment, PatternFill, Border, Side
from openpyxl.utils import get_column_letter

# reportlab（PDF 降级方案）
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import cm, mm
from reportlab.lib.colors import HexColor
from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_JUSTIFY
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle,
    PageBreak, KeepTogether,
)
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.platypus.flowables import HRFlowable

# matplotlib（图表）
import matplotlib
matplotlib.use("Agg")  # 无显示环境
import matplotlib.pyplot as plt
import matplotlib.font_manager as fm
import numpy as np

# ---- 可选库（降级处理） ----
try:
    from weasyprint import HTML as WeasyHTML
    HAS_WEASYPRINT = True
except Exception:
    HAS_WEASYPRINT = False

try:
    import markdown as md_lib
    HAS_MARKDOWN = True
except Exception:
    HAS_MARKDOWN = False

try:
    import pdfplumber
    HAS_PDFPLUMBER = True
except Exception:
    HAS_PDFPLUMBER = False


# ============================================================
# 2. 配置常量
# ============================================================

# 工作区根目录：sandbox/workspace/
WORKSPACE_DIR = Path(__file__).resolve().parent / "workspace"
WORKSPACE_DIR.mkdir(parents=True, exist_ok=True)

# 工作区子目录
DIRS = {
    "教案": WORKSPACE_DIR / "教案",
    "课件": WORKSPACE_DIR / "课件",
    "量规": WORKSPACE_DIR / "量规",
    "临时": WORKSPACE_DIR / "临时",
}
for _d in DIRS.values():
    _d.mkdir(parents=True, exist_ok=True)

# 请求体大小限制（50MB）
MAX_REQUEST_SIZE = 50 * 1024 * 1024

# 代码执行默认超时（秒）
DEFAULT_TIMEOUT = 30

# 代码执行最大超时（秒）
MAX_TIMEOUT = 120

# 主题色（对齐教师端 ink/jade/tan）
THEME = {
    "jade": "4F7A66",       # 竹青主色
    "jade_rgb": (79, 122, 102),
    "ink": "403A30",        # 暖墨色
    "ink_rgb": (64, 58, 48),
    "bg": "FAF8F3",         # 米白背景
    "bg_rgb": (250, 248, 243),
    "tan": "A8814E",        # 秋香赭点缀
    "tan_rgb": (168, 129, 78),
    "light_jade": "E8F0EC",
}

# 中文字体优先级
CN_FONTS = ["Microsoft YaHei", "SimSun", "Noto Sans CJK SC", "WenQuanYi Micro Hei"]

# 系统工具检测
SYSTEM_TOOLS = {}


# ============================================================
# 3. 安全工具函数
# ============================================================

def detect_system_tools():
    """检测系统级工具是否可用（pandoc / libreoffice / tesseract）"""
    tools = {}
    for name, cmd in [("pandoc", "pandoc"), ("libreoffice", "libreoffice"), 
                       ("soffice", "soffice"), ("tesseract", "tesseract")]:
        try:
            result = subprocess.run(
                [cmd, "--version"] if cmd != "libreoffice" and cmd != "soffice" else [cmd, "--version"],
                capture_output=True, text=True, timeout=5,
            )
            tools[name] = {
                "available": result.returncode == 0,
                "version": result.stdout.split("\n")[0] if result.stdout else "",
            }
        except Exception:
            tools[name] = {"available": False, "version": ""}
    # libreoffice/soffice 任一可用即可
    tools["libreoffice"] = {
        "available": tools.get("libreoffice", {}).get("available", False) or tools.get("soffice", {}).get("available", False),
        "version": tools.get("libreoffice", {}).get("version") or tools.get("soffice", {}).get("version", ""),
    }
    return tools


def validate_workspace_path(path: Union[str, Path]) -> Path:
    """
    验证路径在 workspace 目录内，防止路径穿越攻击。
    返回解析后的绝对 Path 对象。
    """
    target = (WORKSPACE_DIR / path).resolve() if not Path(path).is_absolute() else Path(path).resolve()
    workspace_resolved = WORKSPACE_DIR.resolve()
    if not str(target).startswith(str(workspace_resolved)):
        raise HTTPException(status_code=403, detail=f"路径越权：{path} 不在 workspace 目录内")
    return target


def safe_filename(name: str, ext: str) -> str:
    """生成安全文件名：去除非法字符，确保扩展名"""
    s = str(name or "文档").strip()
    # 去除已有扩展名
    s = re.sub(r"\.(docx?|pptx?|xlsx?|pdf|md|txt|html?|csv|png|jpe?g|gif)$", "", s, flags=re.IGNORECASE)
    # 去除 Windows 非法字符
    s = re.sub(r'[\\/:*?"<>|]', "_", s)
    s = s.strip(". ")
    if not s:
        s = "未命名"
    return f"{s}.{ext.lstrip('.')}"


def get_dir_by_ext(ext: str) -> Path:
    """根据扩展名返回默认保存目录"""
    ext_lower = ext.lower().lstrip(".")
    mapping = {
        "docx": DIRS["教案"],
        "pdf": DIRS["教案"],
        "pptx": DIRS["课件"],
        "xlsx": DIRS["量规"],
        "png": DIRS["临时"],
        "jpg": DIRS["临时"],
        "csv": DIRS["临时"],
    }
    return mapping.get(ext_lower, DIRS["临时"])


def file_download_url(rel_path: str) -> str:
    """生成文件下载 URL"""
    return f"/files/{rel_path}"


# ---- RestrictedExec：受限 Python 代码执行 ----

# 危险模块/属性黑名单
DANGEROUS_NAMES = {
    "os", "sys", "subprocess", "shutil", "ctypes", "multiprocessing",
    "socket", "http", "urllib", "requests", "pickle", "marshal",
    "importlib", "builtins", "globals", "locals", "exec", "eval",
    "compile", "open", "exit", "quit", "input",
    "__import__", "__builtins__", "globals", "locals",
}

# 允许使用的模块白名单（教学/数据处理相关）
ALLOWED_MODULES = {
    "math", "random", "statistics", "itertools", "functools",
    "datetime", "time", "re", "json", "collections", "copy",
    "string", "textwrap", "operator", "enum", "typing",
    "dataclasses", "abc", "decimal", "fractions",
    "numpy", "np", "pandas", "pd", "matplotlib", "plt",
    "sympy", "sp", "jieba", "scipy",
    "io", "base64", "hashlib",
}


class RestrictedExec:
    """
    受限的 Python 代码执行器。
    - 限制危险内置函数和模块
    - 捕获 stdout / stderr
    - 超时控制（通过线程）
    """

    @staticmethod
    def _build_globals(workspace: Path) -> dict:
        """构建受限的全局命名空间"""
        # 安全的内置函数子集
        safe_builtins = {
            # 基本类型
            "int": int, "float": float, "str": str, "bool": bool,
            "list": list, "tuple": tuple, "dict": dict, "set": set,
            "frozenset": frozenset, "bytes": bytes, "bytearray": bytearray,
            "complex": complex, "range": range, "type": type, "object": object,
            # 数学相关
            "abs": abs, "round": round, "min": min, "max": max,
            "sum": sum, "pow": pow, "divmod": divmod, "len": len,
            "sorted": sorted, "reversed": reversed, "enumerate": enumerate,
            "zip": zip, "map": map, "filter": filter, "any": any, "all": all,
            # 转换
            "hex": hex, "oct": oct, "bin": bin, "chr": chr, "ord": ord,
            "repr": repr, "format": format, "ascii": ascii,
            # 属性访问
            "getattr": getattr, "hasattr": hasattr, "setattr": setattr,
            "delattr": delattr, "isinstance": isinstance, "issubclass": issubclass,
            "callable": callable, "dir": dir, "vars": vars,
            # 迭代
            "iter": iter, "next": next, "slice": slice,
            # 其他
            "print": print, "id": id, "hash": hash,
            "True": True, "False": False, "None": None,
            "Exception": Exception, "ValueError": ValueError,
            "TypeError": TypeError, "KeyError": KeyError,
            "IndexError": IndexError, "AttributeError": AttributeError,
            "StopIteration": StopIteration, "RuntimeError": RuntimeError,
            "ZeroDivisionError": ZeroDivisionError, "NameError": NameError,
            "ImportError": ImportError, "FileNotFoundError": FileNotFoundError,
            "NotImplementedError": NotImplementedError, "AssertionError": AssertionError,
            "ArithmeticError": ArithmeticError, "LookupError": LookupError,
            "OverflowError": OverflowError, "Warning": Warning,
            "DeprecationWarning": DeprecationWarning,
        }

        g = {
            "__builtins__": safe_builtins,
            "__name__": "__sandbox__",
            "WORKSPACE": str(workspace),
        }

        # 预导入允许的库
        try:
            g["math"] = __import__("math")
        except Exception:
            pass
        try:
            g["random"] = __import__("random")
        except Exception:
            pass
        try:
            g["statistics"] = __import__("statistics")
        except Exception:
            pass
        try:
            g["itertools"] = __import__("itertools")
        except Exception:
            pass
        try:
            g["functools"] = __import__("functools")
        except Exception:
            pass
        try:
            g["datetime"] = __import__("datetime")
        except Exception:
            pass
        try:
            g["re"] = __import__("re")
        except Exception:
            pass
        try:
            g["json"] = __import__("json")
        except Exception:
            pass
        try:
            g["collections"] = __import__("collections")
        except Exception:
            pass
        try:
            import numpy
            g["numpy"] = numpy
            g["np"] = numpy
        except Exception:
            pass
        try:
            import pandas
            g["pandas"] = pandas
            g["pd"] = pandas
        except Exception:
            pass
        try:
            import matplotlib
            matplotlib.use("Agg")
            g["matplotlib"] = matplotlib
            import matplotlib.pyplot as plt_mod
            g["plt"] = plt_mod
        except Exception:
            pass
        try:
            import sympy
            g["sympy"] = sympy
            g["sp"] = sympy
        except Exception:
            pass
        try:
            import jieba
            g["jieba"] = jieba
        except Exception:
            pass
        try:
            import scipy
            g["scipy"] = scipy
        except Exception:
            pass

        return g

    @staticmethod
    def execute(code: str, timeout: int = DEFAULT_TIMEOUT, packages: Optional[List[str]] = None) -> dict:
        """
        执行 Python 代码，返回结果字典。
        
        返回:
            {
                "ok": bool,
                "stdout": str,
                "stderr": str,
                "result": str,
                "error": Optional[str],
                "files": list,
            }
        """
        # 可选包安装
        install_msgs = []
        if packages:
            for pkg in packages:
                try:
                    subprocess.run(
                        [sys.executable, "-m", "pip", "install", pkg, "--quiet"],
                        capture_output=True, text=True, timeout=60,
                    )
                    install_msgs.append(f"已安装: {pkg}")
                except Exception as e:
                    install_msgs.append(f"安装失败 {pkg}: {e}")

        # 捕获 stdout/stderr
        old_stdout = sys.stdout
        old_stderr = sys.stderr
        stdout_buf = io.StringIO()
        stderr_buf = io.StringIO()
        
        result_container = {"value": None, "error": None, "files": []}
        
        def _run():
            sys.stdout = stdout_buf
            sys.stderr = stderr_buf
            try:
                g = RestrictedExec._build_globals(WORKSPACE_DIR)
                # 执行代码
                exec(compile(code, "<sandbox>", "exec"), g)
                # 尝试获取最后一个表达式的值
                result_container["value"] = None
            except Exception as e:
                result_container["error"] = traceback.format_exc()
            finally:
                sys.stdout = old_stdout
                sys.stderr = old_stderr
                # 检查 workspace 下新生成的文件
                result_container["files"] = _collect_new_files()

        def _collect_new_files():
            """收集 workspace 下的文件列表"""
            files = []
            try:
                for root in DIRS.values():
                    if root.exists():
                        for f in root.rglob("*"):
                            if f.is_file():
                                rel = f.relative_to(WORKSPACE_DIR)
                                files.append({
                                    "path": str(rel).replace("\\", "/"),
                                    "size": f.stat().st_size,
                                    "url": file_download_url(str(rel).replace("\\", "/")),
                                })
            except Exception:
                pass
            return files

        # 使用线程执行，支持超时
        thread = threading.Thread(target=_run, daemon=True)
        thread.start()
        thread.join(timeout=timeout)

        if thread.is_alive():
            # 超时
            return {
                "ok": False,
                "stdout": stdout_buf.getvalue(),
                "stderr": stderr_buf.getvalue(),
                "result": "",
                "error": f"代码执行超时（{timeout}秒）",
                "files": [],
            }

        stdout_val = stdout_buf.getvalue()
        stderr_val = stderr_buf.getvalue()
        
        # 加上安装消息
        if install_msgs:
            stdout_val = "\n".join(install_msgs) + "\n" + stdout_val

        error_val = result_container["error"]
        
        return {
            "ok": error_val is None,
            "stdout": stdout_val,
            "stderr": stderr_val,
            "result": str(result_container["value"]) if result_container["value"] is not None else "",
            "error": error_val,
            "files": result_container["files"],
        }


# ============================================================
# 4. 文档生成函数
# ============================================================

# ---- 4.1 Markdown 解析 ----

def parse_markdown_to_blocks(content: str) -> list:
    """
    将 Markdown 文本解析为结构化块列表。
    支持：## 二级标题、### 三级标题、- 无序列表、| 表格 |、普通段落。
    """
    if not content or not content.strip():
        return []
    
    lines = content.split("\n")
    blocks = []
    i = 0
    
    while i < len(lines):
        line = lines[i].rstrip()
        
        # 空行
        if not line.strip():
            i += 1
            continue
        
        # 一级标题 # (作为 h2)
        m = re.match(r"^#\s+(.+)", line)
        if m:
            blocks.append({"type": "h2", "text": m.group(1).strip()})
            i += 1
            continue
        
        # 二级标题 ##
        m = re.match(r"^##\s+(.+)", line)
        if m:
            blocks.append({"type": "h2", "text": m.group(1).strip()})
            i += 1
            continue
        
        # 三级标题 ###
        m = re.match(r"^###\s+(.+)", line)
        if m:
            blocks.append({"type": "h3", "text": m.group(1).strip()})
            i += 1
            continue
        
        # 四级标题 ####
        m = re.match(r"^####\s+(.+)", line)
        if m:
            blocks.append({"type": "h3", "text": m.group(1).strip()})
            i += 1
            continue
        
        # 无序列表 - 或 *
        if re.match(r"^[-*]\s+", line):
            items = []
            while i < len(lines) and re.match(r"^[-*]\s+", lines[i].rstrip()):
                items.append(re.sub(r"^[-*]\s+", "", lines[i].rstrip()).strip())
                i += 1
            blocks.append({"type": "list", "items": items})
            continue
        
        # 有序列表 1. 2.
        if re.match(r"^\d+\.\s+", line):
            items = []
            while i < len(lines) and re.match(r"^\d+\.\s+", lines[i].rstrip()):
                items.append(re.sub(r"^\d+\.\s+", "", lines[i].rstrip()).strip())
                i += 1
            blocks.append({"type": "list", "items": items, "ordered": True})
            continue
        
        # 表格 | ... |
        if line.strip().startswith("|") and i + 1 < len(lines) and re.match(r"^\|[\s\-:|]+\|", lines[i + 1].strip()):
            table_rows = []
            # 表头
            table_rows.append(line.strip())
            i += 1
            # 分隔行
            i += 1
            # 数据行
            while i < len(lines) and lines[i].strip().startswith("|"):
                table_rows.append(lines[i].strip())
                i += 1
            blocks.append({"type": "table", "rows": table_rows})
            continue
        
        # 普通段落（连续非空行合并）
        para_lines = [line.strip()]
        i += 1
        while i < len(lines):
            next_line = lines[i].rstrip()
            if not next_line.strip():
                break
            if re.match(r"^#{1,4}\s+", next_line):
                break
            if re.match(r"^[-*]\s+", next_line):
                break
            if re.match(r"^\d+\.\s+", next_line):
                break
            if next_line.strip().startswith("|"):
                break
            para_lines.append(next_line.strip())
            i += 1
        blocks.append({"type": "text", "text": " ".join(para_lines)})
    
    return blocks


# ---- 4.2 生成 Word 文档 (.docx) ----

def set_cell_shading(cell, color_hex: str):
    """设置 docx 表格单元格背景色"""
    shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{color_hex}"/>')
    cell._tc.get_or_add_tcPr().append(shading)


def generate_docx(title: str, content: str, filename: Optional[str] = None) -> dict:
    """
    用 python-docx 生成 Word 文档。
    解析 Markdown：##标题、###副标题、-列表、|表格|、普通段落。
    设置中文字体（Microsoft YaHei / SimSun）。
    """
    doc = DocxDocument()
    
    # 设置默认字体
    style = doc.styles["Normal"]
    font = style.font
    font.name = "Microsoft YaHei"
    font.size = Pt(12)
    # 设置中文字体
    style.element.rPr.rFonts.set(qn("w:eastAsia"), "Microsoft YaHei")
    
    # 设置页边距
    for section in doc.sections:
        section.top_margin = Cm(2.5)
        section.bottom_margin = Cm(2.5)
        section.left_margin = Cm(2.5)
        section.right_margin = Cm(2.5)
    
    # 标题
    title_para = doc.add_heading(title, level=0)
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    for run in title_para.runs:
        run.font.name = "Microsoft YaHei"
        run.font.color.rgb = RGBColor(*THEME["jade_rgb"])
        run.font.size = Pt(22)
        run._element.rPr.rFonts.set(qn("w:eastAsia"), "Microsoft YaHei")
    
    # 分隔线
    doc.add_paragraph()
    
    # 解析 Markdown 内容
    blocks = parse_markdown_to_blocks(content)
    
    for block in blocks:
        if block["type"] == "h2":
            para = doc.add_heading(level=1)
            run = para.add_run(block["text"])
            run.font.name = "Microsoft YaHei"
            run.font.size = Pt(16)
            run.font.color.rgb = RGBColor(*THEME["jade_rgb"])
            run.bold = True
            run._element.rPr.rFonts.set(qn("w:eastAsia"), "Microsoft YaHei")
        
        elif block["type"] == "h3":
            para = doc.add_heading(level=2)
            run = para.add_run(block["text"])
            run.font.name = "Microsoft YaHei"
            run.font.size = Pt(14)
            run.font.color.rgb = RGBColor(*THEME["ink_rgb"])
            run.bold = True
            run._element.rPr.rFonts.set(qn("w:eastAsia"), "Microsoft YaHei")
        
        elif block["type"] == "list":
            ordered = block.get("ordered", False)
            for idx, item in enumerate(block["items"], 1):
                prefix = f"{idx}. " if ordered else "• "
                para = doc.add_paragraph()
                para.paragraph_format.left_indent = Cm(0.75)
                para.paragraph_format.space_after = Pt(4)
                run = para.add_run(prefix + item)
                run.font.name = "Microsoft YaHei"
                run.font.size = Pt(12)
                run._element.rPr.rFonts.set(qn("w:eastAsia"), "Microsoft YaHei")
        
        elif block["type"] == "table":
            rows_data = block["rows"]
            # 解析表格
            parsed = []
            for row_str in rows_data:
                cells = [c.strip() for c in row_str.strip("|").split("|")]
                parsed.append(cells)
            
            if len(parsed) >= 2:
                # 跳过分隔行（第2行），表头 + 数据行
                header = parsed[0]
                data_rows = parsed[2:] if len(parsed) > 2 else []
                
                table = doc.add_table(rows=1 + len(data_rows), cols=len(header))
                table.style = "Table Grid"
                table.alignment = WD_TABLE_ALIGNMENT.CENTER
                
                # 表头
                for ci, cell_text in enumerate(header):
                    cell = table.rows[0].cells[ci]
                    cell.text = cell_text
                    for para in cell.paragraphs:
                        para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                        for run in para.runs:
                            run.bold = True
                            run.font.name = "Microsoft YaHei"
                            run.font.size = Pt(11)
                            run.font.color.rgb = RGBColor(255, 255, 255)
                            run._element.rPr.rFonts.set(qn("w:eastAsia"), "Microsoft YaHei")
                    set_cell_shading(cell, THEME["jade"])
                
                # 数据行
                for ri, row_cells in enumerate(data_rows):
                    for ci, cell_text in enumerate(row_cells):
                        if ci < len(header):
                            cell = table.rows[ri + 1].cells[ci]
                            cell.text = cell_text
                            for para in cell.paragraphs:
                                for run in para.runs:
                                    run.font.name = "Microsoft YaHei"
                                    run.font.size = Pt(11)
                                    run._element.rPr.rFonts.set(qn("w:eastAsia"), "Microsoft YaHei")
                            # 隔行变色
                            if ri % 2 == 1:
                                set_cell_shading(cell, THEME["light_jade"])
                
                doc.add_paragraph()  # 表格后空行
        
        else:  # text
            para = doc.add_paragraph()
            para.paragraph_format.space_after = Pt(8)
            para.paragraph_format.line_spacing = 1.5
            run = para.add_run(block["text"])
            run.font.name = "Microsoft YaHei"
            run.font.size = Pt(12)
            run._element.rPr.rFonts.set(qn("w:eastAsia"), "Microsoft YaHei")
    
    # 保存
    fname = safe_filename(filename or title, "docx")
    save_dir = DIRS["教案"]
    save_path = save_dir / fname
    doc.save(str(save_path))
    
    rel_path = f"教案/{fname}"
    return {
        "filename": fname,
        "path": rel_path,
        "url": file_download_url(rel_path),
        "size": save_path.stat().st_size,
    }


# ---- 4.3 生成 PPT 课件 (.pptx) ----

def generate_pptx(slides_data: list, filename: Optional[str] = None) -> dict:
    """
    用 python-pptx 生成 PPT 课件。
    16:9 布局，主题色 jade=#4F7A66, ink=#403A30, bg=#FAF8F3。
    支持 cover / content 两种幻灯片类型。
    """
    prs = PptxPresentation()
    
    # 设置 16:9 尺寸
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)
    
    # 使用空白布局
    blank_layout = prs.slide_layouts[6]
    
    jade_color = PptxRGBColor.from_string(THEME["jade"])
    ink_color = PptxRGBColor.from_string(THEME["ink"])
    bg_color = PptxRGBColor.from_string(THEME["bg"])
    tan_color = PptxRGBColor.from_string(THEME["tan"])
    
    for slide_info in slides_data:
        slide = prs.slides.add_slide(blank_layout)
        slide_type = slide_info.get("type", "content")
        title = slide_info.get("title", "")
        subtitle = slide_info.get("subtitle", "")
        bullets = slide_info.get("bullets", [])
        content = slide_info.get("content", "")
        
        # 设置背景色
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
        
        if slide_type == "cover":
            # 封面页
            # 标题
            title_box = slide.shapes.add_textbox(
                PptxInches(0.5), PptxInches(2.4), PptxInches(12.3), PptxInches(1.2)
            )
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title
            p.alignment = PP_ALIGN.CENTER
            p.font.size = PptxPt(40)
            p.font.bold = True
            p.font.color.rgb = jade_color
            p.font.name = "Microsoft YaHei"
            
            # 副标题
            if subtitle:
                sub_box = slide.shapes.add_textbox(
                    PptxInches(0.5), PptxInches(3.8), PptxInches(12.3), PptxInches(0.8)
                )
                tf2 = sub_box.text_frame
                tf2.word_wrap = True
                p2 = tf2.paragraphs[0]
                p2.text = subtitle
                p2.alignment = PP_ALIGN.CENTER
                p2.font.size = PptxPt(20)
                p2.font.color.rgb = ink_color
                p2.font.name = "Microsoft YaHei"
            
            # 装饰线
            line_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                PptxInches(5.6), PptxInches(4.8), PptxInches(2.1), PptxInches(0.06)
            )
            line_shape.fill.solid()
            line_shape.fill.fore_color.rgb = tan_color
            line_shape.line.fill.background()
        
        else:
            # 内容页
            # 标题
            title_box = slide.shapes.add_textbox(
                PptxInches(0.5), PptxInches(0.3), PptxInches(12.3), PptxInches(0.8)
            )
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = PptxPt(28)
            p.font.bold = True
            p.font.color.rgb = jade_color
            p.font.name = "Microsoft YaHei"
            
            # 标题下装饰线
            line_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                PptxInches(0.5), PptxInches(1.15), PptxInches(1.6), PptxInches(0.05)
            )
            line_shape.fill.solid()
            line_shape.fill.fore_color.rgb = tan_color
            line_shape.line.fill.background()
            
            # 内容区域
            content_box = slide.shapes.add_textbox(
                PptxInches(0.6), PptxInches(1.4), PptxInches(12.1), PptxInches(5.6)
            )
            tf3 = content_box.text_frame
            tf3.word_wrap = True
            tf3.vertical_anchor = MSO_ANCHOR.TOP
            
            if bullets:
                # 项目符号列表
                for idx, bullet in enumerate(bullets):
                    if idx == 0:
                        p = tf3.paragraphs[0]
                    else:
                        p = tf3.add_paragraph()
                    p.text = f"•  {bullet}"
                    p.font.size = PptxPt(18)
                    p.font.color.rgb = ink_color
                    p.font.name = "Microsoft YaHei"
                    p.space_after = PptxPt(8)
            
            if content:
                # 纯文本内容
                if not bullets:
                    p = tf3.paragraphs[0] if tf3.paragraphs[0].text == "" else tf3.add_paragraph()
                else:
                    p = tf3.add_paragraph()
                p.text = content
                p.font.size = PptxPt(16)
                p.font.color.rgb = ink_color
                p.font.name = "Microsoft YaHei"
    
    # 保存
    fname = safe_filename(filename or "课件", "pptx")
    save_dir = DIRS["课件"]
    save_path = save_dir / fname
    prs.save(str(save_path))
    
    rel_path = f"课件/{fname}"
    return {
        "filename": fname,
        "path": rel_path,
        "url": file_download_url(rel_path),
        "size": save_path.stat().st_size,
        "slides": len(slides_data),
    }


# ---- 4.4 生成 Excel (.xlsx) ----

def generate_xlsx(rows: list, filename: Optional[str] = None, 
                  sheet_name: Optional[str] = None, dir_name: Optional[str] = None) -> dict:
    """
    用 openpyxl 生成 Excel。
    列宽自适应，表头加粗，保存到 workspace/量规/ 或指定目录。
    """
    wb = Workbook()
    ws = wb.active
    ws.title = sheet_name or "Sheet1"
    
    # 边框样式
    thin_border = Border(
        left=Side(style="thin", color="CCCCCC"),
        right=Side(style="thin", color="CCCCCC"),
        top=Side(style="thin", color="CCCCCC"),
        bottom=Side(style="thin", color="CCCCCC"),
    )
    header_fill = PatternFill(start_color=THEME["jade"], end_color=THEME["jade"], fill_type="solid")
    header_font = XlFont(name="Microsoft YaHei", size=11, bold=True, color="FFFFFF")
    body_font = XlFont(name="Microsoft YaHei", size=11)
    alt_fill = PatternFill(start_color=THEME["light_jade"], end_color=THEME["light_jade"], fill_type="solid")
    center_align = Alignment(horizontal="center", vertical="center", wrap_text=True)
    left_align = Alignment(horizontal="left", vertical="center", wrap_text=True)
    
    # 写入数据
    for ri, row in enumerate(rows):
        for ci, cell_value in enumerate(row):
            cell = ws.cell(row=ri + 1, column=ci + 1, value=cell_value)
            cell.border = thin_border
            cell.font = body_font
            cell.alignment = left_align
            
            if ri == 0:
                # 表头
                cell.font = header_font
                cell.fill = header_fill
                cell.alignment = center_align
            elif ri % 2 == 0:
                # 隔行变色
                cell.fill = alt_fill
    
    # 列宽自适应
    for col_idx in range(1, ws.max_column + 1):
        max_len = 0
        col_letter = get_column_letter(col_idx)
        for row_idx in range(1, ws.max_row + 1):
            cell = ws.cell(row=row_idx, column=col_idx)
            if cell.value:
                # 中文字符算2个宽度
                val_str = str(cell.value)
                length = sum(2 if ord(c) > 127 else 1 for c in val_str)
                if length > max_len:
                    max_len = length
        # 设置列宽（最小8，最大50）
        ws.column_dimensions[col_letter].width = max(8, min(max_len + 4, 50))
    
    # 冻结首行
    if len(rows) > 1:
        ws.freeze_panes = "A2"
    
    # 行高
    ws.row_dimensions[1].height = 28
    for ri in range(2, ws.max_row + 1):
        ws.row_dimensions[ri].height = 22
    
    # 保存
    fname = safe_filename(filename or "表格", "xlsx")
    save_dir = DIRS.get(dir_name, DIRS["量规"]) if dir_name else DIRS["量规"]
    save_dir.mkdir(parents=True, exist_ok=True)
    save_path = save_dir / fname
    wb.save(str(save_path))
    
    rel_dir = dir_name or "量规"
    rel_path = f"{rel_dir}/{fname}"
    return {
        "filename": fname,
        "path": rel_path,
        "url": file_download_url(rel_path),
        "size": save_path.stat().st_size,
        "rows": len(rows),
    }


# ---- 4.5 生成 PDF (.pdf) ----

def _register_cn_font_for_reportlab():
    """尝试注册中文字体到 reportlab"""
    font_candidates = [
        ("NotoSansCJK", "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc"),
        ("NotoSansCJK", "/usr/share/fonts/noto-cjk/NotoSansCJK-Regular.ttc"),
        ("SimSun", "C:/Windows/Fonts/simsun.ttc"),
        ("MSYaHei", "C:/Windows/Fonts/msyh.ttc"),
        ("MSYaHei", "C:/Windows/Fonts/msyh.ttf"),
        ("SimHei", "C:/Windows/Fonts/simhei.ttf"),
        ("WenQuanYi", "/usr/share/fonts/wenquanyi/wqy-microhei/wqy-microhei.ttc"),
    ]
    for font_name, font_path in font_candidates:
        if os.path.exists(font_path):
            try:
                pdfmetrics.registerFont(TTFont(font_name, font_path))
                return font_name
            except Exception:
                continue
    return None


def generate_pdf_weasyprint(content: str, title: str, save_path: Path) -> bool:
    """用 weasyprint 将 HTML/Markdown 转 PDF"""
    if not HAS_WEASYPRINT:
        return False
    
    try:
        # 如果是 Markdown，先转 HTML
        if HAS_MARKDOWN and not content.strip().startswith("<"):
            html_body = md_lib.markdown(content, extensions=["tables", "fenced_code"])
        else:
            html_body = content
        
        full_html = f"""
        <!DOCTYPE html>
        <html>
        <head>
        <meta charset="utf-8">
        <style>
            @page {{ size: A4; margin: 2.5cm; }}
            body {{ 
                font-family: "Microsoft YaHei", "SimSun", "Noto Sans CJK SC", sans-serif; 
                line-height: 1.8; color: #403A30; font-size: 12pt;
            }}
            h1 {{ color: #4F7A66; font-size: 22pt; text-align: center; 
                  border-bottom: 2px solid #4F7A66; padding-bottom: 10px; margin-bottom: 20px; }}
            h2 {{ color: #4F7A66; font-size: 16pt; margin-top: 24px; }}
            h3 {{ color: #2d5f4e; font-size: 14pt; margin-top: 18px; }}
            table {{ border-collapse: collapse; width: 100%; margin: 12px 0; }}
            th {{ background: #4F7A66; color: white; padding: 8px 12px; text-align: center; }}
            td {{ border: 1px solid #ddd; padding: 8px 12px; }}
            tr:nth-child(even) {{ background: #E8F0EC; }}
            ul, ol {{ margin-left: 24px; }}
            p {{ margin: 8px 0; }}
        </style>
        </head>
        <body>
        <h1>{title}</h1>
        {html_body}
        </body>
        </html>
        """
        WeasyHTML(string=full_html).write_pdf(str(save_path))
        return True
    except Exception as e:
        print(f"[sandbox] weasyprint 生成失败: {e}", file=sys.stderr)
        return False


def generate_pdf_reportlab(content: str, title: str, save_path: Path) -> bool:
    """用 reportlab 生成 PDF（降级方案，纯文本）"""
    try:
        cn_font = _register_cn_font_for_reportlab()
        
        doc = SimpleDocTemplate(
            str(save_path), pagesize=A4,
            topMargin=2 * cm, bottomMargin=2 * cm,
            leftMargin=2.5 * cm, rightMargin=2.5 * cm,
        )
        
        styles = getSampleStyleSheet()
        
        # 标题样式
        title_style = ParagraphStyle(
            "CustomTitle", parent=styles["Title"],
            fontSize=22, textColor=HexColor("#" + THEME["jade"]),
            alignment=TA_CENTER, spaceAfter=20, spaceBefore=10,
            fontName=cn_font or "Helvetica",
        )
        h2_style = ParagraphStyle(
            "CustomH2", parent=styles["Heading2"],
            fontSize=16, textColor=HexColor("#" + THEME["jade"]),
            spaceBefore=18, spaceAfter=8,
            fontName=cn_font or "Helvetica",
        )
        h3_style = ParagraphStyle(
            "CustomH3", parent=styles["Heading3"],
            fontSize=14, textColor=HexColor("#" + THEME["ink"]),
            spaceBefore=14, spaceAfter=6,
            fontName=cn_font or "Helvetica",
        )
        body_style = ParagraphStyle(
            "CustomBody", parent=styles["Normal"],
            fontSize=12, leading=20, alignment=TA_JUSTIFY,
            spaceAfter=8, fontName=cn_font or "Helvetica",
        )
        bullet_style = ParagraphStyle(
            "CustomBullet", parent=body_style,
            leftIndent=24, bulletIndent=12, spaceAfter=4,
        )
        
        story = []
        story.append(Paragraph(title, title_style))
        story.append(HRFlowable(width="100%", thickness=2, color=HexColor("#" + THEME["jade"])))
        story.append(Spacer(1, 12))
        
        # 解析 Markdown
        blocks = parse_markdown_to_blocks(content)
        
        for block in blocks:
            if block["type"] == "h2":
                story.append(Paragraph(block["text"], h2_style))
            elif block["type"] == "h3":
                story.append(Paragraph(block["text"], h3_style))
            elif block["type"] == "list":
                for item in block["items"]:
                    # 转义 XML 特殊字符
                    safe_item = item.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                    story.append(Paragraph(f"• {safe_item}", bullet_style))
            elif block["type"] == "table":
                # 解析表格
                rows_data = block["rows"]
                parsed = []
                for row_str in rows_data:
                    cells = [c.strip() for c in row_str.strip("|").split("|")]
                    parsed.append(cells)
                if len(parsed) >= 2:
                    header = parsed[0]
                    data_rows = parsed[2:] if len(parsed) > 2 else []
                    table_data = [header] + data_rows
                    # 转义
                    table_data = [[c.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;") for c in row] for row in table_data]
                    t = Table(table_data, repeatRows=1)
                    t.setStyle(TableStyle([
                        ("BACKGROUND", (0, 0), (-1, 0), HexColor("#" + THEME["jade"])),
                        ("TEXTCOLOR", (0, 0), (-1, 0), HexColor("#FFFFFF")),
                        ("FONTNAME", (0, 0), (-1, -1), cn_font or "Helvetica"),
                        ("FONTSIZE", (0, 0), (-1, 0), 11),
                        ("FONTSIZE", (0, 1), (-1, -1), 11),
                        ("ALIGN", (0, 0), (-1, 0), "CENTER"),
                        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                        ("GRID", (0, 0), (-1, -1), 0.5, HexColor("#CCCCCC")),
                        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [HexColor("#FFFFFF"), HexColor("#" + THEME["light_jade"])]),
                        ("TOPPADDING", (0, 0), (-1, -1), 6),
                        ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
                    ]))
                    story.append(t)
                    story.append(Spacer(1, 8))
            else:
                safe_text = block["text"].replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                story.append(Paragraph(safe_text, body_style))
        
        doc.build(story)
        return True
    except Exception as e:
        print(f"[sandbox] reportlab 生成失败: {e}", file=sys.stderr)
        return False


def generate_pdf(content: str, title: str, filename: Optional[str] = None) -> dict:
    """
    生成 PDF 文档。
    优先用 weasyprint（HTML→PDF，中文友好），降级用 reportlab（纯文本）。
    """
    fname = safe_filename(filename or title or "文档", "pdf")
    save_dir = DIRS["教案"]
    save_path = save_dir / fname
    
    success = False
    method = ""
    
    # 尝试 weasyprint
    if HAS_WEASYPRINT:
        success = generate_pdf_weasyprint(content, title, save_path)
        method = "weasyprint"
    
    # 降级 reportlab
    if not success:
        success = generate_pdf_reportlab(content, title, save_path)
        method = "reportlab"
    
    if not success:
        raise HTTPException(status_code=500, detail="PDF 生成失败：weasyprint 和 reportlab 均不可用或生成出错")
    
    rel_path = f"教案/{fname}"
    return {
        "filename": fname,
        "path": rel_path,
        "url": file_download_url(rel_path),
        "size": save_path.stat().st_size,
        "method": method,
    }


# ============================================================
# 5. 格式转换函数
# ============================================================

def convert_with_pandoc(input_path: Path, output_path: Path, from_fmt: str, to_fmt: str) -> bool:
    """使用 pandoc 进行文档格式转换"""
    try:
        cmd = ["pandoc", str(input_path), "-o", str(output_path), "-f", from_fmt, "-t", to_fmt]
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
        return result.returncode == 0
    except Exception as e:
        print(f"[sandbox] pandoc 转换失败: {e}", file=sys.stderr)
        return False


def convert_with_libreoffice(input_path: Path, output_path: Path, to_fmt: str) -> bool:
    """使用 libreoffice 进行文档格式转换"""
    try:
        # 确定输出目录和格式
        output_dir = output_path.parent
        output_dir.mkdir(parents=True, exist_ok=True)
        
        # libreoffice 格式映射
        fmt_map = {
            "pdf": "pdf",
            "docx": "docx",
            "pptx": "pptx",
            "xlsx": "xlsx",
            "html": "html",
        }
        lo_fmt = fmt_map.get(to_fmt.lower(), to_fmt.lower())
        
        # 尝试 libreoffice 或 soffice
        lo_cmd = "libreoffice" if SYSTEM_TOOLS.get("libreoffice", {}).get("available") else "soffice"
        cmd = [
            lo_cmd, "--headless", "--convert-to", lo_fmt,
            "--outdir", str(output_dir), str(input_path),
        ]
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
        
        if result.returncode == 0:
            # libreoffice 输出文件名是 input 文件名换扩展名
            expected_output = output_dir / (input_path.stem + "." + lo_fmt)
            if expected_output.exists() and expected_output != output_path:
                shutil.move(str(expected_output), str(output_path))
            return output_path.exists()
        return False
    except Exception as e:
        print(f"[sandbox] libreoffice 转换失败: {e}", file=sys.stderr)
        return False


def convert_document(from_format: str, to_format: str, file_content: str, filename: str) -> dict:
    """
    文档格式转换。
    支持 pandoc：md <-> docx <-> html <-> pdf
    支持 libreoffice：pptx->pdf, docx->pdf, xlsx->pdf
    file_content 为 base64 编码的文件内容。
    """
    # 解码文件内容
    try:
        file_bytes = base64.b64decode(file_content)
    except Exception:
        raise HTTPException(status_code=400, detail="file_content 不是有效的 base64 编码")
    
    # 创建临时文件
    with tempfile.TemporaryDirectory() as tmpdir:
        tmpdir_path = Path(tmpdir)
        input_path = tmpdir_path / filename
        input_path.write_bytes(file_bytes)
        
        # 确定输出文件名
        stem = Path(filename).stem
        output_filename = f"{stem}.{to_format}"
        output_path = tmpdir_path / output_filename
        
        success = False
        method = ""
        
        # pandoc 支持的格式
        pandoc_formats = {"md", "docx", "html", "pdf", "rst", "latex", "epub"}
        
        if from_format.lower() in pandoc_formats and to_format.lower() in pandoc_formats:
            if SYSTEM_TOOLS.get("pandoc", {}).get("available"):
                success = convert_with_pandoc(input_path, output_path, from_format, to_format)
                method = "pandoc"
        
        # libreoffice 转换（特别是 pptx/xlsx -> pdf）
        if not success:
            lo_formats = {"pdf", "docx", "pptx", "xlsx", "html"}
            if from_format.lower() in {"pptx", "docx", "xlsx", "odt", "ods", "odp"} and to_format.lower() in lo_formats:
                if SYSTEM_TOOLS.get("libreoffice", {}).get("available"):
                    success = convert_with_libreoffice(input_path, output_path, to_format)
                    method = "libreoffice"
        
        if not success:
            raise HTTPException(
                status_code=500,
                detail=f"格式转换失败：{from_format} -> {to_format}。"
                       f"pandoc 可用: {SYSTEM_TOOLS.get('pandoc', {}).get('available', False)}, "
                       f"libreoffice 可用: {SYSTEM_TOOLS.get('libreoffice', {}).get('available', False)}"
            )
        
        # 读取转换后的文件
        if not output_path.exists():
            raise HTTPException(status_code=500, detail="转换后文件未生成")
        
        output_bytes = output_path.read_bytes()
        output_b64 = base64.b64encode(output_bytes).decode("ascii")
        
        return {
            "filename": output_filename,
            "content": output_b64,
            "size": len(output_bytes),
            "method": method,
            "from_format": from_format,
            "to_format": to_format,
        }


# ---- 5.1 OCR 文字识别 ----

def run_ocr(image_b64: str, lang: str = "chi_sim") -> dict:
    """
    使用 tesseract 进行 OCR 文字识别。
    image_b64: base64 编码的图片内容。
    lang: 识别语言（chi_sim / eng / chi_sim+eng）。
    """
    if not SYSTEM_TOOLS.get("tesseract", {}).get("available"):
        raise HTTPException(
            status_code=503,
            detail="tesseract 未安装或不可用。请安装 tesseract-ocr 和对应语言包。"
        )
    
    try:
        image_bytes = base64.b64decode(image_b64)
    except Exception:
        raise HTTPException(status_code=400, detail="image 不是有效的 base64 编码")
    
    with tempfile.TemporaryDirectory() as tmpdir:
        tmpdir_path = Path(tmpdir)
        # 保存图片到临时文件
        input_path = tmpdir_path / "ocr_input.png"
        input_path.write_bytes(image_bytes)
        
        # 调用 tesseract
        output_base = tmpdir_path / "ocr_output"
        cmd = ["tesseract", str(input_path), str(output_base), "-l", lang]
        
        try:
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
            if result.returncode != 0:
                raise HTTPException(
                    status_code=500,
                    detail=f"OCR 识别失败: {result.stderr}"
                )
        except subprocess.TimeoutExpired:
            raise HTTPException(status_code=504, detail="OCR 识别超时")
        
        # 读取结果
        output_txt = output_base.with_suffix(".txt")
        if not output_txt.exists():
            raise HTTPException(status_code=500, detail="OCR 结果文件未生成")
        
        text = output_txt.read_text(encoding="utf-8")
        
        return {
            "text": text.strip(),
            "lang": lang,
            "char_count": len(text.strip()),
        }


# ---- 5.2 图表生成 ----

def setup_matplotlib_cn_font():
    """设置 matplotlib 中文字体"""
    font_candidates = [
        "Noto Sans CJK SC",
        "Microsoft YaHei",
        "SimHei",
        "SimSun",
        "WenQuanYi Micro Hei",
        "Arial Unicode MS",
    ]
    
    # 查找系统可用字体
    available_fonts = set(f.name for f in fm.fontManager.ttflist)
    
    for font_name in font_candidates:
        if font_name in available_fonts:
            plt.rcParams["font.sans-serif"] = [font_name]
            plt.rcParams["axes.unicode_minus"] = False
            return font_name
    
    # 没有找到中文字体，使用默认
    plt.rcParams["axes.unicode_minus"] = False
    return None


def generate_chart(chart_type: str, data: dict, title: Optional[str] = None, 
                   filename: Optional[str] = None) -> dict:
    """
    用 matplotlib 生成图表并保存为 PNG。
    支持：bar / line / pie / radar / scatter
    """
    setup_matplotlib_cn_font()
    
    # 主题色
    colors = ["#4F7A66", "#A8814E", "#7BA892", "#C4A062", "#3D6B56", "#D4A574", "#5E9882", "#B8956A"]
    
    fig, ax = plt.subplots(figsize=(10, 6), dpi=150)
    fig.patch.set_facecolor("#FAF8F3")
    ax.set_facecolor("#FAF8F3")
    
    chart_type = chart_type.lower()
    
    if chart_type == "bar":
        labels = data.get("labels", [])
        values = data.get("values", [])
        if not labels or not values:
            raise HTTPException(status_code=400, detail="柱状图需要 labels 和 values 数据")
        
        bars = ax.bar(labels, values, color=colors[:len(values)], edgecolor="white", linewidth=0.5)
        ax.set_ylabel("数值", fontsize=12, color="#403A30")
        ax.tick_params(axis="x", rotation=30, labelsize=10, colors="#403A30")
        ax.tick_params(axis="y", labelsize=10, colors="#403A30")
        
        # 在柱子上显示数值
        for bar in bars:
            height = bar.get_height()
            ax.text(bar.get_x() + bar.get_width() / 2., height,
                    f"{height:.1f}", ha="center", va="bottom", fontsize=9, color="#403A30")
    
    elif chart_type == "line":
        labels = data.get("labels", [])
        values = data.get("values", [])
        if not labels or not values:
            raise HTTPException(status_code=400, detail="折线图需要 labels 和 values 数据")
        
        ax.plot(labels, values, color="#4F7A66", linewidth=2.5, marker="o", 
                markersize=7, markerfacecolor="#A8814E", markeredgecolor="white")
        ax.fill_between(range(len(labels)), values, alpha=0.15, color="#4F7A66")
        ax.set_ylabel("数值", fontsize=12, color="#403A30")
        ax.tick_params(axis="x", rotation=30, labelsize=10, colors="#403A30")
        ax.tick_params(axis="y", labelsize=10, colors="#403A30")
        ax.grid(True, linestyle="--", alpha=0.3, color="#A8814E")
    
    elif chart_type == "pie":
        labels = data.get("labels", [])
        values = data.get("values", [])
        if not labels or not values:
            raise HTTPException(status_code=400, detail="饼图需要 labels 和 values 数据")
        
        wedges, texts, autotexts = ax.pie(
            values, labels=labels, autopct="%1.1f%%",
            colors=colors[:len(values)],
            startangle=90, textprops={"fontsize": 10, "color": "#403A30"},
            wedgeprops={"edgecolor": "white", "linewidth": 1.5},
        )
        for autotext in autotexts:
            autotext.set_color("white")
            autotext.set_fontsize(9)
            autotext.set_fontweight("bold")
        ax.axis("equal")
    
    elif chart_type == "radar":
        labels = data.get("labels", [])
        values = data.get("values", [])
        if not labels or not values:
            raise HTTPException(status_code=400, detail="雷达图需要 labels 和 values 数据")
        
        # 雷达图需要闭合
        n = len(labels)
        angles = np.linspace(0, 2 * np.pi, n, endpoint=False).tolist()
        values_closed = values + values[:1]
        angles_closed = angles + angles[:1]
        
        ax.remove()  # 移除原来的 ax
        ax = fig.add_subplot(111, polar=True)
        ax.set_facecolor("#FAF8F3")
        
        ax.fill(angles_closed, values_closed, alpha=0.25, color="#4F7A66")
        ax.plot(angles_closed, values_closed, color="#4F7A66", linewidth=2.5, marker="o", 
                markersize=7, markerfacecolor="#A8814E")
        ax.set_xticks(angles)
        ax.set_xticklabels(labels, fontsize=10, color="#403A30")
        ax.tick_params(axis="y", labelsize=9, colors="#A8814E")
        ax.grid(True, linestyle="--", alpha=0.3)
    
    elif chart_type == "scatter":
        x = data.get("x", [])
        y = data.get("y", [])
        if not x or not y:
            raise HTTPException(status_code=400, detail="散点图需要 x 和 y 数据")
        
        ax.scatter(x, y, c="#4F7A66", s=80, alpha=0.7, edgecolors="white", linewidth=0.5)
        ax.set_xlabel("X", fontsize=12, color="#403A30")
        ax.set_ylabel("Y", fontsize=12, color="#403A30")
        ax.tick_params(axis="both", labelsize=10, colors="#403A30")
        ax.grid(True, linestyle="--", alpha=0.3, color="#A8814E")
    
    else:
        raise HTTPException(status_code=400, detail=f"不支持的图表类型: {chart_type}，支持: bar/line/pie/radar/scatter")
    
    # 设置标题
    if title:
        ax.set_title(title, fontsize=16, color="#4F7A66", fontweight="bold", pad=20)
    
    # 美化边框
    for spine in ax.spines.values():
        spine.set_color("#A8814E")
        spine.set_linewidth(0.5)
    
    plt.tight_layout()
    
    # 保存
    fname = safe_filename(filename or title or "图表", "png")
    save_dir = DIRS["临时"]
    save_path = save_dir / fname
    fig.savefig(str(save_path), dpi=150, bbox_inches="tight", 
                facecolor=fig.get_facecolor(), edgecolor="none")
    plt.close(fig)
    
    rel_path = f"临时/{fname}"
    return {
        "filename": fname,
        "path": rel_path,
        "url": file_download_url(rel_path),
        "size": save_path.stat().st_size,
        "chart_type": chart_type,
    }


# ============================================================
# 6. Pydantic 模型
# ============================================================

class HealthResponse(BaseModel):
    status: str
    timestamp: str
    python_version: str
    packages_status: dict
    system_tools_status: dict
    workspace: str


class ExecuteRequest(BaseModel):
    code: str = Field(..., description="要执行的 Python 代码")
    packages: Optional[List[str]] = Field(None, description="可选安装的 pip 包列表")
    timeout: Optional[int] = Field(DEFAULT_TIMEOUT, description="超时时间（秒），默认30")


class ExecuteResponse(BaseModel):
    ok: bool
    stdout: str
    stderr: str
    result: str
    error: Optional[str] = None
    files: Optional[List[dict]] = None


class DocxRequest(BaseModel):
    title: str = Field(..., description="文档标题")
    content: str = Field("", description="Markdown 格式的文档内容")
    filename: Optional[str] = Field(None, description="文件名（不含扩展名）")


class PptxRequest(BaseModel):
    filename: Optional[str] = Field(None, description="文件名（不含扩展名）")
    slides: List[dict] = Field(..., description="幻灯片列表，每项含 type/title/subtitle/bullets/content")


class XlsxRequest(BaseModel):
    filename: Optional[str] = Field(None, description="文件名（不含扩展名）")
    rows: List[List[Any]] = Field(..., description="二维数组，每行为一行数据")
    sheetName: Optional[str] = Field(None, description="工作表名称")
    dir: Optional[str] = Field(None, description="保存目录（教案/课件/量规/临时）")


class PdfRequest(BaseModel):
    filename: Optional[str] = Field(None, description="文件名（不含扩展名）")
    content: str = Field(..., description="Markdown 或 HTML 格式的内容")
    title: Optional[str] = Field(None, description="文档标题")


class ConvertRequest(BaseModel):
    from_format: str = Field(..., description="源格式：md/docx/html/pptx/xlsx 等")
    to_format: str = Field(..., description="目标格式：md/docx/html/pdf 等")
    file_content: str = Field(..., description="base64 编码的文件内容")
    filename: str = Field(..., description="源文件名")


class OcrRequest(BaseModel):
    image: str = Field(..., description="base64 编码的图片内容")
    lang: Optional[str] = Field("chi_sim", description="识别语言：chi_sim / eng / chi_sim+eng")


class ChartRequest(BaseModel):
    type: str = Field(..., description="图表类型：bar/line/pie/radar/scatter")
    data: dict = Field(..., description="图表数据")
    title: Optional[str] = Field(None, description="图表标题")
    filename: Optional[str] = Field(None, description="文件名（不含扩展名）")


class FetchRequest(BaseModel):
    url: str = Field(..., description="要抓取的网页 URL")
    max_chars: int = Field(4000, description="正文最大字符数")


# ============================================================
# 7. FastAPI 应用 + 路由
# ============================================================

app = FastAPI(
    title="教师智能体沙箱 API",
    description="提供 Python 代码执行、文档生成、格式转换、OCR、图表生成等服务端能力",
    version="1.0.0",
)

# CORS 中间件（开发环境允许所有源）
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# 请求体大小限制中间件
@app.middleware("http")
async def limit_request_size(request: Request, call_next):
    """限制请求体大小为 50MB"""
    content_length = request.headers.get("content-length")
    if content_length and int(content_length) > MAX_REQUEST_SIZE:
        return JSONResponse(
            status_code=413,
            content={"detail": f"请求体过大，最大允许 {MAX_REQUEST_SIZE // (1024 * 1024)}MB"}
        )
    return await call_next(request)


# ---- 7.1 健康检查 ----

@app.get("/health", response_model=HealthResponse, summary="健康检查")
async def health_check():
    """返回服务健康状态、Python 版本、包状态、系统工具状态"""
    # 检查包状态
    packages = {}
    pkg_checks = {
        "python-docx": "docx",
        "python-pptx": "pptx",
        "openpyxl": "openpyxl",
        "reportlab": "reportlab",
        "matplotlib": "matplotlib",
        "numpy": "numpy",
        "pandas": "pandas",
        "weasyprint": "weasyprint",
        "markdown": "markdown",
        "pdfplumber": "pdfplumber",
    }
    for pkg_name, import_name in pkg_checks.items():
        try:
            __import__(import_name)
            packages[pkg_name] = "ok"
        except Exception:
            packages[pkg_name] = "missing"
    
    return HealthResponse(
        status="healthy",
        timestamp=datetime.now().isoformat(),
        python_version=sys.version.split()[0],
        packages_status=packages,
        system_tools_status=SYSTEM_TOOLS,
        workspace=str(WORKSPACE_DIR),
    )


# ---- 7.2 执行 Python 代码 ----

@app.post("/execute", response_model=ExecuteResponse, summary="执行 Python 代码")
async def execute_code(req: ExecuteRequest):
    """
    在隔离的 exec 环境中执行 Python 代码。
    - 捕获 stdout/stderr/返回值
    - 支持可选包安装（pip）
    - 超时控制（默认30秒，最大120秒）
    """
    timeout = min(req.timeout or DEFAULT_TIMEOUT, MAX_TIMEOUT)
    
    result = RestrictedExec.execute(
        code=req.code,
        timeout=timeout,
        packages=req.packages,
    )
    
    return ExecuteResponse(
        ok=result["ok"],
        stdout=result["stdout"],
        stderr=result["stderr"],
        result=result["result"],
        error=result["error"],
        files=result["files"],
    )


# ---- 7.3 生成 Word 文档 ----

@app.post("/generate/docx", summary="生成 Word 文档")
async def gen_docx(req: DocxRequest):
    """用 python-docx 生成 .docx 文件，解析 Markdown 内容"""
    try:
        result = generate_docx(
            title=req.title,
            content=req.content,
            filename=req.filename,
        )
        return {"ok": True, "data": result}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"生成 Word 文档失败: {str(e)}")


# ---- 7.4 生成 PPT 课件 ----

@app.post("/generate/pptx", summary="生成 PPT 课件")
async def gen_pptx(req: PptxRequest):
    """用 python-pptx 生成 .pptx 文件，16:9 布局，主题色 jade/ink/bg"""
    try:
        result = generate_pptx(
            slides_data=req.slides,
            filename=req.filename,
        )
        return {"ok": True, "data": result}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"生成 PPT 课件失败: {str(e)}")


# ---- 7.5 生成 Excel ----

@app.post("/generate/xlsx", summary="生成 Excel 表格")
async def gen_xlsx(req: XlsxRequest):
    """用 openpyxl 生成 .xlsx 文件，列宽自适应，表头加粗"""
    try:
        result = generate_xlsx(
            rows=req.rows,
            filename=req.filename,
            sheet_name=req.sheetName,
            dir_name=req.dir,
        )
        return {"ok": True, "data": result}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"生成 Excel 失败: {str(e)}")


# ---- 7.6 生成 PDF ----

@app.post("/generate/pdf", summary="生成 PDF 文档")
async def gen_pdf(req: PdfRequest):
    """优先用 weasyprint（HTML→PDF），降级用 reportlab 生成 PDF"""
    try:
        result = generate_pdf(
            content=req.content,
            title=req.title or "文档",
            filename=req.filename,
        )
        return {"ok": True, "data": result}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"生成 PDF 失败: {str(e)}")


# ---- 7.7 文档格式转换 ----

@app.post("/convert", summary="文档格式转换")
async def convert_doc(req: ConvertRequest):
    """
    支持 pandoc 转换：md <-> docx <-> html <-> pdf
    支持 libreoffice 转换：pptx->pdf, docx->pdf, xlsx->pdf
    返回转换后的文件 base64 编码。
    """
    try:
        result = convert_document(
            from_format=req.from_format,
            to_format=req.to_format,
            file_content=req.file_content,
            filename=req.filename,
        )
        return {"ok": True, "data": result}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"格式转换失败: {str(e)}")


# ---- 7.8 OCR 文字识别 ----

@app.post("/ocr", summary="OCR 文字识别")
async def ocr_image(req: OcrRequest):
    """调用 tesseract 命令行进行 OCR 识别，返回识别文本"""
    try:
        result = run_ocr(
            image_b64=req.image,
            lang=req.lang,
        )
        return {"ok": True, "data": result}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"OCR 识别失败: {str(e)}")


# ---- 7.9 下载文件 ----

@app.get("/files/{path:path}", summary="下载文件")
async def download_file(path: str):
    """下载 workspace 内的文件"""
    file_path = validate_workspace_path(path)
    
    if not file_path.exists():
        raise HTTPException(status_code=404, detail=f"文件不存在: {path}")
    
    if not file_path.is_file():
        raise HTTPException(status_code=400, detail=f"路径不是文件: {path}")
    
    # 确定媒体类型
    ext = file_path.suffix.lower()
    media_types = {
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ".pdf": "application/pdf",
        ".png": "image/png",
        ".jpg": "image/jpeg",
        ".jpeg": "image/jpeg",
        ".gif": "image/gif",
        ".html": "text/html",
        ".csv": "text/csv",
        ".txt": "text/plain",
        ".md": "text/markdown",
    }
    media_type = media_types.get(ext, "application/octet-stream")
    
    return FileResponse(
        path=str(file_path),
        media_type=media_type,
        filename=file_path.name,
    )


# ---- 7.10 列出文件 ----

@app.get("/files", summary="列出 workspace 文件")
async def list_files():
    """列出 workspace 下所有文件（路径、大小、创建时间）"""
    files = []
    try:
        for root in DIRS.values():
            if root.exists():
                for f in root.rglob("*"):
                    if f.is_file():
                        rel = f.relative_to(WORKSPACE_DIR)
                        stat = f.stat()
                        files.append({
                            "path": str(rel).replace("\\", "/"),
                            "name": f.name,
                            "dir": str(rel.parent).replace("\\", "/"),
                            "size": stat.st_size,
                            "created": datetime.fromtimestamp(stat.st_ctime).isoformat(),
                            "modified": datetime.fromtimestamp(stat.st_mtime).isoformat(),
                            "url": file_download_url(str(rel).replace("\\", "/")),
                        })
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"列出文件失败: {str(e)}")
    
    return {
        "ok": True,
        "data": {
            "workspace": str(WORKSPACE_DIR),
            "total_files": len(files),
            "files": sorted(files, key=lambda x: x["modified"], reverse=True),
        }
    }


# ---- 7.11 生成图表 ----

@app.post("/analyze/chart", summary="生成图表")
async def gen_chart(req: ChartRequest):
    """用 matplotlib 生成图表（bar/line/pie/radar/scatter）并保存为 PNG"""
    try:
        result = generate_chart(
            chart_type=req.type,
            data=req.data,
            title=req.title,
            filename=req.filename,
        )
        return {"ok": True, "data": result}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"生成图表失败: {str(e)}")


# ---- 7.12 网页抓取 ----

def _html_to_text(html: str, max_chars: int) -> str:
    """将 HTML 转为纯文本正文（去脚本/样式/标签）"""
    import re as _re
    text = _re.sub(r"(?is)<(script|style|noscript)[^>]*>.*?</\1>", " ", html)
    text = _re.sub(r"(?s)<!--.*?-->", " ", text)
    text = _re.sub(r"(?s)<[^>]+>", " ", text)
    text = (text.replace("&nbsp;", " ").replace("&lt;", "<")
                .replace("&gt;", ">").replace("&amp;", "&")
                .replace("&quot;", '"').replace("&#39;", "'"))
    text = _re.sub(r"[ \t]+", " ", text)
    text = _re.sub(r"\n\s*\n+", "\n", text)
    return text.strip()[:max_chars]


@app.post("/fetch", summary="抓取网页正文")
async def fetch_page(req: FetchRequest):
    """服务端抓取网页并提取纯文本正文（绕过浏览器 CORS 限制）"""
    import urllib.request
    import urllib.error

    if not req.url.startswith(("http://", "https://")):
        raise HTTPException(status_code=400, detail="URL 必须以 http:// 或 https:// 开头")

    headers = {
        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 "
                      "(KHTML, like Gecko) Chrome/125.0.0.0 Safari/537.36",
        "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8",
        "Accept-Language": "zh-CN,zh;q=0.9,en;q=0.8",
    }
    try:
        request = urllib.request.Request(req.url, headers=headers)
        with urllib.request.urlopen(request, timeout=15) as resp:
            charset = resp.headers.get_content_charset() or "utf-8"
            raw = resp.read(5 * 1024 * 1024)  # 最多读 5MB
            try:
                html = raw.decode(charset, errors="replace")
            except (LookupError, TypeError):
                html = raw.decode("utf-8", errors="replace")
        content = _html_to_text(html, req.max_chars)
        return {"ok": True, "data": {"url": req.url, "content": content, "source": "server-fetch"}}
    except urllib.error.HTTPError as e:
        raise HTTPException(status_code=502, detail=f"目标网页返回 {e.code}")
    except Exception as e:
        raise HTTPException(status_code=502, detail=f"抓取失败: {str(e)}")


# ============================================================
# 8. 异常处理
# ============================================================

@app.exception_handler(HTTPException)
async def http_exception_handler(request: Request, exc: HTTPException):
    """统一 HTTP 异常处理"""
    return JSONResponse(
        status_code=exc.status_code,
        content={
            "ok": False,
            "error": exc.detail,
            "status_code": exc.status_code,
        },
    )


@app.exception_handler(Exception)
async def general_exception_handler(request: Request, exc: Exception):
    """通用异常处理"""
    return JSONResponse(
        status_code=500,
        content={
            "ok": False,
            "error": f"服务器内部错误: {str(exc)}",
            "traceback": traceback.format_exc(),
        },
    )


# ============================================================
# 9. 启动初始化
# ============================================================

@app.on_event("startup")
async def startup_event():
    """服务启动时初始化"""
    global SYSTEM_TOOLS
    print("=" * 60, file=sys.stderr)
    print("教师智能体沙箱 · 服务器端 API 服务", file=sys.stderr)
    print("=" * 60, file=sys.stderr)
    
    # 检测系统工具
    SYSTEM_TOOLS = detect_system_tools()
    
    # 初始化工作区
    for name, d in DIRS.items():
        d.mkdir(parents=True, exist_ok=True)
        print(f"[sandbox] 目录就绪: {name}/", file=sys.stderr)
    
    # 设置 matplotlib 中文字体
    cn_font = setup_matplotlib_cn_font()
    if cn_font:
        print(f"[sandbox] matplotlib 中文字体: {cn_font}", file=sys.stderr)
    else:
        print("[sandbox] 警告: 未找到中文字体，图表中文可能显示为方块", file=sys.stderr)
    
    # 打印系统工具状态
    for tool, info in SYSTEM_TOOLS.items():
        status = "可用" if info.get("available") else "不可用"
        print(f"[sandbox] {tool}: {status} ({info.get('version', '')})", file=sys.stderr)
    
    # 打印包状态
    print(f"[sandbox] weasyprint: {'可用' if HAS_WEASYPRINT else '不可用(将降级为 reportlab)'}", file=sys.stderr)
    print(f"[sandbox] markdown: {'可用' if HAS_MARKDOWN else '不可用'}", file=sys.stderr)
    
    print(f"[sandbox] 工作区: {WORKSPACE_DIR}", file=sys.stderr)
    print(f"[sandbox] 服务就绪，等待请求...", file=sys.stderr)
    print("=" * 60, file=sys.stderr)


# ============================================================
# 10. Main 入口
# ============================================================

if __name__ == "__main__":
    import uvicorn
    
    host = os.environ.get("SANDBOX_HOST", "0.0.0.0")
    port = int(os.environ.get("SANDBOX_PORT", "8000"))
    
    print(f"启动教师智能体沙箱 API 服务: http://{host}:{port}")
    print(f"API 文档: http://{host}:{port}/docs")
    
    uvicorn.run(
        "sandbox_api:app",
        host=host,
        port=port,
        reload=True,
        log_level="info",
    )
